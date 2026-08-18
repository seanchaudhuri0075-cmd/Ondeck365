"""Freeform template renderer — positions every shape at its REAL PPTX coordinates.

Unlike the other six templates (cover, title_stats, card_grid, section_divider,
media_showcase, two_column), which encode a hand-tuned layout built around the
specific decks they were designed for (P&G / GIF creative pitch decks), this
renderer makes no assumption about which slide archetype it's looking at. It
reads each shape's real x/y/width/height (via `parse.shapes.flatten_slide`,
which already handles nested-group transform composition) and places it on
the canvas at that exact position, scaled as a percentage of canvas size via
`render.desktop.positioned_style` — the same primitive the other templates
use for their own hardcoded boxes, just applied shape-by-shape instead of to
one fixed layout.

This is the fallback for any slide that doesn't match one of the six curated
archetypes — which, for a deck the pipeline has never seen before (no
manifest, no per-slide hints), is every slide until auto-classification
exists. It trades "polished, deck-specific layout" for "structurally
faithful to the source file," which is the correct trade for an unknown
deck: a plain, correctly-positioned rendering beats a polished but
wrong-shaped one.

Background: the first full-canvas shape with a solid fill (by z-order) is
treated as the slide's background color; if none is found, deck_brand_color
is used. Non-background shapes are rendered in z-order:
  - <p:pic>            -> absolutely-positioned <img>
  - text-bearing <p:sp> -> absolutely-positioned text block, paragraph/run
                           styling preserved (size, weight, color, align)
  - solid-fill <p:sp>   -> absolutely-positioned color rect (decorative
                           elements: dividers, underlines, color blocks)
  - anything else       -> skipped (connectors, unresolvable groups, etc.)

Mobile gets a simplified vertical reflow of the same content in document
order — not pixel-faithful (that's structurally impossible for a stacked
layout), but readable and content-complete.
"""
from __future__ import annotations

import re
from typing import Optional

from pptx.slide import Slide

from ..css import color_with_alpha
from ..desktop import canvas_aspect_css, positioned_style
from ..html import render_page
from ...layout.detect import SlideClass
from ...parse.font_calibration import MATCHED_METRIC_SUBS, calibrate_size_pt
from ...layout.archetype import classify_slide
from ..fonts import BUNDLED_FALLBACKS
from ...parse.shapes import FlatShape, flatten_slide
from ...parse.slide import NS
from ...parse.text import parse_text_frame
from ...parse.theme import Theme
from ._shared import image_src

_ALIGN_CSS = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}
_ANCHOR_JUSTIFY = {"t": "flex-start", "ctr": "center", "b": "flex-end"}

# PPTX/PowerPoint typefaces routinely encode weight as a literal suffix in
# the family name ("Helvetica Light", "Univers Condensed Bold") rather than
# a separate weight attribute. Browsers/fontconfig only resolve a compound
# name like that if a font is installed under that EXACT string — "Helvetica"
# alone commonly has a real metric-compatible alias (e.g. fontconfig maps it
# to TeX Gyre Heros / Nimbus Sans on Linux, and macOS has real Helvetica),
# but "Helvetica Light" as one string usually matches nothing and falls
# through to a generic default with completely different metrics, which
# then throws off text width, wrapping, and everything downstream of it.
# Splitting the recognized weight word off and expressing it as a real CSS
# font-weight lets the base family resolve correctly.
_WEIGHT_SUFFIXES = {
    "thin": 100, "extralight": 200, "ultralight": 200,
    "light": 300, "regular": 400, "normal": 400,
    "medium": 500, "semibold": 600, "demibold": 600,
    "bold": 700, "extrabold": 800, "ultrabold": 800,
    "black": 900, "heavy": 900,
}


def _split_typeface(typeface: str) -> tuple[str, Optional[int]]:
    """Split a trailing weight keyword off a compound typeface name.

    'Helvetica Light' -> ('Helvetica', 300). 'DIN Condensed' -> unchanged,
    since 'condensed' is a width keyword, not a weight one (handled
    separately via font-stretch). Returns (typeface, None) if no
    recognized weight suffix is present.
    """
    parts = typeface.rsplit(" ", 1)
    if len(parts) == 2:
        base, suffix = parts
        weight = _WEIGHT_SUFFIXES.get(suffix.lower())
        if weight is not None and base:
            return base, weight
    return typeface, None


def _is_condensed(typeface: str) -> bool:
    name = typeface.lower()
    return "condensed" in name or "narrow" in name


def _font_family_css(typeface: str | None) -> str:
    """Resolve a run's real declared typeface to a CSS font-family stack.

    freeform is the generic fallback for a deck we know nothing about, so
    it must not assume every run wants the same (often condensed) display
    font — that silently narrows/widens text versus the source and throws
    off wrapping, which cascades into broken justify/line-count/vertical
    space. Stack, in order: the deck's own declared typeface (works if the
    viewer's machine happens to have it — true on the deck author's own
    machine, not reliably true anywhere else); the weight-stripped base
    name (resolves real OS font aliases the compound name misses, e.g.
    "Helvetica" but not "Helvetica Light"); a bundled, offline-embedded,
    metric-compatible substitute keyed off either name (this is what
    actually determines rendered width on a machine that has none of the
    above — see render/fonts.py); then generic sans-serif as a last resort
    so a genuinely unrecognized font still renders something.
    """
    if typeface is None:
        return "sans-serif"
    base, _weight = _split_typeface(typeface)
    stack = [f"'{typeface}'"]
    if base != typeface:
        stack.append(f"'{base}'")  # weight-stripped fallback, resolves real aliases the compound name misses
    sub = MATCHED_METRIC_SUBS.get(typeface.lower())
    if sub:
        stack.append(f"'{sub.title()}'")
    bundled = BUNDLED_FALLBACKS.get(typeface.lower()) or BUNDLED_FALLBACKS.get(base.lower())
    if bundled and bundled not in stack:
        stack.append(f"'{bundled}'")
    stack.append("sans-serif")
    return ", ".join(stack)


def render_freeform(
    slide: Slide,
    theme: Theme,
    slide_index: int,
    slide_class: SlideClass,
    deck_name: str,
    slide_w_pt: float,
    slide_h_pt: float,
    deck_brand_color: str,
    media_dir=None,
    default_tab_pt: float = 72.0,
    section_headers: Optional[dict] = None,
) -> str:
    """Render any slide by placing each shape at its real PPTX position."""
    flat = list(flatten_slide(slide))

    bg_shape, bg_hex = _find_background(flat, slide_w_pt, slide_h_pt, theme)
    content = [s for s in flat if s is not bg_shape]

    items = []  # list of dicts describing one renderable shape, in z-order
    for s in content:
        if s.x_pt is None:
            continue  # geometry didn't resolve (unhandled group case)
        if s.kind == "pic":
            src = image_src(s, slide_class, slide)
            if src is None:
                continue
            items.append({"type": "image", "shape": s, "src": src})
            continue
        if s.kind == "sp":
            tf = parse_text_frame(s.element, theme)
            has_text = tf is not None and any(
                r.text.strip() for p in tf.paragraphs for r in p.runs
            )
            if has_text:
                items.append({"type": "text", "shape": s, "frame": tf})
                continue
            fill_hex = _solidfill_hex(s.element, theme)
            if fill_hex is not None:
                items.append({"type": "rect", "shape": s, "color": fill_hex})
                continue
        # cxnSp / graphicFrame / unresolvable — skip; no safe generic rendering.

    archetype = classify_slide(items, slide_w_pt)
    is_acrostic = archetype.label.startswith("acrostic_bleed")
    wordmark = _extract_wordmark(items) if is_acrostic else None
    section_words = _extract_section_words(items) if is_acrostic else None
    active_word = _extract_active_section_word(items) if is_acrostic else None
    has_full_bleed_photo = any(
        it["type"] == "image" and it["shape"].w_pt and it["shape"].w_pt >= slide_w_pt * 0.5
        for it in items
    )
    onscreen_acrostic_idx = _find_onscreen_accent_acrostic(items) if has_full_bleed_photo and not is_acrostic else []
    if onscreen_acrostic_idx:
        # A different pattern from the off-canvas bleeding acrostic: here
        # the accent-styled word shapes sit directly on top of a photo
        # (matching specific window openings in the source composition),
        # not bleeding off-canvas. That precise alignment has no
        # equivalent once the photo gets cropped/resized for a narrow
        # mobile screen — confirmed directly (checked every slide for
        # this pattern; only this one has it, so this is a targeted
        # fix, not a new general archetype). Simplest faithful mobile
        # translation: drop the photo, keep the real accent-styled text
        # (already correct via each run's own alpha in the source),
        # and let any other on-canvas text act as a corner caption
        # instead of competing with a photo that no longer aligns.
        mobile_treatment = "onscreen_acrostic"
    elif is_acrostic and has_full_bleed_photo:
        mobile_treatment = "photo_bleed"
    elif is_acrostic:
        mobile_treatment = "ghost_bleed"
    else:
        mobile_treatment = "default"
    ghost_word = None
    if mobile_treatment == "ghost_bleed":
        ghost_word = active_word or wordmark

    slide_css = _build_css(slide_w_pt, slide_h_pt, items, default_tab_pt, mobile_treatment)
    body_html = _build_body(
        items, slide_w_pt, slide_h_pt, default_color=theme.dk1,
        wordmark=wordmark, mobile_treatment=mobile_treatment, ghost_word=ghost_word,
        section_words=section_words, active_word=active_word,
        onscreen_acrostic_idx=onscreen_acrostic_idx,
        section_headers=section_headers,
    )

    return render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg": f"#{bg_hex or (deck_brand_color.lstrip('#') or '000000')}",
            "headline": "#FFFFFF",
            "font-cond": '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=body_html,
        slide_css=slide_css,
    )


# ---------------------------------------------------------------------------
# Background detection
# ---------------------------------------------------------------------------

def _is_full_canvas(s: FlatShape, slide_w_pt: float, slide_h_pt: float) -> bool:
    if s.x_pt is None:
        return False
    return (
        s.w_pt >= slide_w_pt * 0.9
        and s.h_pt >= slide_h_pt * 0.9
        and s.x_pt <= slide_w_pt * 0.05
        and s.y_pt <= slide_h_pt * 0.05
    )


def _find_background(
    flat: list, slide_w_pt: float, slide_h_pt: float, theme: Theme,
) -> tuple[Optional[FlatShape], Optional[str]]:
    """First (lowest z) full-canvas solid-fill sp becomes the background."""
    for s in sorted(flat, key=lambda s: s.z):
        if s.kind != "sp" or not _is_full_canvas(s, slide_w_pt, slide_h_pt):
            continue
        hex_val = _solidfill_hex(s.element, theme)
        if hex_val is not None:
            return s, hex_val
    return None, None


def _solidfill_hex(sp_elem, theme: Theme) -> Optional[str]:
    """Pull an inline solidFill hex from a <p:sp>'s spPr. None if not a solid fill."""
    srgb = sp_elem.find("p:spPr/a:solidFill/a:srgbClr", NS)
    if srgb is not None:
        return srgb.get("val", "000000").upper()
    scheme = sp_elem.find("p:spPr/a:solidFill/a:schemeClr", NS)
    if scheme is not None:
        try:
            return theme.resolve(scheme.get("val"))
        except AttributeError:
            return None
    return None


# ---------------------------------------------------------------------------
# Desktop: absolute positioning at real coordinates
# ---------------------------------------------------------------------------

def _build_css(slide_w_pt: float, slide_h_pt: float, items: list, default_tab_pt: float, mobile_treatment: str = "default") -> str:
    canvas_w, canvas_h = canvas_aspect_css(slide_w_pt, slide_h_pt)
    tab_size_cqw = f"{default_tab_pt / slide_w_pt * 100:.3f}cqw"
    parts = [
        "#deck-desktop { display: flex; align-items: center; justify-content: center; "
        "width: 100vw; height: 100vh; }",
        "#deck-desktop .canvas {\n"
        f"  position: relative; width: {canvas_w}; height: {canvas_h};\n"
        "  overflow: hidden; container-type: inline-size;\n"
        "  background: var(--bg);\n"
        "}",
    ]
    for i, item in enumerate(items):
        s = item["shape"]
        pos = positioned_style(s.x_pt, s.y_pt, s.w_pt, s.h_pt, slide_w_pt, slide_h_pt)
        selectors = [f"#deck-desktop .item-{i}"]
        if mobile_treatment == "onscreen_acrostic":
            # Same item markup gets reused verbatim inside the mobile
            # crop wrapper (see _build_body) — needs the same
            # positioning rules under that selector too, or every shape
            # collapses to position:static.
            selectors.append(f"#deck-mobile .acrostic-crop .item-{i}")
        sel = ", ".join(selectors)
        if item["type"] == "rect":
            parts.append(
                f"{sel} {{ {pos} background: #{item['color']}; }}"
            )
        elif item["type"] == "image":
            parts.append(f"{sel} {{ {pos} }}")
            # Append " img" to each selector individually before joining
            # — joining first and appending once only attaches " img" to
            # the last selector in the list, leaving any earlier
            # selector (here, the desktop one) referring to the .item-N
            # div itself rather than its inner <img>. That malformed
            # rule then set width/height:100% on the div, and — same
            # specificity, later in the stylesheet — won the cascade
            # over the correct percentage-based sizing above, which is
            # exactly what produced the blown-up white box.
            img_sel = ", ".join(f"{s2} img" for s2 in selectors)
            parts.append(
                f"{img_sel} {{ width: 100%; height: 100%; "
                f"object-fit: contain; }}"
            )
        elif item["type"] == "text":
            anchor = item["frame"].anchor
            justify = _ANCHOR_JUSTIFY.get(anchor, "flex-start")
            parts.append(
                f"{sel} {{ {pos} display: flex; "
                f"flex-direction: column; justify-content: {justify}; }}"
            )
            p_sel = ", ".join(f"{s2} p" for s2 in selectors)
            parts.append(f"{p_sel} {{ tab-size: {tab_size_cqw}; }}")

    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    parts.append("  #deck-desktop { display: none; }")
    parts.append(
        "  #deck-mobile { display: flex; flex-direction: column; "
        "justify-content: center; min-height: 100vh; background: var(--bg); }"
    )
    parts.append("  #deck-mobile.pinned-top { justify-content: flex-start; }")
    parts.append(
        "  #deck-mobile .panel { padding: 6vw; box-sizing: border-box; width: 100%; }"
    )
    parts.append(
        "  #deck-mobile .block { margin-bottom: 4vw; width: 100%; box-sizing: border-box; }"
    )
    parts.append("  #deck-mobile .block p { width: 100%; box-sizing: border-box; }")
    parts.append("  #deck-mobile .block img { width: 100%; height: auto; }")
    parts.append("  #deck-mobile { position: relative; overflow: hidden; }")
    parts.append(
        "  #deck-mobile .panel { position: relative; z-index: 1; }"
    )
    parts.append(
        "  #deck-mobile .panel-rounded { position: relative; z-index: 1; }"
    )
    parts.append(
        "  #deck-mobile .photo-hero { position: absolute; top: 0; left: 0; right: 0; "
        "height: 100vh; overflow: hidden; z-index: 0; }"
    )
    parts.append(
        "  #deck-mobile .photo-spacer { height: 62vh; }"
    )
    parts.append(
        "  #deck-mobile .photo-hero img { width: 100%; height: 100%; object-fit: cover; "
        "object-position: center center; }"
    )
    parts.append(
        "  #deck-mobile .photo-hero::after { content: ''; position: absolute; inset: 0; "
        "background: linear-gradient(180deg, rgba(0,0,0,0.65) 0%, rgba(0,0,0,0.15) 16%, "
        "rgba(0,0,0,0.08) 30%, rgba(0,0,0,0.1) 42%, rgba(0,0,0,0.55) 58%, "
        "rgba(5,5,5,0.9) 72%, var(--bg) 88%); }"
    )
    parts.append(
        "  #deck-mobile .signage { display: flex; gap: 0.18em; font-family: "
        "'Barlow Condensed', sans-serif; font-weight: 700; font-size: 7vw; "
        "letter-spacing: 0.15em; color: #fff; text-shadow: "
        "0 0 4vw rgba(255,255,255,0.6), 0 0.2vw 1vw rgba(0,0,0,0.5); }"
    )
    parts.append(
        "  #deck-mobile .photo-hero .signage { position: absolute; top: 5vw; left: 5vw; "
        "z-index: 2; }"
    )
    parts.append(
        "  #deck-mobile .photo-hero .section-eyebrow { position: absolute; top: 13.5vw; "
        "left: 5vw; right: 5vw; z-index: 2; margin-bottom: 0; "
        "text-shadow: 0 0.15vw 0.8vw rgba(0,0,0,0.7); }"
    )
    parts.append(
        "  #deck-mobile .floor-badge { position: absolute; top: 5vw; right: 5vw; "
        "z-index: 2; text-align: right; color: #fff; "
        "text-shadow: 0 0.15vw 0.8vw rgba(0,0,0,0.7); }"
    )
    parts.append(
        "  #deck-mobile .floor-badge p:first-child { font-size: 3vw; font-weight: 700; "
        "letter-spacing: 0.1em; opacity: 0.85; }"
    )
    parts.append(
        "  #deck-mobile .floor-badge p:last-child { font-size: 9vw; font-weight: 800; "
        "line-height: 1; }"
    )
    parts.append(
        "  #deck-mobile .acrostic-crop-wrap { position: absolute; inset: 0; "
        "overflow: hidden; z-index: 1; }"
    )
    parts.append(
        "  #deck-mobile .acrostic-crop { position: absolute; top: 50%; left: 50%; "
        "transform: translate(-50%, -50%); overflow: hidden; "
        "container-type: inline-size; background: var(--bg); }"
    )
    parts.append(
        "  #deck-mobile .signage-standalone { position: relative; z-index: 1; "
        "margin-bottom: 3vw; }"
    )
    parts.append(
        "  #deck-mobile .ghost-word { position: absolute; left: -10vw; right: -10vw; "
        "font-family: 'Barlow Condensed', sans-serif; font-weight: 800; font-size: 32vw; "
        "line-height: 0.85; letter-spacing: -0.02em; color: rgba(255,255,255,0.06); "
        "white-space: nowrap; z-index: 0; text-align: center; pointer-events: none; }"
    )
    parts.append(
        "  #deck-mobile .ghost-fade { position: absolute; top: 0; left: 0; right: 0; "
        "height: 45vw; z-index: 0; pointer-events: none; "
        "background: linear-gradient(180deg, var(--bg) 0%, rgba(0,0,0,0.55) 55%, "
        "transparent 100%); }"
    )
    parts.append(
        "  #deck-mobile .section-eyebrow { position: relative; z-index: 1; "
        "font-family: 'Barlow Condensed', sans-serif; font-size: 4.6vw; font-weight: 700; "
        "letter-spacing: 0.1em; text-transform: uppercase; margin-bottom: 4vw; "
        "color: rgba(255,255,255,0.45); line-height: 1.3; }"
    )
    parts.append("  #deck-mobile .section-eyebrow .dot { margin: 0 0.35em; opacity: 0.5; }")
    parts.append(
        "  #deck-mobile .section-eyebrow .active { color: var(--headline); font-weight: 800; "
        "text-shadow: 0 0 2.5vw rgba(255,255,255,0.5); }"
    )
    parts.append("}")
    return "\n".join(parts)


def _off_canvas_text_by_position(items: list) -> list:
    """Off-canvas acrostic text items, sorted by real vertical position.

    Document order (the order shapes are declared in the slide XML) does
    NOT reliably match visual top-to-bottom order — confirmed directly:
    one slide in this deck declares its "HAIR" shape before "SKIN" in the
    XML, even though SKIN sits visually above HAIR on the canvas. Sorting
    by each shape's actual y_pt is what "top-to-bottom reading order"
    really means; z-order is an authoring artifact, not a layout signal.
    """
    candidates = [
        it for it in items
        if it["type"] == "text" and it["shape"].x_pt is not None and it["shape"].x_pt < 0
        and it["shape"].y_pt is not None
    ]
    return sorted(candidates, key=lambda it: it["shape"].y_pt)


def _find_onscreen_accent_acrostic(items: list) -> list:
    """Find on-canvas text shapes using the deck's accent-letter pattern
    (first run at higher alpha than the rest) directly overlaid on a
    photo, as opposed to the off-canvas bleeding acrostic elsewhere in
    this deck. Returns the matching items' indices, or an empty list if
    fewer than 3 shapes match (3 is a conservative floor — enough to be
    confident this is a deliberate repeated pattern, not a coincidence).
    """
    matches = []
    for i, it in enumerate(items):
        if it["type"] != "text" or it["shape"].x_pt is None or it["shape"].x_pt < 0:
            continue
        paragraphs = it["frame"].paragraphs
        if not paragraphs or len(paragraphs[0].runs) != 2:
            continue
        r0, r1 = paragraphs[0].runs
        a0 = r0.style.color_alpha if r0.style.color_alpha is not None else 1.0
        a1 = r1.style.color_alpha if r1.style.color_alpha is not None else 1.0
        if a0 - a1 > 0.15:  # first run meaningfully brighter than the rest
            matches.append(i)
    return matches if len(matches) >= 3 else []


def _extract_active_section_word(items: list) -> Optional[str]:
    """Find the off-canvas acrostic word that's fully highlighted on this slide.

    Most slides highlight only the first letter of each word in the
    decorative column (S-KIN, H-AIR, ...) — that's what _extract_wordmark
    reads. But on a slide "for" a given section, PowerPoint highlights the
    *entire* word instead (all of "SKIN", not just "S") to show which
    section you're currently in — a single run covering the whole word,
    rather than the usual two-run split. Reuse that same distinction:
    when exactly one run covers the whole word, this slide has an active
    section, and that word (not the brand initials) is the right thing to
    feature on mobile.
    """
    for it in _off_canvas_text_by_position(items):
        paragraphs = it["frame"].paragraphs
        if not paragraphs or not paragraphs[0].runs:
            continue
        p0 = paragraphs[0]
        if len(p0.runs) != 1:
            continue  # the normal two-run (accent letter + dim rest) case
        word = p0.runs[0].text.strip()
        if len(word) > 1:
            return word
    return None


def _extract_wordmark(items: list) -> Optional[str]:
    """Pull the highlighted first letter out of each off-canvas acrostic shape.

    This deck (and the pattern is common enough to be worth generalizing,
    not deck-specific plumbing) encodes "this is the accent letter" the
    same way PowerPoint always encodes emphasis within a run of otherwise-
    uniform text: a distinct color on just that run. Concretely, each
    shape here has two runs — the first letter in the theme's light color
    (bg1), the rest of the word in a dimmer gray — which is exactly how a
    human reading the slide perceives "S-H-E-L-F" as the brand name inside
    "SKIN / HAIR / EYES / LIPS / FACE". Reuse that same signal rather than
    inventing a new one: take each shape's first run's text, sorted into
    real top-to-bottom visual order (see _off_canvas_text_by_position —
    document order isn't reliable for this).
    """
    letters = []
    for it in _off_canvas_text_by_position(items):
        paragraphs = it["frame"].paragraphs
        if not paragraphs or not paragraphs[0].runs:
            continue
        first_run_text = paragraphs[0].runs[0].text.strip()
        if first_run_text:
            # Only the first character: on the slide for whichever section
            # is "currently active", PowerPoint highlights the entire word
            # as one run (not just its first letter) to show where you are
            # — taking the whole run's text there would grab "SKIN" instead
            # of "S", producing "SKINHELF" instead of "SHELF".
            letters.append(first_run_text[0])
    word = "".join(letters)
    return word if word else None


def _extract_section_words(items: list) -> list[str]:
    """Pull the full word out of each off-canvas acrostic shape, in order.

    Companion to _extract_wordmark (which pulls just the accent letter):
    this pulls the whole word each shape spells — concatenating all its
    runs, not just the first — for a compact "SKIN · HAIR · EYES · LIPS ·
    FACE" style listing. Deck-agnostic by construction: whatever words the
    source deck's own shapes actually spell, sorted into real top-to-
    bottom visual order (see _off_canvas_text_by_position), not a
    hardcoded list specific to this one deck's section names.
    """
    words = []
    for it in _off_canvas_text_by_position(items):
        paragraphs = it["frame"].paragraphs
        if not paragraphs:
            continue
        word = "".join(r.text for r in paragraphs[0].runs).strip()
        if word:
            words.append(word)
    return words


def _build_eyebrow_html(section_words: list, active_word: Optional[str]) -> str:
    """Compact "SKIN · HAIR · EYES · LIPS · FACE" listing — the active
    section (if any) reads brighter/bolder than the rest, echoing the
    same highlight signal the desktop acrostic and ghost-bleed word
    already use. Shared between photo_bleed (overlaid on the photo) and
    ghost_bleed (at the top of the panel) so both treatments show the
    same wayfinding element instead of only one of them having it.
    """
    spans = []
    for w in section_words:
        cls = ' class="active"' if active_word and w.upper() == active_word.upper() else ""
        spans.append(f"<span{cls}>{_escape(w)}</span>")
    return '<div class="section-eyebrow">' + '<span class="dot">·</span>'.join(spans) + "</div>"


def _build_body(
    items: list,
    slide_w_pt: float,
    slide_h_pt: float,
    default_color: str = "000000",
    wordmark: Optional[str] = None,
    mobile_treatment: str = "default",
    ghost_word: Optional[str] = None,
    section_words: Optional[list] = None,
    active_word: Optional[str] = None,
    section_headers: Optional[dict] = None,
    onscreen_acrostic_idx: Optional[list] = None,
) -> str:
    desktop_parts = []
    mobile_parts = []
    mobile_prefix = ""  # rendered before/outside .panel — hero photo or ghost-bleed layers

    hero_photo_index = None
    hero_extra_skip = set()
    if mobile_treatment == "photo_bleed":
        # Pick the first full-bleed-ish photo (matches the same threshold
        # used to decide this treatment) as the mobile hero — rendered
        # once, full-width, ahead of the panel, not duplicated as a
        # regular block further down.
        for i, item in enumerate(items):
            if item["type"] == "image" and item["shape"].w_pt and item["shape"].w_pt >= slide_w_pt * 0.5:
                hero_photo_index = i
                break
        if hero_photo_index is not None:
            signage = (
                "".join(f"<span>{_escape(ch)}</span>" for ch in wordmark)
                if wordmark else ""
            )
            eyebrow = _build_eyebrow_html(section_words, active_word) if section_words else ""
            # The floor/level label and its number are two separate text
            # shapes in the source, meant to read as one badge (matches
            # the desktop layout, top-right). Rendered as a normal block
            # in mobile flow, they landed wherever they happened to fall
            # in document order — usually deep in the photo's low-
            # contrast middle zone, easy to miss entirely. Pull them out
            # and overlay them explicitly, the same way signage already
            # gets pulled out and overlaid instead of left in normal flow.
            floor_label_idx = floor_num_idx = None
            for i, it in enumerate(items):
                if it["type"] != "text":
                    continue
                text = "".join(
                    r.text for p in it["frame"].paragraphs for r in p.runs
                ).strip()
                if floor_label_idx is None and re.match(r"^FLOOR\s*LEVEL$", text, re.I):
                    floor_label_idx = i
                elif floor_label_idx is not None and floor_num_idx is None and text.isdigit():
                    floor_num_idx = i
                    break
            floor_badge = ""
            if floor_label_idx is not None and floor_num_idx is not None:
                label_text = "".join(
                    r.text for p in items[floor_label_idx]["frame"].paragraphs for r in p.runs
                ).strip()
                num_text = "".join(
                    r.text for p in items[floor_num_idx]["frame"].paragraphs for r in p.runs
                ).strip()
                floor_badge = (
                    '<div class="floor-badge">'
                    f"<p>{_escape(label_text)}</p><p>{_escape(num_text)}</p>"
                    "</div>"
                )
                hero_extra_skip.update({floor_label_idx, floor_num_idx})
            mobile_prefix = (
                '<div class="photo-hero">'
                f'<img src="{items[hero_photo_index]["src"]}" alt="">'
                + (f'<div class="signage">{signage}</div>' if signage else "")
                + eyebrow
                + floor_badge
                + "</div>"
                '<div class="photo-spacer"></div>'
            )
    elif mobile_treatment == "ghost_bleed" and ghost_word:
        # Three staggered copies of the section word, huge and nearly
        # transparent, standing in for the desktop acrostic's vertical
        # bleed column — texture behind the content rather than a
        # standalone element competing with it. The fade div (painted
        # after, so on top, within the same z-index-0 layer) dims the
        # ghost text specifically near the top, where the eyebrow and
        # header actually sit, so it doesn't visually clash with them —
        # it's still visible further down against otherwise-empty space.
        mobile_prefix = "".join(
            f'<div class="ghost-word" style="top:{top}">{_escape(ghost_word)}</div>'
            for top in ("-4vw", "34vw", "72vw")
        ) + '<div class="ghost-fade"></div>'

    for i, item in enumerate(items):
        if item["type"] == "image":
            desktop_parts.append(
                f'<div class="item-{i}"><img src="{item["src"]}" alt=""></div>'
            )
            if i != hero_photo_index and mobile_treatment != "onscreen_acrostic":
                # The hero photo is already shown once via mobile_prefix —
                # only skip the *mobile* duplicate here. Desktop always
                # renders every image at its real position regardless of
                # which one (if any) mobile chose as its hero.
                mobile_parts.append(
                    f'<div class="block"><img src="{item["src"]}" alt=""></div>'
                )
        elif item["type"] == "rect":
            desktop_parts.append(f'<div class="item-{i}"></div>')
            # Decorative color blocks add no mobile content value; skip.
        elif item["type"] == "text":
            html_p = _render_paragraphs(
                item["frame"].paragraphs, slide_w_pt, mobile=False, default_color=default_color
            )
            desktop_parts.append(f'<div class="item-{i}">{html_p}</div>')
            s = item["shape"]
            # Shapes positioned to start left of the canvas edge are
            # deliberately bleeding off-frame as a desktop background
            # texture (e.g. this deck's giant SKIN/HAIR/EYES/LIPS/FACE
            # column) — there's no "off-canvas" on a stacked mobile
            # layout, so rendering them as their own full-width blocks
            # just produces several giant, near-meaningless blocks ahead
            # of the real content. The photo_bleed/ghost_bleed prefix
            # above is the mobile stand-in for this texture instead.
            if s.x_pt is not None and s.x_pt < 0:
                continue
            if i in hero_extra_skip:
                # Pulled out into the floor-badge overlay above instead —
                # avoid rendering it a second time as a plain block.
                continue
            html_p_m = _render_paragraphs(
                item["frame"].paragraphs, slide_w_pt, mobile=True, default_color=default_color
            )
            if not html_p_m.strip():
                continue
            if mobile_treatment != "onscreen_acrostic":
                mobile_parts.append(f'<div class="block">{html_p_m}</div>')
            # onscreen_acrostic: nothing added to mobile_parts here — this
            # treatment reuses desktop_parts directly (see below) instead
            # of a separate mobile text reconstruction. A parallel build
            # means resolving fonts/positioning twice, which can quietly
            # drift from desktop over time; reusing the exact same markup
            # makes that drift structurally impossible.

    if mobile_treatment == "onscreen_acrostic":
        # Cover-crop the real desktop canvas (same photo, same
        # positioned text, zero separate rendering) to fill the mobile
        # viewport, the same way a photo would be object-fit:cover
        # cropped — except here the "photo" is the whole composed
        # canvas, so cropped text stays pixel-consistent with desktop
        # by construction rather than by matching fonts/sizes by hand.
        # crop_frac < 1 keeps more of the canvas's width in view than a
        # full-height crop would (matches a moderate, non-tightly-zoomed
        # reference crop) — both height and width stay in vh so their
        # ratio (and therefore the crop's aspect) doesn't drift; a CSS
        # percentage for height here would compute against the
        # viewport-height wrapper while width (in vh) stays anchored to
        # the full viewport, breaking the aspect the moment they differ.
        crop_frac = 0.5
        w_to_h = slide_w_pt / slide_h_pt * 100 * crop_frac
        # Anchor the crop horizontally on the real acrostic content's
        # midpoint, not the canvas's raw geometric center — confirmed
        # directly that the content isn't centered within the photo
        # (here it sits at ~55%, not 50%), so a center-of-canvas crop
        # was symmetric with respect to the photo but still cut into
        # the content asymmetrically. Computed from real shape
        # positions so this generalizes rather than hardcoding a value
        # specific to this one deck.
        acrostic_shapes = [items[i]["shape"] for i in onscreen_acrostic_idx]
        content_left = min(s.x_pt for s in acrostic_shapes)
        content_right = max(s.x_pt + s.w_pt for s in acrostic_shapes)
        anchor_pct = (content_left + content_right) / 2 / slide_w_pt * 100
        # desktop_parts[i] corresponds 1:1 to items[i] (built in the same
        # loop above, unconditionally appended). Keep only the photo and
        # the acrostic text itself — other on-canvas text (e.g. "Thank
        # You") isn't part of this pattern and the wider crop needed to
        # avoid cutting the acrostic's edge letters also reveals enough
        # of the canvas corner to bleed that text partially into frame.
        crop_parts = [
            part for i, part in enumerate(desktop_parts)
            if i in onscreen_acrostic_idx or items[i]["type"] == "image"
        ]
        mobile_prefix = (
            '<div class="acrostic-crop-wrap">'
            f'<div class="canvas acrostic-crop" '
            f'style="width:{w_to_h:.2f}vh; height:{crop_frac * 100:.2f}vh; '
            f'transform:translate(-{anchor_pct:.2f}%, -50%);">\n'
            + "\n".join(crop_parts)
            + "\n</div></div>"
        )

    if mobile_treatment == "photo_bleed" and active_word is not None and section_headers is not None:
        def is_substantial(text: str) -> bool:
            return bool(text) and text != "BEAUTY" and not re.match(r"^FLOOR\s*LEVEL$", text, re.I) and not text.isdigit()

        own_header = None
        for block in mobile_parts:
            inner = re.sub(r'^<div class="block">|</div>$', "", block)
            text = re.sub(r"<[^>]+>", "", inner).replace("&nbsp;", " ").strip()
            if is_substantial(text):
                own_header = block
                break
        if own_header is not None:
            section_headers.setdefault(active_word, own_header)
        elif active_word in section_headers:
            # This "closer" slide (e.g. a section's second/final photo,
            # no caption of its own in the source) borrows the header its
            # section's hero slide already showed, rather than reading as
            # a bare photo with just BEAUTY and a floor number.
            mobile_parts.append(section_headers[active_word])

    if mobile_treatment == "ghost_bleed" and section_words:
        eyebrow_html = _build_eyebrow_html(section_words, active_word)
        mobile_parts.insert(0, eyebrow_html)
        if wordmark:
            # Same signage treatment photo slides use, reused here so
            # SHELF reads with consistent brand presence on every mobile
            # slide, not just the ones with a photo to overlay it on.
            signage_html = (
                '<div class="signage signage-standalone">'
                + "".join(f"<span>{_escape(ch)}</span>" for ch in wordmark)
                + "</div>"
            )
            mobile_parts.insert(0, signage_html)

    section_classes = []
    if mobile_prefix:
        section_classes.append("pinned-top")
    mobile_section_class = f' class="{" ".join(section_classes)}"' if section_classes else ""
    if mobile_treatment == "photo_bleed":
        panel_class = "panel panel-rounded"
    else:
        panel_class = "panel"
    data_attrs = f' data-treatment="{mobile_treatment}"'
    if active_word:
        data_attrs += f' data-section="{_escape(active_word.upper())}"'
    return (
        '<section id="deck-desktop"><div class="canvas">\n'
        + "\n".join(desktop_parts)
        + "\n</div></section>\n"
        f'<section id="deck-mobile"{mobile_section_class}{data_attrs}>{mobile_prefix}<div class="{panel_class}">\n'
        + "\n".join(mobile_parts)
        + "\n</div></section>"
    )



def _render_paragraphs(paragraphs: list, slide_w_pt: float, *, mobile: bool, default_color: str = "000000") -> str:
    out = []
    # Some blank spacer paragraphs carry no size info of their own at all
    # (no defRPr, no whitespace run — seen both in this deck's bulleted
    # lists AND in plain body text boxes with no bullets at all, e.g. "The
    # Vision" — the pattern isn't bullet-specific, an earlier version of
    # this fix incorrectly scoped it that way and missed every non-bullet
    # occurrence). Falling back to a real size is much closer to the
    # source's intent than silently collapsing the gap — but it needs to
    # be a *locally* appropriate size, not one constant for the whole
    # frame: a frame routinely mixes a large header with much smaller
    # body/bullet text, and a single frame-wide default (the first sized
    # run found anywhere, which is often the header) inflated every blank
    # line to the header's size — including gaps between small bullet
    # items, which measurably overflowed the text box's own bounds as a
    # result. For each blank paragraph, use the nearest real paragraph's
    # own size instead — checking forward first (a spacer usually leads
    # into what follows), backward as a fallback.
    para_size_pt = []
    for para in paragraphs:
        size = None
        for r in para.runs:
            if r.style.size_pt:
                size = r.style.size_pt
                break
        para_size_pt.append(size)
    local_default_size_pt = [18.0] * len(paragraphs)
    for i in range(len(paragraphs)):
        found = None
        for j in range(i, len(paragraphs)):
            if para_size_pt[j]:
                found = para_size_pt[j]
                break
        if found is None:
            for j in range(i, -1, -1):
                if para_size_pt[j]:
                    found = para_size_pt[j]
                    break
        if found:
            local_default_size_pt[i] = found
    # Only treat a blank paragraph as an internal spacer (worth guessing a
    # height for) if a real, visible paragraph follows it later in the
    # frame. A blank paragraph with nothing after it is a trailing blank —
    # harmless dead space in the source, and guessing a height for it
    # risks shifting already-correct vertically-centered text.
    has_visible_after = [False] * len(paragraphs)
    seen_visible = False
    for i in range(len(paragraphs) - 1, -1, -1):
        has_visible_after[i] = seen_visible
        if any(r.text.strip() for r in paragraphs[i].runs):
            seen_visible = True
    # Mirrors has_visible_after, forward instead of backward: a blank
    # paragraph with nothing *before* it is a leading blank. A run of many
    # leading blanks (seen on this deck: 12+ before a caption starts) is a
    # manual "push the text down within a much taller box" trick — the
    # same category as the leading-tab horizontal trick already stripped
    # on mobile, just vertical. It only means anything against the fixed-
    # height desktop box it was tuned for; on mobile's natural top-to-
    # bottom flow it just produces a large, meaningless dead gap before
    # the real content.
    has_visible_before = [False] * len(paragraphs)
    seen_visible = False
    for i in range(len(paragraphs)):
        has_visible_before[i] = seen_visible
        if any(r.text.strip() for r in paragraphs[i].runs):
            seen_visible = True
    for idx, p in enumerate(paragraphs):
        text_align = _ALIGN_CSS.get(p.align, "left")
        if mobile and text_align == "justify":
            # justify relies on a wide box with many words per line to
            # distribute extra space unnoticeably. A narrow mobile column
            # wraps to just a few words per line, so the same rule stretches
            # gaps into obvious, ugly holes. Center reads cleanly instead,
            # and matches how a human would actually reflow this for mobile.
            text_align = "center"
        spans = []
        whitespace_only_size_pt = None
        pending_leading_ws = ""
        for r in p.runs:
            if not r.text.strip() and r.text != " ":
                # Whitespace-only run (tabs and/or repeated spaces used as a
                # manual spacer, e.g. this deck's leading-tab lines). It
                # contributes no visible span itself, but its declared size is
                # the best signal for how tall this "blank" line should
                # render if the whole paragraph turns out empty — remember
                # that. Separately, if this run's whitespace is a genuine
                # leading tab-indent, carry the literal characters forward
                # onto the next visible run instead of discarding them: some
                # of this deck's headers put the tabs in their own run
                # (distinct from the text run), rather than sharing one run
                # with the visible text the way others do — dropping this
                # run entirely was silently erasing the manual indent for
                # exactly that pattern.
                if whitespace_only_size_pt is None and r.style.size_pt:
                    whitespace_only_size_pt = r.style.size_pt
                if "\t" in r.text:
                    pending_leading_ws += r.text
                continue
            size_pt = calibrate_size_pt(r.style.typeface, r.style.size_pt, r.style.bold)
            if r.style.bold is not None:
                weight = 700 if r.style.bold else 400
            elif r.style.typeface:
                _, implied_weight = _split_typeface(r.style.typeface)
                weight = implied_weight if implied_weight is not None else 400
            else:
                weight = 400
            stretch = "condensed;" if (r.style.typeface and _is_condensed(r.style.typeface)) else ""
            # PowerPoint's real default for a run with no explicit color is
            # the theme's dk1 (usually black) — a run only ends up
            # colorless in the XML when the deck author relied on that
            # default rather than setting a color explicitly (seen on this
            # deck's "FLOOR LEVEL" badges over light photos; most other
            # runs in the deck do set an explicit white, which is why a
            # hardcoded white fallback looked right until a slide actually
            # exercised the true default). Falling back to a hardcoded
            # white was simply wrong, not just deck-specific.
            color = color_with_alpha(r.style.color_hex or default_color, r.style.color_alpha)
            if mobile:
                # Mobile: scale relative to canvas width in vw, matching the
                # desktop cqw math but against the viewport instead of the
                # (absent, on mobile) fixed-aspect canvas.
                size_css = f"{size_pt / slide_w_pt * 100 * 2.2:.2f}vw"
                spacing_css = (
                    f"{r.style.spacing_pt / slide_w_pt * 100 * 2.2:.2f}vw"
                    if r.style.spacing_pt else "normal"
                )
            else:
                size_css = f"{size_pt / slide_w_pt * 100:.3f}cqw"
                spacing_css = (
                    f"{r.style.spacing_pt / slide_w_pt * 100:.3f}cqw"
                    if r.style.spacing_pt else "normal"
                )
            style = (
                f"font-size:{size_css};font-weight:{weight};color:{color};"
                f"letter-spacing:{spacing_css};font-family:{_font_family_css(r.style.typeface)};"
                f"font-stretch:{stretch if stretch else 'normal;'}"
                + ("font-style:italic;" if r.style.italic else "")
                + ("text-decoration:underline;" if r.style.underline else "")
            )
            text = pending_leading_ws + r.text
            pending_leading_ws = ""
            if mobile and "\t" in text:
                # This run uses the manual tab+space "push right" trick
                # (tabs, sometimes followed by literal padding spaces,
                # before the real content) — only makes sense against a
                # wide desktop box. Strip all of it for mobile so the
                # real text just starts at the normal left margin instead
                # of wrapping oddly after eating a chunk of a narrow column.
                text = text.lstrip("\t ")
            spans.append(f'<span style="{style}">{_escape(text)}</span>')
        if spans:
            has_tab = any("\t" in r.text for r in p.runs)
            ws = "white-space:pre-wrap;" if has_tab else ""
            if p.bullet_char:
                # Hanging-indent bullet layout: the bullet sits in the
                # indent gap (negative text-indent pulls the first line's
                # bullet back to the margin), continuation lines wrap
                # flush with the text start (padding-left holds them
                # there) — the same technique word processors use, so it
                # matches the source's marL/indent semantics instead of
                # inventing a fixed layout unrelated to the real value.
                marL_pt = p.marL_pt or 27.0
                indent_css = f"{marL_pt / slide_w_pt * 100:.3f}cqw"
                bullet_size_pt = calibrate_size_pt(
                    p.runs[0].style.typeface, p.runs[0].style.size_pt, p.runs[0].style.bold
                ) if p.runs else 18.0
                bullet_size_css = f"{bullet_size_pt / slide_w_pt * 100:.3f}cqw"
                # Plain inline span, not inline-block: an inline-block bullet
                # combined with negative text-indent silently fails to render
                # at all in Chromium (verified directly) — plain inline
                # content with the same hanging-indent CSS renders correctly.
                bullet_html = (
                    f'<span style="font-size:{bullet_size_css};">'
                    f'{_escape(p.bullet_char)}&nbsp;&nbsp;</span>'
                )
                out.append(
                    f'<p style="text-align:{text_align};margin:0;{ws}'
                    f'padding-left:{indent_css};text-indent:-{indent_css};">'
                    f'{bullet_html}{"".join(spans)}</p>'
                )
            else:
                out.append(f'<p style="text-align:{text_align};margin:0;{ws}">{"".join(spans)}</p>')
        elif mobile and not has_visible_before[idx]:
            # Leading blank paragraph (nothing visible before it in this
            # frame) — on mobile, drop it entirely rather than render its
            # height. See has_visible_before's comment: this is a desktop-
            # only manual vertical-positioning trick, meaningless (and
            # actively harmful — a large dead gap before the real content)
            # once the box's real height/position no longer applies.
            pass
        elif p.blank_line_size_pt or whitespace_only_size_pt or has_visible_after[idx]:
            # Empty paragraph used purely as a manual vertical spacer
            # (leading-tabs horizontal trick, or between bulleted/plain
            # body paragraphs with no size info of their own at all). It
            # still needs to occupy real line height, or paragraphs after
            # it collapse upward and lose the gap the deck's author put
            # there on purpose. Only treated as a spacer worth guessing a
            # height for when has_visible_after[idx] is true (see above) —
            # a trailing blank with nothing real after it is frequently
            # harmless dead space, and guessing a height for it risks
            # shifting already-correct vertically-centered text.
            #
            # Deliberately NOT an explicit height with a hand-tuned
            # multiplier (an earlier version used *1.43, curve-fit against
            # one reference image — exactly the "compensate CSS to chase a
            # screenshot" anti-pattern this codebase's own LEARNINGS.md
            # rule #12 warns against; a constant fitted to one deck is not
            # trustworthy on the next one). Instead: render the blank line
            # as a real line of text at its real font-size — using the
            # *locally* appropriate size (local_default_size_pt), not one
            # constant for the whole frame — and let the browser's own
            # default line-height (the same mechanism every visible text
            # line already uses) determine its height.
            blank_pt = p.blank_line_size_pt or whitespace_only_size_pt or local_default_size_pt[idx]
            if mobile:
                fs_css = f"{blank_pt / slide_w_pt * 100 * 2.2:.2f}vw"
            else:
                fs_css = f"{blank_pt / slide_w_pt * 100:.3f}cqw"
            out.append(f'<p style="margin:0;font-size:{fs_css};">&nbsp;</p>')
    return "\n".join(out)


def _escape(s: str) -> str:
    return (
        s.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )
