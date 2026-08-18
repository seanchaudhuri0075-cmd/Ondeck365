"""Entry-point reader for a .pptx file.

Owns: opening the file, exposing parts (slides, theme, media), reading
slide dimensions. Does NOT own rendering, layout detection, or shape
walking — those modules consume the parts this module exposes.

Modules that need raw XML access (SVG-only blips, group transforms)
reach into `slide.element` themselves; this module deliberately stays
thin so it doesn't lock downstream code into python-pptx idioms.
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterator

import pptx as _pptx
from pptx.presentation import Presentation
from pptx.slide import Slide

EMU_PER_PT = 12700  # 1 point = 12700 English Metric Units (914400 EMU/inch ÷ 72 pt/inch).
# NOTE: 9525 was used here previously — that's EMU-per-PIXEL at 96 DPI
# (914400/96), a completely different conversion that was mislabeled as
# EMU-per-point. Shape x/y/w/h ratios against slide width/height were
# unaffected (numerator and denominator were equally wrong, so the ratio
# canceled out) — but anything mixing a true-unit value against slide
# width/height directly, e.g. font-size or letter-spacing (both declared
# in real hundredths-of-a-point, unrelated to EMU) expressed as a
# percentage of slide width, rendered 25% too small (960/1280) than
# intended. Confirmed against a real declared 48pt headline rendering
# visibly smaller than the source deck's actual export.


class Pptx:
    def __init__(self, path: str | Path):
        self.path = Path(path)
        if not self.path.exists():
            raise FileNotFoundError(self.path)
        self._prs: Presentation = _pptx.Presentation(str(self.path))

    @property
    def slide_size_pt(self) -> tuple[float, float]:
        """Slide width/height in points, read from <p:sldSz>.

        Per the spec: never hardcode dimensions. P&G is 1280x720pt;
        FrameTag and others differ.
        """
        return (
            self._prs.slide_width / EMU_PER_PT,
            self._prs.slide_height / EMU_PER_PT,
        )

    @property
    def slide_count(self) -> int:
        return len(self._prs.slides)

    @property
    def default_tab_pt(self) -> float:
        """Default tab-stop interval in points, from <p:defaultTextStyle defTabSz>.

        Governs how leading/embedded literal tab characters in run text
        advance — some decks use tabs (sometimes combined with literal
        space runs) to manually position text within a box instead of
        setting real alignment. Falls back to the OOXML spec default of
        1 inch (914400 EMU = 72pt) if the attribute is absent.
        """
        el = self._prs.element
        NS = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        default_style = el.find(f"{NS}defaultTextStyle")
        if default_style is not None:
            val = default_style.get("defTabSz")
            if val is not None:
                return int(val) / EMU_PER_PT
        return 914400 / EMU_PER_PT

    def slides(self) -> Iterator[Slide]:
        return iter(self._prs.slides)

    @property
    def raw(self) -> Presentation:
        """Escape hatch: the underlying python-pptx Presentation."""
        return self._prs
