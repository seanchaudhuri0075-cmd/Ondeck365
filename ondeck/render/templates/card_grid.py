"""card_grid template — title + chrome logo + N brand logos + optional bg photo + tint overlay.

Phase 1B Cohort B8 scope: slide 2 (title + 9 brand logos in a 4-deep nested
group + duotoned background photo + 55%-alpha cyan tint overlay).

Structure detection:
- Skip full-canvas TRULY OPAQUE rect matching `canvas_bg` (the rect is
  redundant when alpha is absent — the renderer's `--bg` paints the same
  pixels). A full-canvas rect WITH `<a:alpha>` is NOT skipped — it's a
  deck-author tint overlay sitting on top of background content (slide 2
  Rectangle 1 case: cyan #00B0F0 at alpha=55%).
- Identify a background photo: full-canvas-sized pic at z=0 (always the
  first leaf in flattener order). Slide 2 Picture 62 is duotoned via
  blip-level <a:duotone>+<a:alphaModFix> — both pre-baked into the
  on-disk WebP by transform/image.py before this template renders.
- The first <p:pic> matching `is_logo_pic` geometry is the chrome logo
  (top-right, ~127x46).
- All other pics are brand logos at their flattened world coords.
- All text shapes (post canvas-rect skip) are titles. Slide 2 has one.

Z-order in the rendered DOM (back to front):
  1. Background photo (.photo, full canvas)
  2. Tint overlay rect (.tint-overlay, full canvas, rgba())
  3. Chrome logo (.gif-logo)
  4. Title (.t-title)
  5. Brand logos (.photo-N for N=0..8)

Browser z-order follows source order for absolute-positioned siblings —
no explicit `z-index` needed.

Brand-logo transport: routes through `inline_optimized_data_url`
(transform/image.py-produced WebP, optimized for size). Inheriting the
photo_grid carve-out rule from media_showcase. PNG alpha (transparent
backgrounds on all 9 brand logos) survives the WebP transcode and
base64-inline pipeline end-to-end. Strict contract: caller must pass
`media_dir` so the helper can locate the WebP files.

Brand-logo source resolution is a known limitation: PPTX embeds them at
175-257px (designed for 1280px-canvas display, no 2-3x retina source).
At canvas widths above 1280px (any modern viewport, retina, 4K) the
logos visibly pixelate. Matches the original PPT bundle behavior; not
fixable in the renderer. Phase 2 work to procure higher-DPI brand assets
if greater fidelity is needed. See NOTES.md for the full rationale.

Mobile fallback (sensible defaults — iterate after visual diff):
- Top bar with chrome logo right-aligned (matches title_stats / two_column
  precedent for cyan slides with chrome).
- Title centered below.
- 9 brand logos in a 3-column CSS grid (3×3) at consistent cell size.
- Brand logos render with `width: 100%` of cell, object-fit: contain so
  their native aspect ratios are preserved without crop.
- Background photo NOT rendered on mobile — desktop-only design intent.
  (If a mobile bg becomes needed, route through a separate mobile slot
  with its own data URL so the desktop and mobile copies don't double-pay
  for the inline base64.)
"""
from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx.slide import Slide

from ..css import pt_to_cqw, color_with_alpha
from ..desktop import canvas_aspect_css, positioned_style
from ..html import render_page
from ...layout.detect import SlideClass
from ...parse.font_calibration import calibrate_size_pt
from ...parse.shapes import FlatShape, flatten_slide
from ...parse.slide import NS
from ...parse.text import parse_text_frame
from ...parse.theme import Theme
from ._shared import (
    TYPEFACE_WEIGHT,
    inline_data_url,
    inline_optimized_data_url,
    is_logo_pic,
)


def render_card_grid(
    slide: Slide,
    theme: Theme,
    slide_index: int,
    slide_class: SlideClass,
    deck_name: str,
    slide_w_pt: float,
    slide_h_pt: float,
    deck_brand_color: str = "#000000",
    media_dir: Optional[Path] = None,
) -> str:
    """Render a card_grid slide. Returns HTML string."""
    flat = list(flatten_slide(slide))
    canvas_bg = slide_class.hints.get("canvas_bg") or deck_brand_color or "#FFFFFF"
    logo_invert = bool(slide_class.hints.get("logo_invert", False))
    title_size_pt_hint = slide_class.hints.get("title_size_pt")
    mobile_title_size_vw = slide_class.hints.get("mobile_title_size_vw", 6.0)
    mobile_title_weight = slide_class.hints.get("mobile_title_weight", 500)
    mobile_title_line_height = slide_class.hints.get("mobile_title_line_height", 1.2)

    bg_photo, tint_overlay, chrome_logo, brand_logos, title_frame = _classify(
        flat, slide_w_pt, slide_h_pt, canvas_bg, theme, slide,
    )

    # Title size override — mirrors title_stats / two_column precedent:
    # hint wins over OOXML unconditionally (deck-author intent over the
    # source PPT's run-level sz). Applied even when OOXML resolved a value.
    if title_frame is not None and title_size_pt_hint is not None:
        for para in title_frame["paragraphs"]:
            for run in para["runs"]:
                run["size_pt"] = float(title_size_pt_hint)

    # Cluster brand logos into rows-of-clusters using world coords. Rows
    # group by y-center proximity; clusters within a row group by x-range
    # overlap. Slide 2: 9 logos → row 1 (5 logos in 3 clusters: Secret +
    # Secret Clinical lockup + Secret Whole Body lockup) + row 2 (4 logos
    # in 4 clusters: Always / Gillette / Olay / Ivory). Empty when no
    # brand logos.
    brand_rows = _cluster_logos_by_row(brand_logos)

    # URL resolution. Mobile and desktop reuse the same data URLs (the
    # base64 string is identical → validator's set-based unique-source
    # count collapses correctly; HTML inflation is just the duplicate
    # base64 strings, which gzip well in transit).
    bg_photo_url = (
        inline_optimized_data_url(bg_photo, slide_class, slide, media_dir)
        if bg_photo is not None else None
    )
    chrome_logo_url = inline_data_url(chrome_logo, slide) if chrome_logo is not None else None
    brand_logo_urls = [
        inline_optimized_data_url(p, slide_class, slide, media_dir)
        for p in brand_logos
    ]
    # Per-shape URL map for the cluster mobile emit (URL lookup by
    # shape identity since clusters reorder logos within a row).
    brand_url_by_shape = {
        id(p): u for p, u in zip(brand_logos, brand_logo_urls)
    }

    slide_css = _build_css(
        slide_w_pt, slide_h_pt,
        bg_photo, tint_overlay, chrome_logo, brand_logos, title_frame,
        brand_rows, logo_invert,
        mobile_title_size_vw=mobile_title_size_vw,
        mobile_title_weight=mobile_title_weight,
        mobile_title_line_height=mobile_title_line_height,
    )
    body_html = _build_body(
        bg_photo_url, tint_overlay, chrome_logo_url, brand_logo_urls, title_frame,
        brand_rows, brand_url_by_shape,
    )

    return render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg":        canvas_bg,
            "bg-cyan":   deck_brand_color or "#000000",
            "headline":  "#FFFFFF",
            "font-cond": '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=body_html,
        slide_css=slide_css,
    )


# ────────────────────────────────────────────────────────────────────────────
# Classification

def _classify(flat, slide_w_pt, slide_h_pt, canvas_bg, theme, slide):
    """Return (bg_photo_pic_or_None, tint_overlay_dict_or_None,
    chrome_logo_pic_or_None, brand_logo_pics, title_frame_dict_or_None).

    bg_photo: a full-canvas pic at z=0 (always the first leaf in
        flattener order) — the deck-author background photo. Slide 2
        Picture 62 case. None when no such pic exists.
    tint_overlay: {"r": int, "g": int, "b": int, "alpha": float} for a
        full-canvas semi-transparent rect matching canvas_bg's color.
        Slide 2 Rectangle 1 case (cyan #00B0F0 at alpha=0.55). None when
        no such overlay exists.
    chrome_logo: a logo-geometry pic (top-right) per is_logo_pic.
    brand_logos: all other pics, preserved in flattener (z-order) sequence
        so the rendered HTML reflects deck-author intent for overlapping
        pics.
    title_frame: {"shape": FlatShape, "paragraphs": [...]} or None.

    Truly-opaque full-canvas rects matching canvas_bg are silently skipped
    (redundant with the renderer's --bg paint). See _is_canvas_skip_rect.
    """
    bg_photo = None
    tint_overlay = None
    chrome_logo = None
    brand_logos = []
    title_frame = None
    for s in flat:
        if s.x_pt is None:
            continue
        if s.kind == "pic":
            # z=0 + full-canvas geometry → background photo. The z=0 check
            # is what differentiates an intentional full-bleed bg from a
            # large foreground pic that happens to be near canvas size.
            if (
                bg_photo is None
                and s.z == 0
                and _is_full_canvas(s, slide_w_pt, slide_h_pt)
            ):
                bg_photo = s
                continue
            if chrome_logo is None and is_logo_pic(s, slide_w_pt, slide_h_pt):
                chrome_logo = s
                continue
            brand_logos.append(s)
            continue
        if s.kind != "sp":
            continue
        # Truly opaque full-canvas rect matching canvas_bg → redundant skip.
        if _is_canvas_skip_rect(s, slide_w_pt, slide_h_pt, canvas_bg):
            continue
        # Full-canvas rect matching canvas_bg WITH alpha → tint overlay.
        if tint_overlay is None:
            ov = _read_tint_overlay(s, slide_w_pt, slide_h_pt, canvas_bg)
            if ov is not None:
                tint_overlay = ov
                continue
        paras = _extract_paragraphs(s, theme)
        if not paras or not any(p["runs"] for p in paras):
            continue
        if title_frame is None:
            title_frame = {"shape": s, "paragraphs": paras}
        # Additional text shapes silently dropped — slide 2 has only one;
        # extend if future card_grid slides add subtitle / footnote frames.

    return bg_photo, tint_overlay, chrome_logo, brand_logos, title_frame


def _is_full_canvas(s: FlatShape, slide_w_pt, slide_h_pt) -> bool:
    """True if shape covers ≥90% of canvas in both dimensions."""
    if s.w_pt is None or s.h_pt is None:
        return False
    return (
        s.w_pt >= slide_w_pt * 0.9
        and s.h_pt >= slide_h_pt * 0.9
    )


def _is_canvas_skip_rect(s: FlatShape, slide_w_pt, slide_h_pt, canvas_bg) -> bool:
    """True if shape is a full-canvas TRULY OPAQUE solid rect matching canvas_bg.

    Same alpha-aware heuristic as two_column._is_canvas_skip_rect. A rect
    with `<a:alpha>` inside its srgbClr is NOT skipped here — it routes
    through `_read_tint_overlay` and emits as a semi-transparent overlay.
    """
    if abs(s.x_pt) > 1 or abs(s.y_pt) > 1:
        return False
    if abs(s.w_pt - slide_w_pt) > 1 or abs(s.h_pt - slide_h_pt) > 1:
        return False
    prst = s.element.find("p:spPr/a:prstGeom", NS)
    if prst is None or prst.get("prst") != "rect":
        return False
    target = canvas_bg.lstrip("#").upper()
    sr = s.element.find("p:spPr/a:solidFill/a:srgbClr", NS)
    if sr is None or sr.get("val", "").upper() != target:
        return False
    if sr.find("a:alpha", NS) is not None:
        return False
    return True


def _read_tint_overlay(s: FlatShape, slide_w_pt, slide_h_pt, canvas_bg) -> Optional[dict]:
    """Return {"r","g","b","alpha"} for a full-canvas alpha-bearing rect.

    None when the shape isn't a full-canvas rect, isn't matching
    canvas_bg's color, or has no alpha (truly opaque overlays are
    redundant — handled via _is_canvas_skip_rect's skip path).

    Currently constrained to overlays whose color matches canvas_bg
    (slide 2 Rectangle 1 case). If a future slide overlays a different
    color (e.g., black scrim), generalize by relaxing the color check.
    """
    if abs(s.x_pt) > 1 or abs(s.y_pt) > 1:
        return None
    if abs(s.w_pt - slide_w_pt) > 1 or abs(s.h_pt - slide_h_pt) > 1:
        return None
    prst = s.element.find("p:spPr/a:prstGeom", NS)
    if prst is None or prst.get("prst") != "rect":
        return None
    target = canvas_bg.lstrip("#").upper()
    sr = s.element.find("p:spPr/a:solidFill/a:srgbClr", NS)
    if sr is None or sr.get("val", "").upper() != target:
        return None
    alpha_node = sr.find("a:alpha", NS)
    if alpha_node is None:
        return None
    try:
        alpha = max(0.0, min(1.0, int(alpha_node.get("val", "100000")) / 100000.0))
    except ValueError:
        return None
    val = sr.get("val", "000000")
    try:
        r = int(val[0:2], 16); g = int(val[2:4], 16); b = int(val[4:6], 16)
    except ValueError:
        r, g, b = 0, 0, 0
    return {"r": r, "g": g, "b": b, "alpha": alpha}


# ────────────────────────────────────────────────────────────────────────────
# Brand-logo clustering — 2-stage: rows by y-center proximity, then clusters
# within each row by x-range overlap. Encodes deck-author lockup grouping for
# mobile emission so visual cohesion (e.g. "Secret + sub-element variants in
# row 1, separate brand cards in row 2") survives the desktop→mobile reflow.

# Threshold for splitting rows by y-center gap. 80pt at the 720pt canvas
# height (~11% of canvas) is wider than typical intra-row y-jitter (logos
# in the same visual row vary by ~0-50pt depending on lockup design) and
# narrower than typical inter-row gap (rows separate by 100pt+). Slide 2:
# row-1 y-centers span 228-318 (intra-row delta=90, all clustered);
# row-2 y-centers all 459 (delta=0); inter-row gap=141 → split at 80.
_ROW_GAP_THRESHOLD_PT = 80.0


def _cluster_logos_by_row(logos: list) -> list:
    """Group logos into rows-of-clusters by world-coord proximity.

    Two-stage clustering:
      1. Sort by y-center. Walk the sorted list, starting a new row
         whenever the gap between consecutive y-centers exceeds
         _ROW_GAP_THRESHOLD_PT. Slide 2: produces 2 rows.
      2. Within each row, sort by x-left. Cluster consecutive logos
         whose x-ranges overlap (next.x_left < running_max_x_right).
         Slide 2 row 1: 3 clusters {P43} | {P39, P44} | {P45, P47};
         row 2: 4 clusters (no x-overlap among the 4 brand cards).

    Returns list[list[list[FlatShape]]]: rows → clusters → logos.
    Source order is preserved within clusters so the lockup-base logo
    (typically larger, lower z-index in source) appears before lockup
    overlay sub-elements in the rendered DOM.
    """
    if not logos:
        return []

    # Stage 1: rows by y-center
    by_y_center = sorted(
        logos, key=lambda l: l.y_pt + (l.h_pt or 0) / 2,
    )
    rows: list[list] = []
    current_row: list = []
    last_y_center = None
    for logo in by_y_center:
        y_center = logo.y_pt + (logo.h_pt or 0) / 2
        if last_y_center is not None and (y_center - last_y_center) > _ROW_GAP_THRESHOLD_PT:
            if current_row:
                rows.append(current_row)
            current_row = []
        current_row.append(logo)
        last_y_center = y_center
    if current_row:
        rows.append(current_row)

    # Stage 2: clusters within each row by x-range overlap
    clustered = []
    for row in rows:
        by_x = sorted(row, key=lambda l: l.x_pt)
        clusters: list[list] = []
        current_cluster: list = []
        running_x_right = None
        for logo in by_x:
            x_left = logo.x_pt
            x_right = logo.x_pt + (logo.w_pt or 0)
            if running_x_right is not None and x_left >= running_x_right:
                if current_cluster:
                    clusters.append(current_cluster)
                current_cluster = []
                running_x_right = x_right
            else:
                running_x_right = max(running_x_right or x_right, x_right)
            current_cluster.append(logo)
        if current_cluster:
            clusters.append(current_cluster)
        clustered.append(clusters)
    return clustered


def _cluster_bbox(cluster: list) -> tuple:
    """Return (x_min, y_min, x_max, y_max) bounding box for a cluster.

    Used by the cluster mobile emit to derive each logo's position-percent
    within the cluster cell (preserves the desktop overlap relationship
    e.g. for the Secret Clinical lockup where Picture 44 sits above-right
    of Picture 39 in the same x-range).
    """
    x_min = min(l.x_pt for l in cluster)
    y_min = min(l.y_pt for l in cluster)
    x_max = max(l.x_pt + (l.w_pt or 0) for l in cluster)
    y_max = max(l.y_pt + (l.h_pt or 0) for l in cluster)
    return (x_min, y_min, x_max, y_max)


# ────────────────────────────────────────────────────────────────────────────
# Run / paragraph extraction (single-paragraph, single-run text frames only —
# extend if future card_grid slides need multi-line title support)

def _extract_paragraphs(text_sp: FlatShape, theme: Theme) -> list:
    tf = parse_text_frame(text_sp.element, theme)
    if tf is None:
        return []
    out = []
    for p in tf.paragraphs:
        runs = list(p.runs)
        runs = [r for r in runs if r.text.strip()]
        if not runs:
            out.append({"algn": p.align, "runs": []})
            continue
        first = runs[0]
        weight = TYPEFACE_WEIGHT.get((first.style.typeface or "").lower(), 500)
        if first.style.bold:
            weight = 700
        size_pt = calibrate_size_pt(
            first.style.typeface, first.style.size_pt, first.style.bold,
        )
        color = color_with_alpha(
            first.style.color_hex or "FFFFFF", first.style.color_alpha,
        )
        text = "".join(r.text for r in runs).strip()
        out.append({
            "algn": p.align,
            "runs": [{
                "text": text,
                "size_pt": size_pt,
                "weight": weight,
                "color": color,
                "bold": bool(first.style.bold),
            }],
        })
    return out


# ────────────────────────────────────────────────────────────────────────────
# CSS

def _build_css(slide_w_pt, slide_h_pt, bg_photo, tint_overlay,
               chrome_logo, brand_logos, title_frame,
               brand_rows: list,
               logo_invert: bool,
               mobile_title_size_vw: float = 6.0,
               mobile_title_weight: int = 500,
               mobile_title_line_height: float = 1.2) -> str:
    canvas_w_css, canvas_h_css = canvas_aspect_css(slide_w_pt, slide_h_pt)
    parts = []

    parts.append("/* ---------------- DESKTOP ---------------- */")
    parts.append("#deck-desktop { display: block; }")
    parts.append(
        f"#deck-desktop .canvas {{\n"
        f"  position: relative;\n"
        f"  width: {canvas_w_css};\n"
        f"  height: {canvas_h_css};\n"
        f"  margin: 0 auto;\n"
        f"  container-type: inline-size;\n"
        f"  overflow: hidden;\n"
        f"  background: var(--bg);\n"
        f"}}"
    )
    parts.append("#deck-desktop .slide { position: absolute; inset: 0; }")

    if bg_photo is not None:
        parts.append(
            f"#deck-desktop .photo {{\n"
            f"  {positioned_style(bg_photo.x_pt, bg_photo.y_pt, bg_photo.w_pt, bg_photo.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"}}\n"
            f"#deck-desktop .photo img {{ width: 100%; height: 100%; object-fit: cover; }}"
        )

    if tint_overlay is not None:
        rgba = (
            f"rgba({tint_overlay['r']}, {tint_overlay['g']}, "
            f"{tint_overlay['b']}, {tint_overlay['alpha']:.3f})"
        )
        parts.append(
            f"#deck-desktop .tint-overlay {{\n"
            f"  position: absolute; inset: 0;\n"
            f"  background-color: {rgba};\n"
            f"}}"
        )

    if chrome_logo is not None:
        logo_filter = "  filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            f"#deck-desktop .gif-logo {{\n"
            f"  {positioned_style(chrome_logo.x_pt, chrome_logo.y_pt, chrome_logo.w_pt, chrome_logo.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"{logo_filter}"
            f"}}\n"
            f"#deck-desktop .gif-logo img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )

    if title_frame is not None:
        s = title_frame["shape"]
        first_run = next((r for p in title_frame["paragraphs"] for r in p["runs"]), None)
        algn = next((p["algn"] for p in title_frame["paragraphs"] if p["runs"]), None)
        text_align = {"ctr": "center", "r": "right", "just": "justify"}.get(algn, "left")
        size_css = (
            pt_to_cqw(first_run["size_pt"], slide_w_pt)
            if first_run and first_run["size_pt"]
            else "1.875cqw"  # 24pt at 1280pt canvas fallback
        )
        weight = first_run["weight"] if first_run else 500
        color = first_run["color"] if first_run else "#FFFFFF"
        parts.append(
            f"#deck-desktop .t-title {{\n"
            f"  {positioned_style(s.x_pt, s.y_pt, s.w_pt, s.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  font-size: {size_css};\n"
            f"  font-weight: {weight};\n"
            f"  line-height: 1.2;\n"
            f"  color: {color};\n"
            f"  text-align: {text_align};\n"
            f"  letter-spacing: 0.005em;\n"
            f"  margin: 0;\n"
            f"  display: flex; align-items: center; justify-content: {text_align};\n"
            f"}}"
        )

    for i, p in enumerate(brand_logos):
        parts.append(
            f"#deck-desktop .photo-{i} {{\n"
            f"  {positioned_style(p.x_pt, p.y_pt, p.w_pt, p.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"}}\n"
            f"#deck-desktop .photo-{i} img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )

    # Mobile: bg photo + tint overlay (full-panel) under top-bar chrome,
    # title, and a 2-row stack of brand-logo clusters (preserves desktop
    # row-y grouping; clusters within each row preserve x-overlap lockups).
    parts.append("\n/* ---------------- MOBILE ---------------- */")
    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    parts.append("  #deck-desktop { display: none; }")
    parts.append("  #deck-mobile { display: block; background: var(--bg); }")
    # Panel becomes position:relative so absolute-positioned bg + overlay
    # sit beneath the natural flow content. Content inside top-bar /
    # title / brand-rows declares position:relative so they layer above
    # the inset:0 absolutes via stacking context (no z-index needed —
    # source order wins for siblings within the same stacking context).
    parts.append(
        "  #deck-mobile .panel {\n"
        "    min-height: 100vh;\n"
        "    color: #FFFFFF;\n"
        "    padding: 0 0 12vw;\n"
        "    position: relative;\n"
        "    overflow: hidden;\n"
        "  }"
    )
    if bg_photo is not None:
        parts.append(
            "  #deck-mobile .photo-mobile {\n"
            "    position: absolute; inset: 0;\n"
            "  }\n"
            "  #deck-mobile .photo-mobile img {\n"
            "    width: 100%; height: 100%; object-fit: cover;\n"
            "  }"
        )
    if tint_overlay is not None:
        rgba = (
            f"rgba({tint_overlay['r']}, {tint_overlay['g']}, "
            f"{tint_overlay['b']}, {tint_overlay['alpha']:.3f})"
        )
        parts.append(
            f"  #deck-mobile .tint-overlay-mobile {{\n"
            f"    position: absolute; inset: 0;\n"
            f"    background-color: {rgba};\n"
            f"  }}"
        )
    if chrome_logo is not None:
        parts.append(
            "  #deck-mobile .top-bar {\n"
            "    display: flex; justify-content: flex-end;\n"
            "    padding: 5vw;\n"
            "    position: relative;\n"
            "  }\n"
            "  #deck-mobile .gif-logo-mobile {\n"
            "    width: 22vw;\n"
            "    filter: brightness(0) invert(1);\n"
            "  }\n"
            "  #deck-mobile .gif-logo-mobile img { width: 100%; }"
        )
    if title_frame is not None:
        parts.append(
            f"  #deck-mobile .title-mobile {{\n"
            f"    padding: 6vw 6vw 8vw;\n"
            f"    text-align: center;\n"
            f"    font-size: {mobile_title_size_vw}vw;\n"
            f"    font-weight: {mobile_title_weight};\n"
            f"    line-height: {mobile_title_line_height};\n"
            f"    color: #FFFFFF;\n"
            f"    letter-spacing: 0.01em;\n"
            f"    margin: 0;\n"
            f"    position: relative;\n"
            f"  }}"
        )
    if brand_rows:
        # Per-row aspect ratio derived from cluster bounding boxes in the
        # row — keeps multi-img clusters (e.g. Secret Clinical lockup with
        # P39 + P44 overlapping) at their desktop relative shape on mobile.
        # Each cluster cell uses position:relative + per-img absolute
        # positioning to faithfully reproduce desktop overlap geometry.
        parts.append(
            "  #deck-mobile .brand-rows-mobile {\n"
            "    display: flex; flex-direction: column;\n"
            "    gap: 6vw;\n"
            "    padding: 0 6vw;\n"
            "    position: relative;\n"
            "  }\n"
            "  #deck-mobile .brand-row-mobile {\n"
            "    display: flex; align-items: center;\n"
            "    gap: 4vw;\n"
            "  }\n"
            "  #deck-mobile .brand-cluster-mobile {\n"
            "    flex: 1 1 0;\n"
            "    position: relative;\n"
            "  }\n"
            "  #deck-mobile .brand-cluster-mobile img {\n"
            "    position: absolute;\n"
            "    object-fit: contain;\n"
            "  }"
        )
        # Per-row aspect ratios + per-cluster aspect ratios + per-img
        # within-cluster positioning. Computed from world coords so the
        # mobile layout proportionally matches desktop without manual
        # tuning. Indexed by row + cluster + img position so future deck
        # slides can use the same pattern with different geometries.
        for row_idx, clusters in enumerate(brand_rows):
            # Row aspect ratio = sum(cluster widths) : tallest cluster
            # height (cluster widths normalized so they share a baseline).
            for cluster_idx, cluster in enumerate(clusters):
                bx0, by0, bx1, by1 = _cluster_bbox(cluster)
                cw = bx1 - bx0
                ch = by1 - by0
                aspect = (cw / ch) if ch else 1.0
                parts.append(
                    f"  #deck-mobile .brand-row-mobile-{row_idx} "
                    f".brand-cluster-mobile-{cluster_idx} {{\n"
                    f"    aspect-ratio: {aspect:.4f};\n"
                    f"  }}"
                )
                for img_idx, logo in enumerate(cluster):
                    left_pct = (logo.x_pt - bx0) / cw * 100 if cw else 0
                    top_pct = (logo.y_pt - by0) / ch * 100 if ch else 0
                    width_pct = (logo.w_pt or 0) / cw * 100 if cw else 100
                    height_pct = (logo.h_pt or 0) / ch * 100 if ch else 100
                    parts.append(
                        f"  #deck-mobile .brand-row-mobile-{row_idx} "
                        f".brand-cluster-mobile-{cluster_idx} "
                        f".brand-cluster-img-{img_idx} {{\n"
                        f"    left: {left_pct:.4f}%; top: {top_pct:.4f}%;\n"
                        f"    width: {width_pct:.4f}%; height: {height_pct:.4f}%;\n"
                        f"  }}"
                    )

    parts.append("}")
    return "\n".join(parts)


# ────────────────────────────────────────────────────────────────────────────
# Body HTML

def _build_body(bg_photo_url, tint_overlay, chrome_logo_url, brand_logo_urls, title_frame,
                brand_rows: list, brand_url_by_shape: dict) -> str:
    desktop_lines = []
    if bg_photo_url:
        desktop_lines.append(f'      <div class="photo"><img src="{bg_photo_url}" alt=""></div>')
    if tint_overlay is not None:
        desktop_lines.append('      <div class="tint-overlay"></div>')
    if chrome_logo_url:
        desktop_lines.append(f'      <div class="gif-logo"><img src="{chrome_logo_url}" alt=""></div>')
    if title_frame is not None:
        for para in title_frame["paragraphs"]:
            for run in para["runs"]:
                desktop_lines.append(
                    f'      <div class="t-title">{_escape(run["text"])}</div>'
                )
    for i, url in enumerate(brand_logo_urls):
        if url:
            desktop_lines.append(f'      <div class="photo-{i}"><img src="{url}" alt=""></div>')

    mobile_lines = []
    # Mobile bg + overlay: same z-stack pattern as desktop. Shares the
    # data URLs with the desktop emit (validator's set-based unique-src
    # count collapses correctly).
    if bg_photo_url:
        mobile_lines.append(f'    <div class="photo-mobile"><img src="{bg_photo_url}" alt=""></div>')
    if tint_overlay is not None:
        mobile_lines.append('    <div class="tint-overlay-mobile"></div>')
    if chrome_logo_url:
        mobile_lines.append('    <div class="top-bar">')
        mobile_lines.append(f'      <div class="gif-logo-mobile"><img src="{chrome_logo_url}" alt=""></div>')
        mobile_lines.append('    </div>')
    if title_frame is not None:
        for para in title_frame["paragraphs"]:
            for run in para["runs"]:
                mobile_lines.append(
                    f'    <div class="title-mobile">{_escape(run["text"])}</div>'
                )
    if brand_rows:
        mobile_lines.append('    <div class="brand-rows-mobile">')
        for row_idx, clusters in enumerate(brand_rows):
            mobile_lines.append(
                f'      <div class="brand-row-mobile brand-row-mobile-{row_idx}">'
            )
            for cluster_idx, cluster in enumerate(clusters):
                mobile_lines.append(
                    f'        <div class="brand-cluster-mobile '
                    f'brand-cluster-mobile-{cluster_idx}">'
                )
                for img_idx, logo in enumerate(cluster):
                    url = brand_url_by_shape.get(id(logo))
                    if url:
                        mobile_lines.append(
                            f'          <img class="brand-cluster-img-{img_idx}" '
                            f'src="{url}" alt="">'
                        )
                mobile_lines.append('        </div>')
            mobile_lines.append('      </div>')
        mobile_lines.append('    </div>')

    return (
        "<!-- DESKTOP -->\n"
        "<section id=\"deck-desktop\">\n"
        "  <div class=\"canvas\">\n"
        "    <div class=\"slide\">\n"
        f"{chr(10).join(desktop_lines)}\n"
        "    </div>\n"
        "  </div>\n"
        "</section>\n\n"
        "<!-- MOBILE -->\n"
        "<section id=\"deck-mobile\">\n"
        "  <div class=\"panel\">\n"
        f"{chr(10).join(mobile_lines)}\n"
        "  </div>\n"
        "</section>"
    )


def _escape(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
