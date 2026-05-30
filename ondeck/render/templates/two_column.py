"""two_column template — title + N adjacent text columns + optional chrome logo.

Phase 1B Cohort B6 scope: slide 22 (title + 2 name columns, 30 names total).
Slide 23 is also tagged template=two_column but has a different internal
layout (Identity SVG + asymmetric contact info + email mailto: hyperlinks);
it'll route through this template too but may need additional handling.
Defer slide 23 to B7.

Structure detection:
- Skip full-canvas rectangle whose fill matches `canvas_bg` (deck-author
  chrome that's redundant once the renderer paints `--bg`).
- The first <p:pic> matching `is_logo_pic` geometry is the chrome logo
  (top-right). Inline as data URL (logos are tiny).
- All other text-bearing <p:sp> shapes are content frames. Top-most
  (smallest y) = title; remaining = body columns sorted left-to-right.

Hint behavior (mirrors title_stats.py precedent):
- `title_size_pt`: when set, OVERRIDES every run in the title frame.
- `body_size_pt`: when set, fills inherited (`size_pt is None`) runs in body
  columns. Explicit per-run sizes are preserved.
- Without hints, body inherits via resolve_inherited_size which walks to
  the master's <p:otherStyle> for txBox="1" frames (slide 22 case).

Run-level rules (hardcoded, no per-slide hints):
- Concatenate all runs within a paragraph into a single string (no <br>
  between runs — multi-run paragraphs are author-side text splits, not
  formatting changes).
- Strip trailing whitespace from the concatenated paragraph (handles the
  trailing whitespace-only runs in slide 22 right col paragraphs 2 + 4).
- Preserve U+00A0 (non-breaking space) verbatim — author intent (e.g.
  "Angel\\xa0 Bui" on slide 22).
- When <a:pPr algn> is None on a body paragraph, emit `text-align: left`
  in CSS rather than relying on browser default (paragraph alignment is
  master-inherited by OOXML; the master defaults to left, but we want
  the rendered HTML to be self-evident).
"""
from __future__ import annotations

from typing import Optional

from pptx.slide import Slide

from ..css import pt_to_pct_x, pt_to_pct_y, pt_to_cqw, color_with_alpha
from ..desktop import canvas_aspect_css, positioned_style
from ..html import render_page
from ...layout.detect import SlideClass
from ...parse.font_calibration import calibrate_size_pt
from ...parse.shapes import FlatShape, flatten_slide
from ...parse.slide import NS
from ...parse.text import parse_text_frame
from ...parse.theme import Theme
from ._shared import (
    TYPEFACE_WEIGHT,
    inline_data_url,
    is_logo_pic,
    resolve_inherited_size,
)
from ...parse.svg import has_svg_blip


def render_two_column(
    slide: Slide,
    theme: Theme,
    slide_index: int,
    slide_class: SlideClass,
    deck_name: str,
    slide_w_pt: float,
    slide_h_pt: float,
    deck_brand_color: str = "#000000",
) -> str:
    """Render a two_column slide. Returns HTML string."""
    flat = list(flatten_slide(slide))
    canvas_bg = slide_class.hints.get("canvas_bg") or deck_brand_color or "#FFFFFF"
    layout = slide_class.hints.get("two_column_layout", "stakeholders")
    if layout == "contact":
        return _render_contact_variant(
            flat, slide, theme, slide_index, slide_class,
            deck_name, slide_w_pt, slide_h_pt, canvas_bg, deck_brand_color,
        )
    title_size_pt_hint = slide_class.hints.get("title_size_pt")
    title_weight_hint = slide_class.hints.get("title_weight")
    body_size_pt_hint = slide_class.hints.get("body_size_pt")
    mobile_title_size_vw = slide_class.hints.get("mobile_title_size_vw", 6.0)
    mobile_title_weight = slide_class.hints.get("mobile_title_weight", 700)
    mobile_title_line_height = slide_class.hints.get("mobile_title_line_height", 1.0)
    logo_invert = bool(slide_class.hints.get("logo_invert", False))

    logo_pic, text_frames = _classify(
        flat, slide_w_pt, slide_h_pt, canvas_bg, theme, slide,
    )

    title_frame = text_frames[0] if text_frames else None
    body_columns = text_frames[1:] if len(text_frames) > 1 else []

    # Apply hint overrides (matches title_stats precedent: hint wins over OOXML
    # unconditionally — applied even when the OOXML resolution produced a
    # value, since the hint is the deck-author's design intent override).
    if title_frame is not None and title_size_pt_hint is not None:
        for para in title_frame["paragraphs"]:
            for run in para["runs"]:
                run["size_pt"] = float(title_size_pt_hint)
    if title_frame is not None and title_weight_hint is not None:
        for para in title_frame["paragraphs"]:
            for run in para["runs"]:
                run["weight"] = int(title_weight_hint)
    if body_size_pt_hint is not None:
        for col in body_columns:
            for para in col["paragraphs"]:
                for run in para["runs"]:
                    run["size_pt"] = float(body_size_pt_hint)

    logo_url = inline_data_url(logo_pic, slide) if logo_pic is not None else None

    slide_css = _build_css(
        slide_w_pt, slide_h_pt,
        title_frame, body_columns, logo_pic,
        logo_invert,
        mobile_title_size_vw=mobile_title_size_vw,
        mobile_title_weight=mobile_title_weight,
        mobile_title_line_height=mobile_title_line_height,
    )
    body_html = _build_body(title_frame, body_columns, logo_url)

    return render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg":        canvas_bg,
            "bg-cyan":   deck_brand_color or "#000000",
            "headline":  "#FFFFFF",
            "font-cond": '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=body_html,
        slide_css=slide_css,
    )


# ────────────────────────────────────────────────────────────────────────────
# Classification

def _classify(flat, slide_w_pt, slide_h_pt, canvas_bg, theme, slide):
    """Return (logo_pic_or_None, text_frames_sorted_by_y_then_x).

    text_frames is a list of dicts:
      {"shape": FlatShape, "paragraphs": [{"algn": str|None, "runs": [...]}]}
    """
    logo_pic = None
    text_frames = []
    for s in flat:
        if s.x_pt is None:
            continue
        if s.kind == "pic":
            if logo_pic is None and is_logo_pic(s, slide_w_pt, slide_h_pt):
                logo_pic = s
            # Other pics in two_column slides handled when they appear
            # (slide 23 will need the SVG hero — separate B7 work).
            continue
        if s.kind != "sp":
            continue
        # Skip full-canvas redundant background rect
        if _is_canvas_skip_rect(s, slide_w_pt, slide_h_pt, canvas_bg):
            continue
        paras = _extract_paragraphs(s, theme, slide)
        if not paras or not any(p["runs"] for p in paras):
            continue
        text_frames.append({"shape": s, "paragraphs": paras})

    text_frames.sort(key=lambda f: (round(f["shape"].y_pt), round(f["shape"].x_pt)))
    return logo_pic, text_frames


def _is_canvas_skip_rect(s: FlatShape, slide_w_pt, slide_h_pt, canvas_bg) -> bool:
    """True if shape is a full-canvas TRULY OPAQUE solid rect matching canvas_bg.

    Rectangle 5 on slides 22, 23 is exactly this — a deck-author cyan
    overlay that would visually duplicate `--bg`. Skipping at classify
    time means the renderer doesn't need a guard in the body emitter.

    A rect with `<a:alpha>` inside its srgbClr is NOT skipped: it's a
    semi-transparent overlay (slide 2 Rectangle 1 case — cyan with
    alpha=55000 = 55% opacity sitting on top of a duotoned background
    photo). Skipping it would erase the deck-author's tint pass.
    """
    if abs(s.x_pt) > 1 or abs(s.y_pt) > 1:
        return False
    if abs(s.w_pt - slide_w_pt) > 1 or abs(s.h_pt - slide_h_pt) > 1:
        return False
    prst = s.element.find("p:spPr/a:prstGeom", NS)
    if prst is None or prst.get("prst") != "rect":
        return False
    target = canvas_bg.lstrip("#").upper()
    sr = s.element.find("p:spPr/a:solidFill/a:srgbClr", NS)
    if sr is None or sr.get("val", "").upper() != target:
        return False
    # Color matches canvas_bg — but only skip if truly opaque.
    if sr.find("a:alpha", NS) is not None:
        return False
    return True


# ────────────────────────────────────────────────────────────────────────────
# Run / paragraph extraction

def _extract_paragraphs(text_sp: FlatShape, theme: Theme, slide: Slide) -> list:
    """Return [{"algn": str|None, "runs": [...]}] for the shape's text body.

    Each run dict: {text, size_pt, weight, color, bold}. size_pt is the
    declared pt (None when inherited; caller resolves via hints or
    resolve_inherited_size). Trailing whitespace-only runs in a paragraph
    are dropped.
    """
    tf = parse_text_frame(text_sp.element, theme)
    if tf is None:
        return []
    inherited_size = resolve_inherited_size(text_sp, slide)
    out = []
    for p in tf.paragraphs:
        # Drop trailing whitespace-only runs (preserves interior whitespace runs
        # such as "Rigbe" + " " + "Zemicael" — only strips dangling " ").
        runs = list(p.runs)
        while runs and not runs[-1].text.strip():
            runs.pop()
        if not runs:
            out.append({"algn": p.align, "runs": []})
            continue
        # Concatenate and rstrip the joined string to handle paragraphs
        # whose final non-empty run carries trailing whitespace ("Damilola "
        # + "Oyebanjo" + " " becomes "Damilola Oyebanjo" after the last "
        # " run is dropped above; but if the final non-empty run itself ends
        # in whitespace, rstrip catches it).
        concat = "".join(r.text for r in runs).rstrip()
        if not concat:
            out.append({"algn": p.align, "runs": []})
            continue

        # Inherit styling from the first non-whitespace run for the concatenated string.
        # Multi-run paragraphs in this deck always share styling across runs
        # (author-side text splits, not formatting changes), so taking the first
        # run's style is sound.
        first = runs[0]
        weight = TYPEFACE_WEIGHT.get((first.style.typeface or "").lower(), 500)
        if first.style.bold:
            weight = 700
        # Prefer resolved master inheritance over calibrate_size_pt's hardcoded
        # 16/18pt fallback when run's declared size is None. resolve_inherited_size
        # walks to <p:otherStyle>/<a:lvl1pPr>/<a:defRPr sz="..."> for txBox shapes
        # — authoritative source. calibrate_size_pt's INHERITED_SIZE_REG_PT=16.0
        # fallback would otherwise shadow this for any non-bold inherited run.
        declared = first.style.size_pt
        if declared is None and inherited_size is not None:
            declared = inherited_size
        size_pt = calibrate_size_pt(
            first.style.typeface, declared, first.style.bold,
        )
        color = color_with_alpha(
            first.style.color_hex or "FFFFFF", first.style.color_alpha,
        )
        out.append({
            "algn": p.align,
            "runs": [{
                "text": concat,
                "size_pt": size_pt,
                "weight": weight,
                "color": color,
                "bold": bool(first.style.bold),
            }],
        })
    return out


# ────────────────────────────────────────────────────────────────────────────
# CSS

# OOXML <a:bodyPr> default insets (ECMA-376 §20.1.10.6):
#   lIns = 91440 EMU
#   tIns = 45720 EMU
#   rIns = 91440 EMU
#   bIns = 45720 EMU
#
# Convert EMU to the codebase's "pt" unit. Note: the codebase uses
# EMU_PER_PT = 9525 throughout (parse/pptx.py:20), which is actually
# EMU-per-pixel-at-96-DPI rather than true ECMA points (12700 EMU/pt).
# All other geometry in the renderer (frame x/y/w/h, slide_w_pt,
# slide_h_pt) is in this same px-at-96-DPI unit, so insets must be too
# to stay dimensionally consistent. Using 9525 here makes the OOXML
# defaults convert to 9.6 / 4.8 codebase-pt (= 7.2 / 3.6 true pt),
# matching how every other shape's geometry is computed.
_EMU_PER_CODEBASE_PT = 9525
_DEFAULT_LINS_PT = 91440 / _EMU_PER_CODEBASE_PT  # 9.6
_DEFAULT_TINS_PT = 45720 / _EMU_PER_CODEBASE_PT  # 4.8
_DEFAULT_RINS_PT = 91440 / _EMU_PER_CODEBASE_PT  # 9.6
_DEFAULT_BINS_PT = 45720 / _EMU_PER_CODEBASE_PT  # 4.8


def _read_bodypr_insets(sp_elem) -> tuple:
    """Return (lIns, tIns, rIns, bIns) in codebase-pt, honoring OOXML defaults.

    PPT renders text with these insets between the frame edge and the text
    content. Slide 22's frames omit the attributes entirely, so the OOXML
    defaults apply. Reading explicit values when set lets future slides
    override (e.g. zero-inset chrome elements).
    """
    bp = sp_elem.find("p:txBody/a:bodyPr", NS)
    if bp is None:
        return (_DEFAULT_LINS_PT, _DEFAULT_TINS_PT, _DEFAULT_RINS_PT, _DEFAULT_BINS_PT)
    def _ins(attr, default):
        val = bp.get(attr)
        return int(val) / _EMU_PER_CODEBASE_PT if val is not None else default
    return (
        _ins("lIns", _DEFAULT_LINS_PT),
        _ins("tIns", _DEFAULT_TINS_PT),
        _ins("rIns", _DEFAULT_RINS_PT),
        _ins("bIns", _DEFAULT_BINS_PT),
    )


def _build_css(slide_w_pt, slide_h_pt, title_frame, body_columns, logo_pic,
               logo_invert: bool,
               mobile_title_size_vw: float = 6.0,
               mobile_title_weight: int = 700,
               mobile_title_line_height: float = 1.0) -> str:
    canvas_w_css, canvas_h_css = canvas_aspect_css(slide_w_pt, slide_h_pt)
    parts = []

    parts.append("/* ---------------- DESKTOP ---------------- */")
    parts.append("#deck-desktop { display: block; }")
    parts.append(
        f"#deck-desktop .canvas {{\n"
        f"  position: relative;\n"
        f"  width: {canvas_w_css};\n"
        f"  height: {canvas_h_css};\n"
        f"  margin: 0 auto;\n"
        f"  container-type: inline-size;\n"
        f"  overflow: hidden;\n"
        f"  background: var(--bg);\n"
        f"}}"
    )
    parts.append("#deck-desktop .slide { position: absolute; inset: 0; }")

    if logo_pic is not None:
        logo_filter = "  filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            f"#deck-desktop .gif-logo {{\n"
            f"  {positioned_style(logo_pic.x_pt, logo_pic.y_pt, logo_pic.w_pt, logo_pic.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"{logo_filter}"
            f"}}\n"
            f"#deck-desktop .gif-logo img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )

    if title_frame is not None:
        s = title_frame["shape"]
        # Block layout (not flex) — flex items default to min-width:auto
        # which prevents wrap when content exceeds container width. Block
        # children honor the parent's width and wrap text naturally.
        # Insets emit OOXML bodyPr defaults (lIns=7.2pt, tIns=3.6pt) so text
        # starts at frame.x + lIns (matches PPT visual reference).
        l_ins, t_ins, r_ins, b_ins = _read_bodypr_insets(s.element)
        parts.append(
            f"#deck-desktop .title-frame {{\n"
            f"  {positioned_style(s.x_pt, s.y_pt, s.w_pt, s.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  box-sizing: border-box;\n"
            f"  padding: {pt_to_cqw(t_ins, slide_w_pt)} {pt_to_cqw(r_ins, slide_w_pt)} {pt_to_cqw(b_ins, slide_w_pt)} {pt_to_cqw(l_ins, slide_w_pt)};\n"
            f"  text-align: left;\n"
            f"}}"
        )
        for run in _flat_runs(title_frame):
            parts.append(
                f"#deck-desktop .title-frame .t-title {{\n"
                f"  font-size: {pt_to_cqw(run['size_pt'], slide_w_pt) if run['size_pt'] else '2.5cqw'};\n"
                f"  font-weight: {run['weight']};\n"
                f"  line-height: 1.0;\n"
                f"  color: {run['color']};\n"
                f"  letter-spacing: 0.005em;\n"
                f"  margin: 0;\n"
                f"}}"
            )
            break  # single styling block; title is one paragraph

    for i, col in enumerate(body_columns):
        s = col["shape"]
        # Block layout. Insets honor bodyPr defaults so names start at
        # frame.x + lIns (7.2pt) — matches PPT visual reference.
        l_ins, t_ins, r_ins, b_ins = _read_bodypr_insets(s.element)
        parts.append(
            f"#deck-desktop .body-col-{i} {{\n"
            f"  {positioned_style(s.x_pt, s.y_pt, s.w_pt, s.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  box-sizing: border-box;\n"
            f"  padding: {pt_to_cqw(t_ins, slide_w_pt)} {pt_to_cqw(r_ins, slide_w_pt)} {pt_to_cqw(b_ins, slide_w_pt)} {pt_to_cqw(l_ins, slide_w_pt)};\n"
            f"  text-align: left;\n"
            f"}}"
        )
        # Per-run font-size emission. Two_column body paragraphs in this
        # deck share styling within a column; emit one rule covering the
        # column's body lines. font-weight: 700 + line-height: 1.15 to
        # match slide-3 heading-item visual density (Barlow Condensed
        # Medium at 22pt was reading too light vs the PPT visual reference
        # at the same nominal size — see NOTES.md font-substitution gap).
        any_run = next((r for p in col["paragraphs"] for r in p["runs"]), None)
        if any_run:
            parts.append(
                f"#deck-desktop .body-col-{i} .t-body {{\n"
                f"  font-size: {pt_to_cqw(any_run['size_pt'], slide_w_pt) if any_run['size_pt'] else '1.5cqw'};\n"
                f"  font-weight: 700;\n"
                f"  line-height: 1.15;\n"
                f"  color: {any_run['color']};\n"
                f"  margin: 0;\n"
                f"}}"
            )

    # Mobile: single-column flow, title at top, body columns stack vertically.
    parts.append("\n/* ---------------- MOBILE ---------------- */")
    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    parts.append("  #deck-desktop { display: none; }")
    parts.append("  #deck-mobile { display: block; background: var(--bg); }")
    parts.append(
        "  #deck-mobile .panel {\n"
        "    min-height: 100vh;\n"
        "    color: #FFFFFF;\n"
        "    padding: 0;\n"
        "  }"
    )
    if logo_pic is not None:
        # Mobile top-bar — cyan with always-on invert (matches the
        # logo_invert hint behavior across templates).
        parts.append(
            "  #deck-mobile .top-bar {\n"
            "    display: flex; justify-content: flex-end;\n"
            "    padding: 5vw;\n"
            "    background: var(--bg-cyan);\n"
            "  }\n"
            "  #deck-mobile .gif-logo-mobile {\n"
            "    width: 22vw;\n"
            "    filter: brightness(0) invert(1);\n"
            "  }\n"
            "  #deck-mobile .gif-logo-mobile img { width: 100%; }"
        )
    if title_frame is not None:
        parts.append(
            "  #deck-mobile .title-mobile {\n"
            "    padding: 8vw 5vw 4vw;\n"
            "    text-align: left;\n"
            "  }\n"
            "  #deck-mobile .title-mobile .t-title {\n"
            f"    font-size: {mobile_title_size_vw}vw;\n"
            f"    font-weight: {mobile_title_weight};\n"
            f"    line-height: {mobile_title_line_height};\n"
            "    color: #FFFFFF;\n"
            "    letter-spacing: 0.01em;\n"
            "    margin: 0;\n"
            "  }"
        )
    if body_columns:
        parts.append(
            "  #deck-mobile .body-mobile {\n"
            "    padding: 4vw 6vw 8vw;\n"
            "    display: flex; flex-direction: column;\n"
            "    gap: 6vw;\n"
            "  }\n"
            "  #deck-mobile .body-col-mobile {\n"
            "    display: flex; flex-direction: column;\n"
            "  }\n"
            "  #deck-mobile .body-col-mobile .t-body {\n"
            "    font-size: 4.4vw;\n"
            "    font-weight: 400;\n"
            "    line-height: 1.0;\n"
            "    color: #FFFFFF;\n"
            "    margin: 0;\n"
            "  }"
        )

    parts.append("}")
    return "\n".join(parts)


def _flat_runs(frame: dict):
    for para in frame["paragraphs"]:
        for run in para["runs"]:
            yield run


# ────────────────────────────────────────────────────────────────────────────
# Body HTML

def _build_body(title_frame, body_columns, logo_url) -> str:
    desktop_lines = []
    if logo_url:
        desktop_lines.append(f'      <div class="gif-logo"><img src="{logo_url}" alt=""></div>')
    if title_frame is not None:
        desktop_lines.append('      <div class="title-frame">')
        for para in title_frame["paragraphs"]:
            for run in para["runs"]:
                desktop_lines.append(
                    f'        <div class="t t-title">{_escape(run["text"])}</div>'
                )
        desktop_lines.append('      </div>')
    for i, col in enumerate(body_columns):
        desktop_lines.append(f'      <div class="body-col-{i}">')
        for para in col["paragraphs"]:
            if not para["runs"]:
                continue
            for run in para["runs"]:
                desktop_lines.append(
                    f'        <div class="t t-body">{_escape(run["text"])}</div>'
                )
        desktop_lines.append('      </div>')

    mobile_lines = []
    if logo_url:
        mobile_lines.append('    <div class="top-bar">')
        mobile_lines.append(f'      <div class="gif-logo-mobile"><img src="{logo_url}" alt=""></div>')
        mobile_lines.append('    </div>')
    if title_frame is not None:
        mobile_lines.append('    <div class="title-mobile">')
        for para in title_frame["paragraphs"]:
            for run in para["runs"]:
                mobile_lines.append(
                    f'      <div class="t t-title">{_escape(run["text"])}</div>'
                )
        mobile_lines.append('    </div>')
    if body_columns:
        mobile_lines.append('    <div class="body-mobile">')
        for col in body_columns:
            mobile_lines.append('      <div class="body-col-mobile">')
            for para in col["paragraphs"]:
                if not para["runs"]:
                    continue
                for run in para["runs"]:
                    mobile_lines.append(
                        f'        <div class="t t-body">{_escape(run["text"])}</div>'
                    )
            mobile_lines.append('      </div>')
        mobile_lines.append('    </div>')

    return (
        "<!-- DESKTOP -->\n"
        "<section id=\"deck-desktop\">\n"
        "  <div class=\"canvas\">\n"
        "    <div class=\"slide\">\n"
        f"{chr(10).join(desktop_lines)}\n"
        "    </div>\n"
        "  </div>\n"
        "</section>\n\n"
        "<!-- MOBILE -->\n"
        "<section id=\"deck-mobile\">\n"
        "  <div class=\"panel\">\n"
        f"{chr(10).join(mobile_lines)}\n"
        "  </div>\n"
        "</section>"
    )


def _escape(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


# ────────────────────────────────────────────────────────────────────────────
# Contact variant (slide 23)
#
# Variant dispatched via `two_column_layout: "contact"` hint. Differs from
# the stakeholders variant in three ways:
#   1. Has a left-half SVG hero (the Identity logo via SVG-only blip) — no
#      stakeholders slide has this.
#   2. Has a content-position image (the GIF wordmark) in the right column,
#      which is NOT the standard top-right chrome logo (`is_logo_pic` rejects
#      it on geometry). The content image emits as <div class="photo-0"> so
#      the asset validator reads it as role=photo_grid.
#   3. The text frames don't split into a single title + N body columns;
#      they're a stack of N adjacent contact blocks (bold headers + plain
#      detail lines) all in the right column. Per-paragraph bold flag drives
#      weight (no slide-wide title/body distinction).
#
# Hyperlink handling: mailto: hyperlinks present on email runs (rId4/rId5
# on slide 23) are NOT extracted by parse_text_frame and render as plain
# text. Clickability deferred to Phase 2 alongside the broader hlinkClick
# extraction work.

def _render_contact_variant(
    flat, slide, theme, slide_index, slide_class,
    deck_name, slide_w_pt, slide_h_pt, canvas_bg, deck_brand_color,
):
    hero_pic, wordmark_pic, text_frames = _classify_contact(
        flat, slide_w_pt, slide_h_pt, canvas_bg, theme, slide,
    )

    hero_url = inline_data_url(hero_pic, slide) if hero_pic is not None else None
    wordmark_url = inline_data_url(wordmark_pic, slide) if wordmark_pic is not None else None

    slide_css = _build_css_contact(
        slide_w_pt, slide_h_pt,
        hero_pic, wordmark_pic, text_frames,
    )
    body_html = _build_body_contact(hero_url, wordmark_url, text_frames)

    return render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg":        canvas_bg,
            "bg-cyan":   deck_brand_color or "#000000",
            "headline":  "#FFFFFF",
            "font-cond": '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=body_html,
        slide_css=slide_css,
    )


def _classify_contact(flat, slide_w_pt, slide_h_pt, canvas_bg, theme, slide):
    """Return (hero_pic, wordmark_pic, text_frames).

    Pic disambiguation:
      - hero: pic with an SVG-only blip (has_svg_blip) — slide 23 Graphic 4.
      - wordmark: a non-SVG pic that is NOT geometric chrome logo — slide 23
        Picture 10. There's only ever one such pic on a contact slide; if
        future contact slides carry more, this picks the first encountered
        and the rest are silently skipped (caller can extend).

    Text frames are returned as the full y-then-x sorted list — no title /
    body split.
    """
    hero_pic = None
    wordmark_pic = None
    text_frames = []
    for s in flat:
        if s.x_pt is None:
            continue
        if s.kind == "pic":
            if hero_pic is None and has_svg_blip(s.element):
                hero_pic = s
                continue
            if wordmark_pic is None and not is_logo_pic(s, slide_w_pt, slide_h_pt):
                wordmark_pic = s
            continue
        if s.kind != "sp":
            continue
        if _is_canvas_skip_rect(s, slide_w_pt, slide_h_pt, canvas_bg):
            continue
        paras = _extract_paragraphs(s, theme, slide)
        if not paras or not any(p["runs"] for p in paras):
            continue
        text_frames.append({"shape": s, "paragraphs": paras})

    text_frames.sort(key=lambda f: (round(f["shape"].y_pt), round(f["shape"].x_pt)))
    return hero_pic, wordmark_pic, text_frames


def _build_css_contact(slide_w_pt, slide_h_pt, hero_pic, wordmark_pic, text_frames):
    canvas_w_css, canvas_h_css = canvas_aspect_css(slide_w_pt, slide_h_pt)
    parts = []

    parts.append("/* ---------------- DESKTOP ---------------- */")
    parts.append("#deck-desktop { display: block; }")
    parts.append(
        f"#deck-desktop .canvas {{\n"
        f"  position: relative;\n"
        f"  width: {canvas_w_css};\n"
        f"  height: {canvas_h_css};\n"
        f"  margin: 0 auto;\n"
        f"  container-type: inline-size;\n"
        f"  overflow: hidden;\n"
        f"  background: var(--bg);\n"
        f"}}"
    )
    parts.append("#deck-desktop .slide { position: absolute; inset: 0; }")

    if hero_pic is not None:
        parts.append(
            f"#deck-desktop .hero {{\n"
            f"  {positioned_style(hero_pic.x_pt, hero_pic.y_pt, hero_pic.w_pt, hero_pic.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"}}\n"
            f"#deck-desktop .hero img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )
    if wordmark_pic is not None:
        parts.append(
            f"#deck-desktop .photo-0 {{\n"
            f"  {positioned_style(wordmark_pic.x_pt, wordmark_pic.y_pt, wordmark_pic.w_pt, wordmark_pic.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"}}\n"
            f"#deck-desktop .photo-0 img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )

    for i, tf in enumerate(text_frames):
        s = tf["shape"]
        l_ins, t_ins, r_ins, b_ins = _read_bodypr_insets(s.element)
        # Per-frame size: sample first run (slide 23 contact frames are
        # single-styled within a frame). Fall back to inherited if size_pt
        # is None on the resolved run dict.
        first_run = next((r for p in tf["paragraphs"] for r in p["runs"]), None)
        size_css = (
            pt_to_cqw(first_run["size_pt"], slide_w_pt)
            if first_run and first_run["size_pt"]
            else "1.406cqw"  # 18pt at 1280pt canvas
        )
        color_css = first_run["color"] if first_run else "#FFFFFF"
        parts.append(
            f"#deck-desktop .text-block-{i} {{\n"
            f"  {positioned_style(s.x_pt, s.y_pt, s.w_pt, s.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  box-sizing: border-box;\n"
            f"  padding: {pt_to_cqw(t_ins, slide_w_pt)} {pt_to_cqw(r_ins, slide_w_pt)} {pt_to_cqw(b_ins, slide_w_pt)} {pt_to_cqw(l_ins, slide_w_pt)};\n"
            f"  font-size: {size_css};\n"
            f"  line-height: 1.15;\n"
            f"  color: {color_css};\n"
            f"  text-align: left;\n"
            f"}}\n"
            f"#deck-desktop .text-block-{i} .t {{ margin: 0; }}"
        )

    # Mobile: vertical stack — hero, wordmark, text blocks in source order.
    parts.append("\n/* ---------------- MOBILE ---------------- */")
    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    parts.append("  #deck-desktop { display: none; }")
    parts.append("  #deck-mobile { display: block; background: var(--bg); }")
    parts.append(
        "  #deck-mobile .panel {\n"
        "    min-height: 100vh;\n"
        "    color: #FFFFFF;\n"
        "    padding: 8vw 6vw 12vw;\n"
        "    display: flex; flex-direction: column;\n"
        "    gap: 6vw;\n"
        "  }"
    )
    if hero_pic is not None:
        parts.append(
            "  #deck-mobile .hero-mobile {\n"
            "    width: 70vw; align-self: center;\n"
            "  }\n"
            "  #deck-mobile .hero-mobile img { width: 100%; }"
        )
    if wordmark_pic is not None:
        parts.append(
            "  #deck-mobile .photo-m {\n"
            "    width: 50vw; align-self: center; margin-top: 4vw;\n"
            "  }\n"
            "  #deck-mobile .photo-m img { width: 100%; }"
        )
    if text_frames:
        parts.append(
            "  #deck-mobile .contact-stack-mobile {\n"
            "    display: flex; flex-direction: column;\n"
            "    gap: 4vw;\n"
            "    margin-top: 4vw;\n"
            "  }\n"
            "  #deck-mobile .contact-block-mobile {\n"
            "    display: flex; flex-direction: column;\n"
            "  }\n"
            "  #deck-mobile .contact-block-mobile .t {\n"
            "    font-size: 4.4vw;\n"
            "    line-height: 1.2;\n"
            "    color: #FFFFFF;\n"
            "    margin: 0;\n"
            "  }"
        )

    parts.append("}")
    return "\n".join(parts)


def _build_body_contact(hero_url, wordmark_url, text_frames) -> str:
    desktop_lines = []
    if hero_url:
        desktop_lines.append(f'      <div class="hero"><img src="{hero_url}" alt=""></div>')
    if wordmark_url:
        desktop_lines.append(f'      <div class="photo-0"><img src="{wordmark_url}" alt=""></div>')
    for i, tf in enumerate(text_frames):
        desktop_lines.append(f'      <div class="text-block-{i}">')
        for para in tf["paragraphs"]:
            for run in para["runs"]:
                weight_attr = f' style="font-weight: {run["weight"]};"' if run.get("weight") else ""
                desktop_lines.append(
                    f'        <div class="t"{weight_attr}>{_escape(run["text"])}</div>'
                )
        desktop_lines.append('      </div>')

    mobile_lines = []
    if hero_url:
        mobile_lines.append(f'    <div class="hero-mobile"><img src="{hero_url}" alt=""></div>')
    if wordmark_url:
        mobile_lines.append(f'    <div class="photo-m"><img src="{wordmark_url}" alt=""></div>')
    if text_frames:
        mobile_lines.append('    <div class="contact-stack-mobile">')
        for tf in text_frames:
            mobile_lines.append('      <div class="contact-block-mobile">')
            for para in tf["paragraphs"]:
                for run in para["runs"]:
                    weight_attr = f' style="font-weight: {run["weight"]};"' if run.get("weight") else ""
                    mobile_lines.append(
                        f'        <div class="t"{weight_attr}>{_escape(run["text"])}</div>'
                    )
            mobile_lines.append('      </div>')
        mobile_lines.append('    </div>')

    return (
        "<!-- DESKTOP -->\n"
        "<section id=\"deck-desktop\">\n"
        "  <div class=\"canvas\">\n"
        "    <div class=\"slide\">\n"
        f"{chr(10).join(desktop_lines)}\n"
        "    </div>\n"
        "  </div>\n"
        "</section>\n\n"
        "<!-- MOBILE -->\n"
        "<section id=\"deck-mobile\">\n"
        "  <div class=\"panel\">\n"
        f"{chr(10).join(mobile_lines)}\n"
        "  </div>\n"
        "</section>"
    )
