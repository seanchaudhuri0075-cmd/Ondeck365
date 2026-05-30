"""title_stats template — title + structured text list, three variants.

Variants, auto-detected from body-shape structure:

  plain_list                    — 1 body shape, single text list (default
                                  for future decks; no current P&G slide).
  dotted_leader_paired          — 2 body shapes with equal paragraph count.
                                  Each para[N] in shape A = label, in shape
                                  B = value; rendered as flex rows with a
                                  CSS dotted-leader spacer between them.
                                  (P&G slide 21 — turnaround times.)
  dotted_leader_continuation    — 2 body shapes with unequal paragraph
                                  count. Paragraphs concatenated as one
                                  list, rendered as a 2-column CSS layout.
                                  (P&G slide 3 — type of brand assets.)

Detection heuristic (no manifest hint required for the canonical cases):
  - 0 or 1 body shape           → plain_list
  - 2 body shapes, n_paras equal → dotted_leader_paired
  - 2 body shapes, n_paras differ → dotted_leader_continuation
  - 3+ body shapes              → plain_list with all shapes concatenated
                                  (defensive fallback; no canonical case yet)

Manifest opt-in: hints.title_stats_layout = "<variant>" forces a specific
variant when the heuristic misfires (e.g., paired data where one entry is
missing in the value shape, leaving the counts unequal).

DOTTED LEADER DETECTION IS CONTENT-BASED, NOT ATTRIBUTE-BASED. The bundle
authors typed Unicode dot characters (……, …, ...) directly into the text
content rather than using OOXML's <a:tab leader='dot'/> tab-stop mechanism.
We scan label runs for ≥3 consecutive dot characters of any common Unicode
variant — … (horizontal ellipsis), · (middle dot), ‧
(hyphenation point), ⋯ (midline horizontal ellipsis), • (bullet),
ASCII period, etc. — and strip the dot-run from the rendered text. The
visual leader is then drawn by a CSS pseudo-element span with
border-bottom: dotted, so the leader scales to whatever space remains
between label and value (no literal '....' characters in the HTML, no
fragility when labels are long).

ROLE CLASSIFICATION (geometry-based, mirrors other templates):
  title:    text shape in top 30% of canvas, single paragraph, large font.
            Largest such shape if multiple candidates.
  body:     text shapes that aren't title or footnote — middle of canvas.
            Sorted left-to-right by x_pt for deterministic A-vs-B labelling.
  footnote: text shape in bottom 15% of canvas, small font; optional.
            Rendered as a small italic line if present.
  logo:     <p:pic> matching _shared.is_logo_pic (top-right corner).

DESIGN DEFAULTS (locked-in answers from the design discussion):
  - Plain-list bullet character: emit no bullet by default (OOXML's
    omission of <a:buChar> is treated as a signal). hints.bullet_char
    overrides per slide if a future deck wants visible bullets.
  - Mobile dotted-leader: rows stay flex even when long labels wrap
    (the leader spacer's flex:1 stretches across whatever's left). No
    detect-and-stack behavior; that's a Phase 2 follow-up if needed.
  - Title vs list disambiguation: handled only when title and body are
    distinct shapes. Single-shape slides where para 0 = title and 1..N
    = list aren't supported; defer until a slide hits that case.
"""
from __future__ import annotations

import re
from typing import Optional

from pptx.slide import Slide

from ..css import pt_to_pct_x, pt_to_pct_y, pt_to_cqw, color_with_alpha
from ..desktop import canvas_aspect_css
from ..html import render_page
from ...layout.detect import SlideClass
from ...parse.font_calibration import calibrate_size_pt
from ...parse.shapes import FlatShape, flatten_slide
from ...parse.slide import NS
from ...parse.text import parse_text_frame
from ...parse.theme import Theme
from ._shared import image_src, inline_data_url, is_logo_pic


# Unicode dot characters that count as "leader fill" when ≥3 in a row.
# Horizontal ellipsis, ASCII period, middle dot, hyphenation point, midline
# horizontal ellipsis, bullet operator, bullet, two-dot leader, three-dot leader.
_DOT_CHARS = "….·‧⋯∙•‥․"
_DOT_RUN_RE = re.compile(rf"\s*[{re.escape(_DOT_CHARS)}]{{3,}}\s*$")


# Maps typeface name (lowercased) → numeric weight, mirrors section_divider.
TYPEFACE_WEIGHT = {
    "univers":                  500,
    "univers light":            300,
    "univers bold":             700,
    "univers condensed":        500,
    "univers condensed light":  300,
    "univers condensed bold":   700,
    "barlow condensed":         500,
    "barlow condensed light":   300,
    "barlow condensed bold":    700,
}


def render_title_stats(
    slide: Slide,
    theme: Theme,
    slide_index: int,
    slide_class: SlideClass,
    deck_name: str,
    slide_w_pt: float,
    slide_h_pt: float,
    deck_brand_color: str = "#000000",
) -> str:
    """Render a title_stats slide. Dispatches on detected body-shape variant.

    Markup class names mirror the bundle (`.L`/`.gif-logo` on desktop;
    `.top-row`+`.item` for paired mobile, `.header-row`+`.uc-item-mobile`
    for continuation mobile) so any future Deck Editor compatibility work
    has a clean parity surface against the original bundle output.
    """
    flat = list(flatten_slide(slide))
    title_sp, body_sps, footnote_sp, logo_pic = _classify(flat, slide_w_pt, slide_h_pt)
    variant = _detect_variant(body_sps, slide_class)

    hints = slide_class.hints
    title_size_pt_hint   = hints.get("title_size_pt")
    heading_size_pt_hint = hints.get("heading_size_pt")
    body_size_pt_hint    = hints.get("body_size_pt")
    logo_invert          = bool(hints.get("logo_invert", False))
    hide_footnote        = bool(hints.get("hide_footnote", False))

    title_run = _extract_title_run(title_sp, theme) if title_sp else None
    if title_run is not None and title_size_pt_hint is not None:
        title_run["size_pt"] = float(title_size_pt_hint)

    footnote_run = None
    if footnote_sp is not None and not hide_footnote:
        footnote_run = _extract_title_run(footnote_sp, theme)

    logo_url = inline_data_url(logo_pic, slide) if logo_pic else None
    canvas_bg = hints.get("canvas_bg") or deck_brand_color or "#000000"

    if variant == "dotted_leader_paired":
        desktop_body_html, desktop_body_css, mobile_panel_html, mobile_panel_css = (
            _build_paired(
                body_sps, theme, slide_w_pt, slide_h_pt,
                title_run, logo_url,
                body_size_pt=float(body_size_pt_hint) if body_size_pt_hint else 16.0,
            )
        )
    elif variant == "dotted_leader_continuation":
        desktop_body_html, desktop_body_css, mobile_panel_html, mobile_panel_css = (
            _build_continuation(
                body_sps, theme, slide_w_pt, slide_h_pt,
                title_run, logo_url,
                heading_size_pt=float(heading_size_pt_hint) if heading_size_pt_hint else 18.0,
                body_size_pt=float(body_size_pt_hint) if body_size_pt_hint else 16.0,
            )
        )
    else:  # plain_list
        desktop_body_html, desktop_body_css, mobile_panel_html, mobile_panel_css = (
            _build_plain_list(
                body_sps, theme, slide_w_pt, slide_h_pt, slide_class,
                title_run, logo_url,
                body_size_pt=float(body_size_pt_hint) if body_size_pt_hint else 18.0,
            )
        )

    footnote_size_cqw = (
        pt_to_cqw(footnote_run["size_pt"], slide_w_pt)
        if footnote_run else "1.094cqw"
    )

    slide_css = _build_css(
        slide_w_pt, slide_h_pt,
        title_sp, footnote_sp, logo_pic,
        desktop_body_css, mobile_panel_css, canvas_bg,
        title_run, footnote_run, footnote_size_cqw, logo_invert,
    )
    full_body_html = _build_body(
        title_run, desktop_body_html, footnote_run, logo_url,
        mobile_panel_html, logo_invert,
    )

    return render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg":         canvas_bg,
            "bg-cyan":    deck_brand_color or "#000000",
            "headline":   "#FFFFFF",
            "font-cond":  '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=full_body_html,
        slide_css=slide_css,
    )


# ────────────────────────────────────────────────────────────────────────────
# Role classification

def _classify(shapes, slide_w_pt: float, slide_h_pt: float):
    """Return (title_sp, body_sps, footnote_sp, logo_pic).

    Geometry rules:
      - title:    text shape, single paragraph, y < slide_h * 0.3.
                  Largest by w_pt if multiple candidates.
      - footnote: text shape, single paragraph, y > slide_h * 0.85, small.
      - body:     text shapes not classified as title/footnote, sorted L→R.
      - logo:     pic in top-right (via is_logo_pic).
    """
    title_sp = None
    body_sps: list[FlatShape] = []
    footnote_sp = None
    logo_pic = None

    candidates_with_text = []
    for s in shapes:
        if s.x_pt is None:
            continue
        if s.kind == "pic":
            if is_logo_pic(s, slide_w_pt, slide_h_pt):
                logo_pic = s
            continue
        if s.kind != "sp" or not _has_text(s.element):
            continue
        candidates_with_text.append(s)

    for s in candidates_with_text:
        n_paras = _count_paragraphs(s.element)
        in_top = s.y_pt < slide_h_pt * 0.30
        in_bottom = s.y_pt > slide_h_pt * 0.85

        if in_top and n_paras == 1:
            if title_sp is None or (s.w_pt or 0) > (title_sp.w_pt or 0):
                title_sp = s
        elif in_bottom and n_paras == 1:
            footnote_sp = s
        else:
            body_sps.append(s)

    body_sps.sort(key=lambda s: s.x_pt)
    return title_sp, body_sps, footnote_sp, logo_pic


def _has_text(sp_elem) -> bool:
    """True if the shape carries at least one non-blank <a:t>."""
    for t in sp_elem.findall(".//a:t", NS):
        if t.text and t.text.strip():
            return True
    return False


def _count_paragraphs(sp_elem) -> int:
    """Number of <a:p> children in this shape's <p:txBody>."""
    tx = sp_elem.find("p:txBody", NS)
    if tx is None:
        return 0
    return len(tx.findall("a:p", NS))


# ────────────────────────────────────────────────────────────────────────────
# Variant detection

def _detect_variant(body_sps, slide_class: SlideClass) -> str:
    """Pick a layout variant. Manifest hint > heuristic."""
    hint = slide_class.hints.get("title_stats_layout")
    if hint in ("plain_list", "dotted_leader_paired", "dotted_leader_continuation"):
        return hint
    if len(body_sps) <= 1:
        return "plain_list"
    if len(body_sps) == 2:
        a, b = body_sps
        n1 = _count_paragraphs(a.element)
        n2 = _count_paragraphs(b.element)
        return "dotted_leader_paired" if n1 == n2 else "dotted_leader_continuation"
    return "plain_list"  # 3+ shapes — defensive fallback


# ────────────────────────────────────────────────────────────────────────────
# Run + paragraph extraction

def _extract_title_run(title_sp, theme: Theme) -> Optional[dict]:
    """Single-paragraph title text + first-run weight + size."""
    tf = parse_text_frame(title_sp.element, theme)
    if tf is None or not tf.paragraphs:
        return None
    runs_text = []
    weight = 500
    size_pt = None
    color = "#FFFFFF"
    for p in tf.paragraphs:
        for r in p.runs:
            if r.text.strip():
                runs_text.append(r.text)
                if r.style.bold:
                    weight = 700
                elif r.style.typeface:
                    tw = TYPEFACE_WEIGHT.get(r.style.typeface.lower())
                    if tw is not None:
                        weight = tw
                if size_pt is None:
                    size_pt = calibrate_size_pt(
                        r.style.typeface, r.style.size_pt, r.style.bold,
                    )
                if r.style.color_hex:
                    color = color_with_alpha(r.style.color_hex, r.style.color_alpha)
    if not runs_text:
        return None
    return {
        "text": "".join(runs_text).strip(),
        "weight": weight,
        "size_pt": size_pt or 18.0,
        "color": color,
    }


def _strip_dot_suffix(text: str) -> tuple[str, bool]:
    """Strip a trailing ≥3-char dot run. Returns (clean_text, had_leader)."""
    new = _DOT_RUN_RE.sub("", text)
    return new.rstrip(), new != text


def _extract_value_paragraphs(text_sp, theme: Theme) -> list[list[tuple]]:
    """Per paragraph: a list of (text, weight, size_pt) tuples — one per run.

    Used for the value column in the paired layout, where the bundle's
    per-run formatting matters: the lead run is bold (e.g. "3-5 days")
    and the trail run is the regular parenthetical context. Flattening
    to a single span loses that distinction.
    """
    tf = parse_text_frame(text_sp.element, theme)
    if tf is None:
        return []
    out = []
    for p in tf.paragraphs:
        runs = []
        for r in p.runs:
            if not r.text:
                continue
            stripped = r.text.strip()
            is_dot_only = (
                len(stripped) >= 3
                and all(c in _DOT_CHARS for c in stripped)
            )
            if is_dot_only:
                continue
            weight = 500
            if r.style.bold:
                weight = 700
            elif r.style.bold is False:
                weight = 400
            elif r.style.typeface:
                tw = TYPEFACE_WEIGHT.get(r.style.typeface.lower())
                if tw is not None:
                    weight = tw
            size_pt = calibrate_size_pt(
                r.style.typeface, r.style.size_pt, r.style.bold,
            )
            runs.append((r.text, weight, size_pt or 16.0))
        if not any(t.strip() for (t, _, _) in runs):
            continue
        out.append(runs)
    return out


def _union_box_pt(shapes) -> tuple:
    """Return (x, y, w, h) in points covering all shapes' bounding boxes."""
    x0 = min(s.x_pt for s in shapes)
    y0 = min(s.y_pt for s in shapes)
    x1 = max(s.x_pt + s.w_pt for s in shapes)
    y1 = max(s.y_pt + s.h_pt for s in shapes)
    return (x0, y0, x1 - x0, y1 - y0)


def _box_to_pct(box_pt, slide_w_pt: float, slide_h_pt: float) -> dict:
    """Format a (x, y, w, h) box as a CSS positioning dict in percent."""
    x, y, w, h = box_pt
    return {
        "left":   pt_to_pct_x(x, slide_w_pt),
        "top":    pt_to_pct_y(y, slide_h_pt),
        "width":  pt_to_pct_x(w, slide_w_pt),
        "height": pt_to_pct_y(h, slide_h_pt),
    }


def _extract_paragraphs(text_sp, theme: Theme) -> list[dict]:
    """Per paragraph: {'text', 'is_bold', 'has_leader', 'size_pt'}.
    Concatenates runs in document order; strips trailing dot-leader runs."""
    tf = parse_text_frame(text_sp.element, theme)
    if tf is None:
        return []
    out = []
    for p in tf.paragraphs:
        # Bold-ness = any non-dot-run is bold (heuristic: dot runs use light weight,
        # the actual content runs carry the paragraph's true weight).
        text_parts = []
        is_bold = False
        size_pt = None
        for r in p.runs:
            if not r.text:
                continue
            text_parts.append(r.text)
            stripped = r.text.strip()
            is_dot_only = (
                len(stripped) >= 3
                and all(c in _DOT_CHARS for c in stripped)
            )
            if not is_dot_only:
                if r.style.bold:
                    is_bold = True
                if size_pt is None:
                    size_pt = calibrate_size_pt(
                        r.style.typeface, r.style.size_pt, r.style.bold,
                    )
        text = "".join(text_parts)
        clean, had_leader = _strip_dot_suffix(text)
        if not clean.strip() and not had_leader:
            continue  # blank paragraph (P&G slides have stray empties)
        out.append({
            "text": clean.strip(),
            "is_bold": is_bold,
            "has_leader": had_leader,
            "size_pt": size_pt or 16.0,
        })
    return out


# ────────────────────────────────────────────────────────────────────────────
# Variant body builders — each returns (inner_html, variant_specific_css)

def _build_paired(body_sps, theme: Theme,
                  slide_w_pt: float, slide_h_pt: float,
                  title_run, logo_url,
                  body_size_pt: float) -> tuple[str, str, str, str]:
    """Slide-21 style. Returns (desktop_body_html, desktop_body_css,
    mobile_panel_html, mobile_panel_css).

    Desktop layout uses a single absolute-positioned `.paired-body` at
    the union bbox of body_sps[0]+body_sps[1], with one `.uc-row` flex
    row per entry containing `.l-name` + `.l-leader` + `.r-value`. The
    leader span fills dynamically between the label content and the
    value column on each row, so the dotted line ends right at the
    value text regardless of label/value content lengths. The value
    column's flex-basis is fixed at sp[1].w / union.w so its right
    edge mirrors the source sp[3] right edge across rows. Per-run
    weight is preserved inside `.r-value` via `.rt-bold` / `.rt-light`
    spans.

    Trade-off vs the source's two-shape geometry: sp[0] and sp[1] are
    at slightly different y-positions in source (12pt offset on slide
    21); collapsing both into one container aligns all rows to sp[0]'s
    y-position. Accepted divergence — visually negligible vs the
    bundle's per-row leader/value gap that this fix removes.

    Mobile is stacked, not row-flex: `<div class="item">` per entry with
    `<div class="name">` (label) on top and `<div class="duration">` below
    containing `.b` / `.l` per-run spans. The dotted-line visual on
    mobile is a single `border-bottom` on the entire `.item`, NOT a
    per-row leader.
    """
    labels = _extract_paragraphs(body_sps[0], theme)
    value_paras = _extract_value_paragraphs(body_sps[1], theme)
    n = min(len(labels), len(value_paras))

    union = _union_box_pt(body_sps)
    pos = _box_to_pct(union, slide_w_pt, slide_h_pt)
    value_col_pct = body_sps[1].w_pt / union[2] * 100
    body_size_cqw = pt_to_cqw(body_size_pt, slide_w_pt)

    # Desktop body — single container, per-entry flex row
    rows = []
    for i in range(n):
        lab = labels[i]
        val_runs = value_paras[i]
        value_inner = "".join(
            (f'<span class="rt-bold">{_escape(t)}</span>' if w >= 700
             else f'<span class="rt-light">{_escape(t)}</span>')
            for (t, w, _sz) in val_runs
        )
        rows.append(
            '        <div class="uc-row">'
            f'<span class="l-name">{_escape(lab["text"])}</span>'
            '<span class="l-leader"></span>'
            f'<span class="r-value">{value_inner}</span>'
            '</div>'
        )

    desktop_html = (
        '      <div class="paired-body">\n'
        + "\n".join(rows)
        + '\n      </div>'
    )

    desktop_css = (
        "#deck-desktop .paired-body {\n"
        "  position: absolute;\n"
        f"  left: {pos['left']}; top: {pos['top']};\n"
        f"  width: {pos['width']}; height: {pos['height']};\n"
        "  display: flex; flex-direction: column; gap: 1.4cqw;\n"
        f"  font-size: {body_size_cqw}; color: var(--headline);\n"
        "}\n"
        "#deck-desktop .paired-body .uc-row {\n"
        "  display: flex; align-items: baseline; gap: 0.5cqw;\n"
        "  line-height: 1.3;\n"
        "}\n"
        "#deck-desktop .paired-body .l-name { flex: 0 0 auto; font-weight: 700; white-space: nowrap; }\n"
        "#deck-desktop .paired-body .l-leader {\n"
        "  flex: 1 1 auto; min-width: 1cqw;\n"
        "  border-bottom: 0.15cqw dotted rgba(255,255,255,0.95);\n"
        "  position: relative; top: -0.25cqw;\n"
        "}\n"
        f"#deck-desktop .paired-body .r-value {{ flex: 0 0 {value_col_pct:.4f}%; }}\n"
        "#deck-desktop .paired-body .r-value .rt-bold { font-weight: 700; }\n"
        "#deck-desktop .paired-body .r-value .rt-light { font-weight: 300; }\n"
    )

    # Mobile — stacked .item blocks, title + logo in a top-row above
    item_html = []
    for i in range(n):
        lab = labels[i]
        val_runs = value_paras[i]
        duration_html = "".join(
            (f'<span class="b">{_escape(t)}</span>' if w >= 700
             else f'<span class="l">{_escape(t)}</span>')
            for (t, w, _sz) in val_runs
        )
        item_html.append(
            '      <div class="item">\n'
            f'        <div class="name">{_escape(lab["text"])}</div>\n'
            f'        <div class="duration">{duration_html}</div>\n'
            '      </div>'
        )

    title_text = _escape(title_run["text"]) if title_run else ""
    logo_img = (
        f'<div class="gif-logo-mobile"><img src="{logo_url}" alt=""></div>'
        if logo_url else ""
    )
    mobile_html = (
        '    <div class="top-row">\n'
        f'      <div class="L L-mobile"><div class="t">{title_text}</div></div>\n'
        f'      {logo_img}\n'
        '    </div>\n'
        + "\n".join(item_html)
    )

    mobile_css = (
        "  #deck-mobile .panel {\n"
        "    background: var(--bg-cyan);\n"
        "    color: var(--headline);\n"
        "    padding: 6vw 5vw 12vw;\n"
        "    min-height: 100vh;\n"
        "  }\n"
        "  #deck-mobile .top-row {\n"
        "    display: flex; justify-content: space-between;\n"
        "    align-items: flex-start; margin-bottom: 6vw;\n"
        "  }\n"
        "  #deck-mobile .L-mobile { flex: 1; padding-right: 4vw; }\n"
        "  #deck-mobile .L-mobile .t {\n"
        "    font-size: 8vw; font-weight: 400; line-height: 1.0;\n"
        "    text-transform: uppercase;\n"
        "  }\n"
        "  #deck-mobile .gif-logo-mobile { width: 18vw; }\n"
        "  #deck-mobile .gif-logo-mobile img { width: 100%; }\n"
        "  #deck-mobile .item {\n"
        "    margin-bottom: 4vw; padding-bottom: 3vw;\n"
        "    border-bottom: 1px dotted rgba(255,255,255,0.5);\n"
        "  }\n"
        "  #deck-mobile .item .name {\n"
        "    font-weight: 700; font-size: 4.5vw; margin-bottom: 1vw;\n"
        "  }\n"
        "  #deck-mobile .item .duration { font-size: 4vw; line-height: 1.3; }\n"
        "  #deck-mobile .item .duration .b { font-weight: 700; }\n"
        "  #deck-mobile .item .duration .l { font-weight: 300; }\n"
    )

    return (desktop_html, desktop_css, mobile_html, mobile_css)


def _group_continuation_items(paragraphs) -> list[dict]:
    """Group a flat paragraph list into [{heading, sub}] items.

    A bold paragraph starts a new item; following non-bold paragraphs
    concatenate as that item's sub. The bundle's slide-3 output joins
    multi-paragraph sub continuations into one ucb line, so e.g. a
    heading "Package Design Concepts & Artwork" followed by two sub
    paragraphs collapses to one item.
    """
    items: list[dict] = []
    current: dict | None = None
    for p in paragraphs:
        if p["is_bold"]:
            if current is not None:
                items.append(current)
            current = {"heading": p["text"], "subs": []}
        else:
            if current is None:
                current = {"heading": "", "subs": [p["text"]]}
            else:
                current["subs"].append(p["text"])
    if current is not None:
        items.append(current)
    return items


def _build_continuation(body_sps, theme: Theme,
                        slide_w_pt: float, slide_h_pt: float,
                        title_run, logo_url,
                        heading_size_pt: float,
                        body_size_pt: float) -> tuple[str, str, str, str]:
    """Slide-3 style. Returns (desktop_body_html, desktop_body_css,
    mobile_panel_html, mobile_panel_css).

    Desktop matches bundle: two absolute-positioned `.col` siblings (one
    per body shape, at the source sp[N] geometry), each containing
    `.uc-item` blocks with `.uct` heading + `.ucb` sub. Vertical gap
    between items is hand-calibrated in the bundle (2.96cqw) — we use
    the same value, since deriving it from PPT screenshots isn't
    something the OOXML carries.

    Mobile flattens both columns into one stacked list (`.uc-item-mobile`
    with `.uct-mobile` + `.ucb-mobile`), preceded by a single
    `.headline-mobile` and a `.header-row` containing only the logo.
    """
    cols = []
    for sp in body_sps:
        items = _group_continuation_items(_extract_paragraphs(sp, theme))
        cols.append((sp, items))

    heading_cqw = pt_to_cqw(heading_size_pt, slide_w_pt)
    body_cqw    = pt_to_cqw(body_size_pt, slide_w_pt)

    # Desktop body — two absolute-positioned columns
    col_blocks = []
    for col_idx, (sp, items) in enumerate(cols, start=1):
        pos = _box_to_pct(
            (sp.x_pt, sp.y_pt, sp.w_pt, sp.h_pt),
            slide_w_pt, slide_h_pt,
        )
        item_html = []
        for it in items:
            sub_text = " ".join(s for s in it["subs"] if s)
            uct = (
                f'<div class="uct">{_escape(it["heading"])}</div>'
                if it["heading"] else ""
            )
            ucb = f'<div class="ucb">{_escape(sub_text)}</div>' if sub_text else ""
            item_html.append(
                f'        <div class="uc-item">\n'
                f'          {uct}\n'
                f'          {ucb}\n'
                f'        </div>'
            )
        col_blocks.append(
            f'      <div class="col col-{col_idx}" '
            f'style="left:{pos["left"]}; top:{pos["top"]}; width:{pos["width"]};">\n'
            + "\n".join(item_html)
            + '\n      </div>'
        )
    desktop_html = "\n".join(col_blocks)

    desktop_css = (
        "#deck-desktop .col {\n"
        "  position: absolute;\n"
        "  display: flex; flex-direction: column;\n"
        "  gap: 2.96cqw;\n"
        "}\n"
        "#deck-desktop .uc-item {\n"
        "  display: flex; flex-direction: column; gap: 0.89cqw;\n"
        "}\n"
        "#deck-desktop .uct {\n"
        f"  font-size: {heading_cqw}; font-weight: 700; line-height: 1.15;\n"
        "  color: var(--headline);\n"
        "}\n"
        "#deck-desktop .ucb {\n"
        f"  font-size: {body_cqw}; font-weight: 300; line-height: 1.25;\n"
        "  color: var(--headline);\n"
        "}\n"
    )

    # Mobile — stacked across both columns, headline above
    flat_items: list[dict] = []
    for _sp, items in cols:
        flat_items.extend(items)

    item_lines = []
    for it in flat_items:
        sub_text = " ".join(s for s in it["subs"] if s)
        uct = f'<div class="uct-mobile">{_escape(it["heading"])}</div>' if it["heading"] else ""
        ucb = f'<div class="ucb-mobile">{_escape(sub_text)}</div>' if sub_text else ""
        item_lines.append(
            '      <div class="uc-item-mobile">\n'
            f'        {uct}\n'
            f'        {ucb}\n'
            '      </div>'
        )

    title_text = _escape(title_run["text"]) if title_run else ""
    logo_img = (
        f'<div class="gif-logo-mobile" style="margin-left:auto"><img src="{logo_url}" alt=""></div>'
        if logo_url else ""
    )
    mobile_html = (
        '    <div class="header-row">\n'
        f'      {logo_img}\n'
        '    </div>\n'
        f'    <div class="headline-mobile">{title_text}</div>\n'
        + "\n".join(item_lines)
    )

    mobile_css = (
        "  #deck-mobile .panel {\n"
        "    min-height: 100vh;\n"
        "    background: var(--bg-cyan);\n"
        "    color: var(--headline);\n"
        "    padding: 8vw 6vw 12vw;\n"
        "  }\n"
        "  #deck-mobile .header-row {\n"
        "    display: flex; justify-content: space-between;\n"
        "    align-items: flex-start; margin-bottom: 6vw;\n"
        "  }\n"
        "  #deck-mobile .gif-logo-mobile { width: 22vw; }\n"
        "  #deck-mobile .gif-logo-mobile img { width: 100%; }\n"
        "  #deck-mobile .headline-mobile {\n"
        "    font-size: 9vw; font-weight: 500; text-transform: uppercase;\n"
        "    line-height: 1.05; margin: 6vw 0 8vw; letter-spacing: 0.01em;\n"
        "  }\n"
        "  #deck-mobile .uc-item-mobile { margin-bottom: 5vw; }\n"
        "  #deck-mobile .uct-mobile {\n"
        "    font-size: 4.4vw; font-weight: 700;\n"
        "    line-height: 1.2; margin-bottom: 1vw;\n"
        "  }\n"
        "  #deck-mobile .ucb-mobile {\n"
        "    font-size: 4.2vw; font-weight: 300; line-height: 1.3;\n"
        "  }\n"
    )

    return (desktop_html, desktop_css, mobile_html, mobile_css)


def _build_plain_list(body_sps, theme: Theme,
                      slide_w_pt: float, slide_h_pt: float,
                      slide_class: SlideClass,
                      title_run, logo_url,
                      body_size_pt: float) -> tuple[str, str, str, str]:
    """Single-shape plain list. No canonical P&G case yet — defensive
    fallback. Returns the same 4-tuple shape as the other variants so
    `render_title_stats` can dispatch uniformly.
    """
    entries: list[dict] = []
    for sp in body_sps:
        entries.extend(_extract_paragraphs(sp, theme))

    bullet = slide_class.hints.get("bullet_char", "")
    bullet_span = (
        f'<span class="bullet">{_escape(bullet)}</span>' if bullet else ""
    )

    if body_sps:
        pos = _box_to_pct(_union_box_pt(body_sps), slide_w_pt, slide_h_pt)
    else:
        pos = {"left": "8%", "top": "28%", "width": "84%", "height": "54%"}

    body_cqw = pt_to_cqw(body_size_pt, slide_w_pt)

    lines = ['      <div class="body plain">']
    for e in entries:
        lines.append(
            f'        <div class="entry">{bullet_span}{_escape(e["text"])}</div>'
        )
    lines.append('      </div>')
    desktop_html = "\n".join(lines)

    desktop_css = (
        "#deck-desktop .body.plain {\n"
        "  position: absolute;\n"
        f"  left: {pos['left']}; top: {pos['top']};\n"
        f"  width: {pos['width']}; height: {pos['height']};\n"
        "  display: flex; flex-direction: column; gap: 0.4em;\n"
        "  overflow: hidden;\n"
        "}\n"
        "#deck-desktop .body.plain .entry {\n"
        f"  font-size: {body_cqw}; line-height: 1.3;\n"
        "  display: flex; align-items: baseline; gap: 0.4em;\n"
        "}\n"
        "#deck-desktop .body.plain .bullet { flex: 0 0 auto; opacity: 0.8; }\n"
    )

    # Mobile — single stacked list mirroring the continuation pattern
    title_text = _escape(title_run["text"]) if title_run else ""
    logo_img = (
        f'<div class="gif-logo-mobile" style="margin-left:auto"><img src="{logo_url}" alt=""></div>'
        if logo_url else ""
    )
    item_lines = [
        f'      <div class="uc-item-mobile">{bullet_span}{_escape(e["text"])}</div>'
        for e in entries
    ]
    mobile_html = (
        '    <div class="header-row">\n'
        f'      {logo_img}\n'
        '    </div>\n'
        f'    <div class="headline-mobile">{title_text}</div>\n'
        + "\n".join(item_lines)
    )
    mobile_css = (
        "  #deck-mobile .panel {\n"
        "    min-height: 100vh;\n"
        "    background: var(--bg-cyan);\n"
        "    color: var(--headline);\n"
        "    padding: 8vw 6vw 12vw;\n"
        "  }\n"
        "  #deck-mobile .header-row {\n"
        "    display: flex; justify-content: flex-end;\n"
        "    margin-bottom: 6vw;\n"
        "  }\n"
        "  #deck-mobile .gif-logo-mobile { width: 22vw; }\n"
        "  #deck-mobile .gif-logo-mobile img { width: 100%; }\n"
        "  #deck-mobile .headline-mobile {\n"
        "    font-size: 9vw; font-weight: 500; text-transform: uppercase;\n"
        "    line-height: 1.05; margin: 6vw 0 8vw;\n"
        "  }\n"
        "  #deck-mobile .uc-item-mobile {\n"
        "    font-size: 4.4vw; line-height: 1.3; margin-bottom: 4vw;\n"
        "  }\n"
        "  #deck-mobile .uc-item-mobile .bullet { opacity: 0.8; margin-right: 0.4em; }\n"
    )

    return (desktop_html, desktop_css, mobile_html, mobile_css)


# ────────────────────────────────────────────────────────────────────────────
# CSS + body assembly (frame: title, body wrapper, footnote, logo)
#
# Body container positioning is owned by each variant builder — it derives
# the box from the source shape geometry rather than a hardcoded percentage.

def _build_css(slide_w_pt: float, slide_h_pt: float,
               title_sp, footnote_sp, logo_pic,
               body_variant_css: str, mobile_panel_css: str,
               canvas_bg: str,
               title_run, footnote_run,
               footnote_size_cqw: str,
               logo_invert: bool) -> str:
    canvas_w_css, canvas_h_css = canvas_aspect_css(slide_w_pt, slide_h_pt)
    parts = []

    parts.append("/* ---------------- DESKTOP ---------------- */")
    parts.append("#deck-desktop { display: block; }")
    parts.append(
        "#deck-desktop .canvas {\n"
        f"  width: {canvas_w_css}; height: {canvas_h_css};\n"
        "  margin: 0 auto;\n"
        "  background: var(--bg);\n"
        "  position: relative; container-type: inline-size;\n"
        "  overflow: hidden;\n"
        "}"
    )
    parts.append("#deck-desktop .slide { position: absolute; inset: 0; color: var(--headline); }")

    if title_sp is not None:
        title_size_cqw = (
            pt_to_cqw(title_run["size_pt"], slide_w_pt)
            if title_run else "3.438cqw"
        )
        parts.append(_positioned("#deck-desktop .L", title_sp, slide_w_pt, slide_h_pt))
        parts.append(
            "#deck-desktop .L .t {\n"
            f"  font-size: {title_size_cqw}; font-weight: 400;\n"
            "  color: var(--headline); line-height: 1.05;\n"
            "  letter-spacing: 0.005em; text-transform: uppercase;\n"
            "}"
        )

    if logo_pic is not None:
        logo_filter = "  filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            f"#deck-desktop .gif-logo {{ position: absolute; "
            f"left: {pt_to_pct_x(logo_pic.x_pt, slide_w_pt)}; "
            f"top: {pt_to_pct_y(logo_pic.y_pt, slide_h_pt)}; "
            f"width: {pt_to_pct_x(logo_pic.w_pt, slide_w_pt)}; "
            f"height: {pt_to_pct_y(logo_pic.h_pt, slide_h_pt)};\n"
            f"{logo_filter}"
            "}"
        )
        parts.append("#deck-desktop .gif-logo img { width: 100%; height: 100%; object-fit: contain; }")

    parts.append(body_variant_css.rstrip())

    if footnote_run is not None and footnote_sp is not None:
        parts.append(_positioned("#deck-desktop .footnote", footnote_sp, slide_w_pt, slide_h_pt))
        parts.append(
            "#deck-desktop .footnote {\n"
            f"  font-size: {footnote_size_cqw}; font-style: italic; opacity: 0.7;\n"
            "  line-height: 1.3;\n"
            "}"
        )

    # Mobile — frame only. The variant owns the panel's content rules.
    parts.append("\n/* ---------------- MOBILE ---------------- */")
    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    parts.append("  #deck-desktop { display: none; }")
    parts.append("  #deck-mobile { display: block; background: var(--bg-cyan); }")
    if logo_invert:
        parts.append(
            "  #deck-mobile .gif-logo-mobile { filter: brightness(0) invert(1); }"
        )
    parts.append(mobile_panel_css.rstrip())
    parts.append("}")
    return "\n".join(parts)


def _positioned(selector: str, shape: FlatShape,
                slide_w_pt: float, slide_h_pt: float) -> str:
    return (
        f"{selector} {{ position: absolute; "
        f"left: {pt_to_pct_x(shape.x_pt, slide_w_pt)}; "
        f"top: {pt_to_pct_y(shape.y_pt, slide_h_pt)}; "
        f"width: {pt_to_pct_x(shape.w_pt, slide_w_pt)}; "
        f"height: {pt_to_pct_y(shape.h_pt, slide_h_pt)}; }}"
    )


def _build_body(title_run, desktop_body_html, footnote_run, logo_url,
                mobile_panel_html, logo_invert) -> str:
    """Assemble the complete body HTML — desktop + mobile sections.

    Desktop: canvas/slide wrappers + .L title + .gif-logo + variant body
    + optional .footnote. Mobile: variant-owned panel content (the
    variant builder already emits the bundle's per-variant top-row /
    header-row / item structure).
    """
    desktop_lines = []
    if title_run:
        desktop_lines.append(
            f'      <div class="L"><div class="t">{_escape(title_run["text"])}</div></div>'
        )
    if logo_url:
        desktop_lines.append(f'      <div class="gif-logo"><img src="{logo_url}" alt=""></div>')
    desktop_lines.append(desktop_body_html)
    if footnote_run:
        desktop_lines.append(f'      <div class="footnote">{_escape(footnote_run["text"])}</div>')

    return (
        "<!-- DESKTOP -->\n"
        "<section id=\"deck-desktop\">\n"
        "  <div class=\"canvas\">\n"
        "    <div class=\"slide\">\n"
        f"{chr(10).join(desktop_lines)}\n"
        "    </div>\n"
        "  </div>\n"
        "</section>\n\n"
        "<!-- MOBILE -->\n"
        "<section id=\"deck-mobile\">\n"
        "  <div class=\"panel\">\n"
        f"{mobile_panel_html}\n"
        "  </div>\n"
        "</section>"
    )


def _escape(s: str) -> str:
    return (
        s.replace("&", "&amp;")
         .replace("<", "&lt;")
         .replace(">", "&gt;")
    )
