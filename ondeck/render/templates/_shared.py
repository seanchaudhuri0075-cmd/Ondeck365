"""Shared helpers used across multiple template builders.

Per the NOTES.md guidance ("lift to a shared module when the third
template needs it, not before"), this is the home for cross-template
utilities.

Currently lifted:
  - is_logo_pic / image_src / inline_data_url       (cover, section_divider,
                                                     media_showcase, title_stats)
  - parse_gradient / resolve_color / shadow_to_cqw  (section_divider standard,
                                                     media_showcase large-headline)
  - extract_headline_runs                            (section_divider standard,
                                                     media_showcase large-headline)
  - TYPEFACE_WEIGHT (font-weight by typeface map)    (all 4 templates)
  - EMU_PER_PT (PowerPoint EMU → pt conversion)      (all callers of OOXML
                                                     dist/blur attributes)

Still pending lift on third-use: _data_url, _escape, _solidfill_hex,
_has_text — local to ≤2 templates so far. Lift on demand, not preemptively.
"""
from __future__ import annotations

import math
from pathlib import Path
from typing import Optional

from pptx.slide import Slide

from ...layout.detect import SlideClass
from ...parse.font_calibration import calibrate_size_pt
from ...parse.images import extract_image_ref
from ...parse.shapes import FlatShape
from ...parse.slide import NS
from ...parse.text import parse_text_frame
from ...parse.theme import Theme, apply_lum_mod_off
from ..css import color_with_alpha


# PowerPoint stores distances in English Metric Units (1pt = 9525 EMU).
# Used by shadow + gradient extraction below.
EMU_PER_PT = 9525

# Empirical typeface → font-weight map. Keys lowercased for case-insensitive
# lookup. Univers Condensed → 500 (Medium), Light → 300, Bold → 700; Barlow
# Condensed mirrors the same mapping (matched-metric substitution).
TYPEFACE_WEIGHT = {
    "univers condensed":        500,
    "univers condensed light":  300,
    "univers condensed bold":   700,
    "barlow condensed":         500,
    "barlow condensed light":   300,
    "barlow condensed bold":    700,
}


def is_logo_pic(
    shape,
    slide_w_pt: float,
    slide_h_pt: float,
    *,
    x_min_frac: float = 0.7,
    y_max_frac: float = 0.15,
    w_max_frac: float = 0.2,
) -> bool:
    """Geometry-based logo classifier: small <p:pic> in the top-right corner.

    Defaults match the P&G GIF logo placement (top-right corner, < 20%
    canvas width, sitting in the top 15% vertical band). Overrides via
    keyword args let future templates loosen / tighten thresholds without
    forking the helper. The thresholds were established by cover,
    section_divider, and media_showcase consistently using the same trio
    — lifting on the fourth use (title_stats).

    Returns False for shapes with no geometry (group descendants whose
    transform didn't resolve, etc.) and for non-pic shape kinds.
    """
    if shape.x_pt is None or shape.kind != "pic":
        return False
    return (
        shape.x_pt > slide_w_pt * x_min_frac
        and shape.y_pt < slide_h_pt * y_max_frac
        and shape.w_pt < slide_w_pt * w_max_frac
    )


def image_src(shape, slide_class: SlideClass, slide: Slide) -> Optional[str]:
    """Resolve an image shape to a src URL: manifest filename or inline data URL.

    Lookup order:
      1. The shape's OOXML cNvPr/@id is matched against entries in
         slide_class.media.get("images", []). On hit, returns the entry's
         filename (relative path, written by transform/image.py).
      2. On miss — typically because the shape is an SVG-only blip
         (transform/image.py skips SVGs by design) or transform hasn't
         run yet — falls back to inline base64 via parse/images.py
         extract_image_ref. SVG content_type passes through; the caller
         decides how to embed it.

    Returns None when neither a manifest entry nor a resolvable image blob
    is available (callers can decide whether to skip the shape or raise).

    The fallback path means this helper is forgiving — a missing manifest
    entry doesn't crash render. If you want strict mode (manifest required),
    check slide_class.media yourself before calling this.
    """
    cnv = shape.element.find("p:nvPicPr/p:cNvPr", NS)
    src_id = cnv.get("id") if cnv is not None else None
    if src_id is not None:
        for entry in slide_class.media.get("images", []):
            if entry.get("src_id") == src_id:
                return entry.get("filename")
    # Fallback: inline data URL (SVG path, or transform not yet run).
    return inline_data_url(shape, slide)


def inline_data_url(shape, slide: Slide) -> Optional[str]:
    """Return a data: URL for the shape's image blob — SVG preferred over raster.

    Used both (a) as the fallback path for image_src when the manifest
    doesn't carry an entry for the shape (typical for SVG-only shapes,
    since transform/image.py skips SVGs by design) and (b) as the
    direct call for shapes that should ALWAYS travel inline regardless
    of manifest entry — currently the GIF logo on every slide that
    has one. Logos are tiny (≈5KB raster) and inlining them keeps the
    HTML self-contained for AirDrop / single-file delivery to iOS,
    which is the primary mobile-preview path. Sibling-file references
    break under that delivery model because only the HTML transfers.

    The SVG-first ordering matches the cover template's prior
    _hero_data_url behavior — for shapes with BOTH an SVG extension
    and a PNG raster fallback (PowerPoint's mixed mode), the SVG is
    preferred for vector fidelity.
    """
    import base64

    from ...parse.svg import extract_svg_ref

    svg = extract_svg_ref(shape.element, slide)
    if svg is not None and svg.blob is not None:
        b64 = base64.b64encode(svg.blob).decode("ascii")
        return f"data:{svg.content_type};base64,{b64}"

    ref = extract_image_ref(shape.element, slide)
    if ref is None or ref.blob is None:
        return None
    b64 = base64.b64encode(ref.blob).decode("ascii")
    return f"data:{ref.content_type};base64,{b64}"


def resolve_inherited_size(shape, slide: Slide, paragraph_lvl: int = 0) -> Optional[float]:
    """Walk the OOXML font-size inheritance chain for a shape's text body.

    Returns the inherited size in pt, or None if no resolution found.

    The OOXML inheritance chain for run-level size_pt when <a:rPr sz=""> is
    absent (verified via slide-22 audit, 2026-05-02):

      1. Per-paragraph <a:pPr><a:defRPr sz="..."/> on the run's paragraph.
      2. Shape-level <p:txBody><a:lstStyle><a:lvl{N}pPr><a:defRPr sz="..."/>
         where N is 1-indexed paragraph level (0 → lvl1).
      3. Slide-layout text styles for placeholder shapes (rare; not present
         in this deck's slideLayout7 — most layouts inherit straight to
         master).
      4. Slide-master <p:txStyles>:
         - Shapes with <p:ph type="title"|"ctrTitle">: read <p:titleStyle>
         - Other placeholders (<p:ph> present, type != title): read <p:bodyStyle>
         - Freestanding text boxes (txBox="1", no <p:ph>): read <p:otherStyle>
         Within the chosen style, look up <a:lvl{N}pPr><a:defRPr sz="..."/>.

    Critical correction over the naïve "everything inherits from bodyStyle"
    assumption: txBox shapes resolve via otherStyle, NOT bodyStyle. Office
    default master has bodyStyle/lvl1=28pt vs otherStyle/lvl1=18pt. Using
    the wrong style overflows narrow text frames.
    """
    # Step 1: per-paragraph defRPr (rare, not in slide 22 — defensive)
    # Skipped: would require the paragraph element, not the shape; callers
    # that need per-paragraph resolution can handle that at the run-extraction site.

    # Step 2: shape-level lstStyle
    lst = shape.element.find("p:txBody/a:lstStyle", NS)
    if lst is not None:
        lvl_elem = lst.find(f"a:lvl{paragraph_lvl + 1}pPr", NS)
        if lvl_elem is not None:
            drpr = lvl_elem.find("a:defRPr", NS)
            if drpr is not None and drpr.get("sz"):
                return int(drpr.get("sz")) / 100

    # Step 3: layout-level text styles — skipped for now. None of the
    # currently-shipped slides use layouts that override body sizes.
    # When that day comes, slide.slide_layout.element / find p:txStyles.

    # Step 4: master-level <p:txStyles>
    try:
        master_elem = slide.slide_layout.slide_master.element
    except AttributeError:
        return None
    tx_styles = master_elem.find("p:txStyles", NS)
    if tx_styles is None:
        return None

    ph = shape.element.find("p:nvSpPr/p:nvPr/p:ph", NS)
    if ph is not None:
        ph_type = ph.get("type", "")
        if ph_type in ("title", "ctrTitle"):
            style = tx_styles.find("p:titleStyle", NS)
        else:
            style = tx_styles.find("p:bodyStyle", NS)
    else:
        # Freestanding txBox or shape without placeholder ref
        style = tx_styles.find("p:otherStyle", NS)

    if style is None:
        return None

    lvl_elem = style.find(f"a:lvl{paragraph_lvl + 1}pPr", NS)
    if lvl_elem is not None:
        drpr = lvl_elem.find("a:defRPr", NS)
        if drpr is not None and drpr.get("sz"):
            return int(drpr.get("sz")) / 100

    return None


def inline_optimized_data_url(
    shape,
    slide_class: SlideClass,
    slide: Slide,
    media_dir: Path,
) -> Optional[str]:
    """Return a data: URL using the optimized WebP from media_dir.

    Used by media_showcase photo-grid + large-headline variants where
    photos are content (not decoration) but the rendered HTML must
    still be self-contained for iOS Quick Look's file:// sandbox,
    which blocks all sibling-file fetches even when the files arrive
    in the same directory (confirmed 2026-05-01 — see NOTES.md
    Operating principles entry on AirDrop verification).

    Inlining the optimized WebP (already transcoded by
    transform/image.py at quality 85, with EXIF rotation and any
    per-image `flatten_on_canvas` applied) keeps HTML inflation
    bounded — typically 4× smaller than inlining the raw PPTX PNG
    via inline_data_url, since PowerPoint embeds full-resolution
    source images that the transform stage compresses.

    Strict contract: media_dir is required. Callers must pass the
    directory containing the optimized WebP files. Raises ValueError
    when media_dir is None — the silent raw-PPTX fallback was
    deliberately removed because it produced unbounded HTML inflation
    when a driver forgot to thread the param.

    Lookup order:
      1. cNvPr/@id matched against slide_class.media["images"] — on
         hit, reads {media_dir}/{filename} and returns
         `data:image/webp;base64,...`.
      2. Falls back to inline_data_url (raw PPTX blob) ONLY when the
         shape has no manifest entry or its file is missing on disk.
         This narrow fallback covers SVG-only blips that
         transform/image.py skips by design — it does NOT cover the
         no-media_dir case.
    """
    if media_dir is None:
        raise ValueError(
            "inline_optimized_data_url requires media_dir; "
            "the silent raw-PPTX fallback for missing media_dir was removed "
            "to prevent unbounded HTML inflation"
        )
    cnv = shape.element.find("p:nvPicPr/p:cNvPr", NS)
    src_id = cnv.get("id") if cnv is not None else None
    if src_id is not None:
        for entry in slide_class.media.get("images", []):
            if entry.get("src_id") == src_id:
                fn = entry.get("filename")
                if fn:
                    p = Path(media_dir) / fn
                    if p.exists():
                        import base64
                        b64 = base64.b64encode(p.read_bytes()).decode("ascii")
                        return f"data:image/webp;base64,{b64}"
    return inline_data_url(shape, slide)


# ────────────────────────────────────────────────────────────────────────────
# Gradient / shadow / headline extraction — lifted from section_divider.py
# at the third-use point per the NOTES.md rule. Used by both
# section_divider standard variant and media_showcase large-headline variant.

def parse_gradient(overlay_sp: FlatShape, theme: Theme) -> Optional[dict]:
    """Extract gradient stops + direction + base color from <a:gradFill>.

    Returns a dict: {color_css, color_zero_css, css_angle_deg, stops}
    where stops is a list of (pos_pct, alpha) — color is implicit (the base).
    Only the FIRST stop's color is used as the canonical overlay color;
    secondary stops carry the same hue and only differ in alpha.
    """
    grad = overlay_sp.element.find("p:spPr/a:gradFill", NS)
    if grad is None:
        return None

    base_hex = None
    stops = []
    for gs in grad.findall("a:gsLst/a:gs", NS):
        pos = int(gs.get("pos", "0")) / 1000  # per-mille → percent
        clr = resolve_color(gs, theme)
        if clr is None:
            continue
        hex_color, alpha = clr
        if base_hex is None:
            base_hex = hex_color
        # treat alpha=None as opaque
        a = 1.0 if alpha is None else alpha
        stops.append((pos, a))

    if base_hex is None:
        return None

    # Direction
    lin = grad.find("a:lin", NS)
    pptx_deg = (int(lin.get("ang", "0")) / 60000) if lin is not None else 0
    css_deg = (pptx_deg + 90) % 360

    r, g, b = (int(base_hex[i : i + 2], 16) for i in (0, 2, 4))
    color_css = f"#{base_hex.upper()}"
    color_zero_css = f"rgba({r},{g},{b},0)"

    # Build the stop list as CSS strings using rgba so per-stop alpha applies.
    stop_css = []
    for pos, a in stops:
        if a >= 1.0:
            stop_css.append(f"#{base_hex.upper()} {pos:.0f}%")
        else:
            stop_css.append(f"rgba({r},{g},{b},{a:.2f}) {pos:.0f}%")

    return {
        "color_css": color_css,
        "color_zero_css": color_zero_css,
        "css_angle_deg": css_deg,
        "stops": stops,
        "stop_css": stop_css,
    }


def resolve_color(gs_elem, theme: Theme) -> Optional[tuple]:
    """Resolve a gradient stop's color (srgbClr or schemeClr + modifiers) to (hex, alpha)."""
    srgb = gs_elem.find("a:srgbClr", NS)
    scheme = gs_elem.find("a:schemeClr", NS)
    if srgb is not None:
        base = srgb.get("val").upper()
        node = srgb
    elif scheme is not None:
        try:
            base = theme.resolve(scheme.get("val"))
        except AttributeError:
            return None
        node = scheme
    else:
        return None

    lum_mod = node.find("a:lumMod", NS)
    lum_off = node.find("a:lumOff", NS)
    base = apply_lum_mod_off(
        base,
        int(lum_mod.get("val")) if lum_mod is not None else None,
        int(lum_off.get("val")) if lum_off is not None else None,
    )

    alpha_node = node.find("a:alpha", NS)
    alpha = (int(alpha_node.get("val")) / 100000) if alpha_node is not None else None
    return (base, alpha)


def shadow_to_cqw(effect_lst, slide_w_pt: float) -> Optional[str]:
    """Convert <a:effectLst><a:outerShdw> to CSS text-shadow (cqw units)."""
    shdw = effect_lst.find("a:outerShdw", NS)
    if shdw is None:
        return None
    blur_pt = int(shdw.get("blurRad", "0")) / EMU_PER_PT
    dist_pt = int(shdw.get("dist", "0")) / EMU_PER_PT
    dir_deg = int(shdw.get("dir", "0")) / 60000
    rad = math.radians(dir_deg)
    off_x_pt = dist_pt * math.cos(rad)
    off_y_pt = dist_pt * math.sin(rad)

    color_node = shdw.find("a:srgbClr", NS) or shdw.find("a:prstClr", NS)
    alpha_node = color_node.find("a:alpha", NS) if color_node is not None else None
    alpha = (int(alpha_node.get("val")) / 100000) if alpha_node is not None else 1.0
    # Phase 1B: assume black drop shadows (matches all P&G section dividers)
    return (
        f"{off_x_pt / slide_w_pt * 100:.2f}cqw "
        f"{off_y_pt / slide_w_pt * 100:.2f}cqw "
        f"{blur_pt / slide_w_pt * 100:.2f}cqw "
        f"rgba(0,0,0,{alpha:.2f})"
    )


def extract_headline_runs(headline_sp: FlatShape, theme: Theme, slide_w_pt: float) -> list:
    """Return list of dicts: {text, size_pt, weight, color, shadow_cqw}.

    Phase 1B: in section dividers and media_showcase large-headline overlays,
    every run on the headline carries the same effect (drop shadow). We
    extract the shape's first effectLst and apply it to all runs uniformly.
    """
    tf = parse_text_frame(headline_sp.element, theme)
    if tf is None:
        return []

    first_effect = headline_sp.element.find(".//a:rPr/a:effectLst", NS)
    shadow_cqw_str = shadow_to_cqw(first_effect, slide_w_pt) if first_effect is not None else None

    runs = []
    for p in tf.paragraphs:
        for r in p.runs:
            if not r.text.strip():
                continue
            weight = TYPEFACE_WEIGHT.get((r.style.typeface or "").lower(), 500)
            if r.style.bold:
                weight = 700
            cal = calibrate_size_pt(r.style.typeface, r.style.size_pt, r.style.bold)
            color = color_with_alpha(r.style.color_hex or "FFFFFF", r.style.color_alpha)
            runs.append({
                "text": r.text.strip(),
                "size_pt": cal,
                "weight": weight,
                "color": color,
                "shadow_cqw": shadow_cqw_str,
            })
    return runs
