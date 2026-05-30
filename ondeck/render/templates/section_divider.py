"""Section divider template — two sub-variants.

  STANDARD variant (slides 4, 8, 10, 13, 18)
    z=0  pic   full-bleed photo backdrop (raster, slightly inset)
    z=1  sp    gradient overlay (covers part of canvas; readability layer)
    z=2  sp    headline text — 2-3 stacked ALL CAPS lines at 88pt
    z=N  pic   small logo top-right

  BADGE variant (slide 5 — WNBA)
    z=0..3  pic   4 photos in a 2×2 grid (full canvas coverage)
    z=N     pic   small logo top-right
    z=M     sp    cyan-filled "circle" (rendered with border-radius:50%)
    z=M+1   sp    small WNBA text positioned to overlap the circle

Variant detection: STANDARD has a gradient-fill sp; BADGE does not. Both
share the same six-template entry under `section_divider` per the locked
manifest schema.

Two PPTX features that other templates don't (yet) need:

  GRADIENT EXTRACTION
    <a:gradFill> on the overlay sp carries:
      - <a:gsLst><a:gs pos="...">  stops, in per-mille of the gradient line
      - <a:srgbClr> or <a:schemeClr> with optional lumMod/lumOff/alpha
      - <a:lin ang="...">  flow direction, degrees × 60000
    PPTX angle convention: 0=right, 90=down, 180=left, 270=up.
    CSS angle convention: 0=up, 90=right, 180=down, 270=left.
    Conversion: css_angle = (pptx_angle + 90) mod 360.

  DROP SHADOW EXTRACTION
    <a:outerShdw blurRad= dist= dir= alpha=> on a run's <a:effectLst>:
      - blurRad / dist in EMU (÷9525 → pt)
      - dir in degrees × 60000; PPTX math angle (0=right, 90=down)
    CSS text-shadow: <off-x> <off-y> <blur> <color>
      off-x = dist × cos(dir),  off-y = dist × sin(dir)
"""
from __future__ import annotations

import base64
from typing import Optional

from pptx.slide import Slide

from ..css import pt_to_pct_x, pt_to_pct_y, pt_to_cqw, color_with_alpha
from ..desktop import canvas_aspect_css, positioned_style
from ..html import render_page
from ...layout.detect import SlideClass
from ...parse.font_calibration import calibrate_size_pt
from ...parse.images import extract_image_ref
from ...parse.shapes import FlatShape, flatten_slide
from ...parse.slide import NS
from ...parse.text import parse_text_frame
from ...parse.theme import Theme
from ._shared import (
    TYPEFACE_WEIGHT,
    extract_headline_runs,
    image_src,
    inline_data_url,
    is_logo_pic,
    parse_gradient,
)


def render_section_divider(
    slide: Slide,
    theme: Theme,
    slide_index: int,
    slide_class: SlideClass,
    deck_name: str,
    slide_w_pt: float,
    slide_h_pt: float,
    deck_brand_color: str = "#000000",
) -> str:
    """Render a section_divider slide. Dispatches to the right variant.

    `deck_brand_color` fills any visible canvas area not covered by shapes.
    """
    flat = list(flatten_slide(slide))
    variant = _detect_variant(flat)
    if variant == "badge":
        return _render_badge_variant(
            slide, theme, flat, slide_index, slide_class,
            deck_name, slide_w_pt, slide_h_pt, deck_brand_color,
        )
    return _render_standard_variant(
        slide, theme, flat, slide_index, slide_class,
        deck_name, slide_w_pt, slide_h_pt, deck_brand_color,
    )


def _detect_variant(shapes) -> str:
    """STANDARD = has a gradient-fill sp (the readability overlay).
    BADGE    = no gradient sp; uses photo grid + solid-fill circle instead."""
    for s in shapes:
        if s.kind == "sp" and s.element.find("p:spPr/a:gradFill", NS) is not None:
            return "standard"
    return "badge"


def _render_standard_variant(
    slide, theme, flat, slide_index, slide_class,
    deck_name, slide_w_pt, slide_h_pt, deck_brand_color,
) -> str:
    photo, overlay, headline, logo = _classify(flat, slide_w_pt, slide_h_pt)

    photo_url = inline_data_url(photo, slide) if photo else None
    logo_url = inline_data_url(logo, slide) if logo else None
    logo_invert = bool(slide_class.hints.get("logo_invert", False))

    overlay_info = parse_gradient(overlay, theme) if overlay else None
    headline_runs = extract_headline_runs(headline, theme, slide_w_pt) if headline else []
    headline_class = slide_class.hints.get("headline_class", "")

    slide_css = _build_css(
        slide_w_pt, slide_h_pt,
        photo, overlay, headline, logo,
        overlay_info, headline_runs, headline_class, logo_invert,
    )
    body_html = _build_body(photo_url, logo_url, headline_runs, headline_class)

    overlay_color_css = overlay_info["color_css"] if overlay_info else "rgba(0,0,0,0)"
    overlay_zero_css = overlay_info["color_zero_css"] if overlay_info else "rgba(0,0,0,0)"

    return render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg":         deck_brand_color or "#000000",
            "headline":   "#FFFFFF",
            "overlay":    overlay_color_css,
            "overlay-0":  overlay_zero_css,
            "font-cond":  '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=body_html,
        slide_css=slide_css,
    )


# ────────────────────────────────────────────────────────────────────────────
# Role detection

def _classify(shapes, slide_w, slide_h):
    """Return (photo_pic, overlay_sp, headline_sp, logo_pic) by geometry + features."""
    photo = None
    overlay = None
    headline = None
    logo = None
    for s in shapes:
        if s.x_pt is None:
            continue
        if s.kind == "pic":
            if is_logo_pic(s, slide_w, slide_h):
                logo = s
            elif photo is None or s.w_pt * s.h_pt > photo.w_pt * photo.h_pt:
                photo = s
        elif s.kind == "sp":
            grad = s.element.find("p:spPr/a:gradFill", NS)
            if grad is not None:
                if overlay is None or s.w_pt * s.h_pt > overlay.w_pt * overlay.h_pt:
                    overlay = s
            elif _has_text(s.element):
                if headline is None:
                    headline = s
    return photo, overlay, headline, logo


def _has_text(sp_elem) -> bool:
    """Quick text-presence check (no full parse / no theme dependency)."""
    for t in sp_elem.findall(".//a:t", NS):
        if t.text and t.text.strip():
            return True
    return False


# Gradient/shadow/headline-run extraction now live in `_shared.py` (lifted at
# their third-use point per the NOTES.md "lift on third use" rule). Imported
# at the top of this module: parse_gradient, extract_headline_runs.

# ────────────────────────────────────────────────────────────────────────────
# CSS + body assembly

def _build_css(slide_w_pt, slide_h_pt, photo, overlay, headline, logo,
               overlay_info, headline_runs, headline_class,
               logo_invert: bool = False) -> str:
    canvas_w_css, canvas_h_css = canvas_aspect_css(slide_w_pt, slide_h_pt)
    parts = []

    parts.append("/* ---------------- DESKTOP ---------------- */")
    parts.append("#deck-desktop { display: block; }")
    parts.append(
        f"#deck-desktop .canvas {{\n"
        f"  position: relative;\n"
        f"  width: {canvas_w_css};\n"
        f"  height: {canvas_h_css};\n"
        f"  margin: 0 auto;\n"
        f"  container-type: inline-size;\n"
        f"  overflow: hidden;\n"
        f"  background: var(--bg);\n"
        f"}}"
    )
    parts.append("#deck-desktop .slide { position: absolute; inset: 0; }")

    if photo:
        parts.append(
            f"#deck-desktop .photo {{ {positioned_style(photo.x_pt, photo.y_pt, photo.w_pt, photo.h_pt, slide_w_pt, slide_h_pt)} }}\n"
            f"#deck-desktop .photo img {{ width: 100%; height: 100%; object-fit: cover; }}"
        )

    if overlay and overlay_info:
        gradient_css = (
            f"linear-gradient({overlay_info['css_angle_deg']:.0f}deg, "
            + ", ".join(overlay_info["stop_css"])
            + ")"
        )
        parts.append(
            f"#deck-desktop .overlay {{\n"
            f"  {positioned_style(overlay.x_pt, overlay.y_pt, overlay.w_pt, overlay.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  background: {gradient_css};\n"
            f"}}"
        )

    if logo:
        logo_filter = "  filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            f"#deck-desktop .gif-logo {{\n"
            f"  {positioned_style(logo.x_pt, logo.y_pt, logo.w_pt, logo.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"{logo_filter}"
            f"}}\n"
            f"#deck-desktop .gif-logo img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )

    if headline and headline_runs:
        parts.append(
            f"#deck-desktop .L {{\n"
            f"  position: absolute;\n"
            f"  left: {pt_to_pct_x(headline.x_pt, slide_w_pt)};\n"
            f"  top: {pt_to_pct_y(headline.y_pt, slide_h_pt)};\n"
            f"  width: {pt_to_pct_x(headline.w_pt, slide_w_pt)};\n"
            f"  height: auto;\n"
            f"  display: flex;\n"
            f"  flex-direction: column;\n"
            f"}}"
        )
        # All runs share styling for section_divider (same size, weight, color)
        first = headline_runs[0]
        shadow_decl = f"\n  text-shadow: {first['shadow_cqw']};" if first.get("shadow_cqw") else ""
        parts.append(
            f"#deck-desktop .L .t {{\n"
            f"  font-size: {pt_to_cqw(first['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {first['weight']};\n"
            f"  letter-spacing: 0.005em;\n"
            f"  line-height: 1.0;\n"
            f"  color: {first['color']};\n"
            f"  text-transform: uppercase;{shadow_decl}\n"
            f"}}"
        )

    # Mobile
    parts.append("\n/* ---------------- MOBILE ---------------- */")
    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    parts.append("  html, body { height: 100%; }")
    parts.append("  #deck-desktop { display: none; }")
    parts.append("  #deck-mobile { display: block; height: 100%; }")
    parts.append(
        "  #deck-mobile .panel {\n"
        "    height: 100svh;\n"
        "    height: 100dvh;\n"
        "    position: relative;\n"
        "    overflow: hidden;\n"
        "    background: var(--bg);\n"
        "  }"
    )
    parts.append(
        "  #deck-mobile .photo-mobile {\n"
        "    position: absolute; inset: 0;\n"
        "    background-size: cover; background-position: center;\n"
        "    z-index: 0;\n"
        "  }"
    )
    if overlay_info:
        parts.append(
            "  #deck-mobile .overlay-mobile {\n"
            "    position: absolute; inset: 0;\n"
            f"    background: linear-gradient(to top, var(--overlay) 0%, var(--overlay) 28%, var(--overlay-0) 65%);\n"
            "    z-index: 1;\n"
            "    pointer-events: none;\n"
            "  }"
        )
    if logo:
        logo_filter_m = "    filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            "  #deck-mobile .top-bar {\n"
            "    position: absolute;\n"
            "    top: 0; right: 0;\n"
            "    z-index: 3;\n"
            "    display: flex; justify-content: flex-end;\n"
            "    padding: 5vw;\n"
            "  }\n"
            "  #deck-mobile .gif-logo-mobile {\n"
            "    width: 22vw;\n"
            f"{logo_filter_m}"
            "  }\n"
            "  #deck-mobile .gif-logo-mobile img { width: 100%; }"
        )
    if headline_runs:
        first = headline_runs[0]
        shadow_m = ""
        if first.get("shadow_cqw"):
            shadow_m = f"\n    text-shadow: 2px 2px 6px rgba(0,0,0,0.3);"
        parts.append(
            "  #deck-mobile .headline-mobile {\n"
            "    position: absolute;\n"
            "    left: 0; right: 0; bottom: 16dvh;\n"
            "    padding: 0 1.5rem calc(env(safe-area-inset-bottom));\n"
            "    z-index: 2;\n"
            "  }\n"
            "  #deck-mobile .headline-mobile .t {\n"
            "    font-size: clamp(3rem, 13vw, 4rem);\n"
            f"    font-weight: {first['weight']};\n"
            "    line-height: 0.95;\n"
            "    text-transform: uppercase;"
            f"{shadow_m}\n"
            "  }"
        )
    parts.append("}")
    return "\n".join(parts)


def _build_body(photo_url, logo_url, headline_runs, headline_class) -> str:
    desktop_lines = []
    if photo_url:
        desktop_lines.append(f'      <div class="photo"><img src="{photo_url}" alt=""></div>')
    if any(True for _ in headline_runs):  # has overlay only when we have headline
        desktop_lines.append('      <div class="overlay"></div>')
    if headline_runs:
        cls = f' {headline_class}' if headline_class else ''
        desktop_lines.append(f'      <div class="L{cls}">')
        for r in headline_runs:
            desktop_lines.append(f'        <div class="t">{_escape(r["text"])}</div>')
        desktop_lines.append('      </div>')
    if logo_url:
        desktop_lines.append(f'      <div class="gif-logo"><img src="{logo_url}" alt=""></div>')

    mobile_lines = []
    if photo_url:
        mobile_lines.append(f'    <div class="photo-mobile" style=\'background-image: url("{photo_url}")\'></div>')
    if headline_runs:
        mobile_lines.append('    <div class="overlay-mobile"></div>')
    if logo_url:
        mobile_lines.append('    <div class="top-bar">')
        mobile_lines.append(f'      <div class="gif-logo-mobile"><img src="{logo_url}" alt=""></div>')
        mobile_lines.append('    </div>')
    if headline_runs:
        cls = f' {headline_class}' if headline_class else ''
        mobile_lines.append(f'    <div class="L L-mobile headline-mobile{cls}">')
        for r in headline_runs:
            mobile_lines.append(f'      <div class="t">{_escape(r["text"])}</div>')
        mobile_lines.append('    </div>')

    return (
        "<!-- DESKTOP -->\n"
        "<section id=\"deck-desktop\">\n"
        "  <div class=\"canvas\">\n"
        "    <div class=\"slide\">\n"
        f"{chr(10).join(desktop_lines)}\n"
        "    </div>\n"
        "  </div>\n"
        "</section>\n\n"
        "<!-- MOBILE -->\n"
        "<section id=\"deck-mobile\">\n"
        "  <div class=\"panel\">\n"
        f"{chr(10).join(mobile_lines)}\n"
        "  </div>\n"
        "</section>"
    )


def _data_url(shape: FlatShape, slide: Slide) -> Optional[str]:
    img_ref = extract_image_ref(shape.element, slide)
    if img_ref is None:
        return None
    b64 = base64.b64encode(img_ref.blob).decode("ascii")
    return f"data:{img_ref.content_type};base64,{b64}"


def _escape(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


# ────────────────────────────────────────────────────────────────────────────
# BADGE variant (slide 5)

def _render_badge_variant(
    slide, theme, flat, slide_index, slide_class,
    deck_name, slide_w_pt, slide_h_pt, deck_brand_color,
) -> str:
    """Render the badge sub-variant: 4-photo grid + cyan circle + small label."""
    canvas_bg = slide_class.hints.get("canvas_bg") or deck_brand_color or "#000000"
    photos, logo, circle, text_sp = _classify_badge(flat, slide_w_pt, slide_h_pt)

    photo_urls = [inline_data_url(p, slide) for p in photos]
    logo_url = inline_data_url(logo, slide) if logo else None
    logo_invert = bool(slide_class.hints.get("logo_invert", False))

    # Circle fill color from the sp's solidFill
    circle_color = _solidfill_hex(circle.element) if circle else deck_brand_color

    # Two text runs in the badge: "WNBA" (or equivalent — bold, larger) + "Theme design" (smaller, regular)
    headline_class = slide_class.hints.get("headline_class", "")
    text_runs = _extract_badge_runs(text_sp, theme, headline_class) if text_sp else []
    main_run = text_runs[0] if text_runs else None
    sub_run = text_runs[1] if len(text_runs) > 1 else None

    canvas_w_css, canvas_h_css = canvas_aspect_css(slide_w_pt, slide_h_pt)

    # CSS — desktop: 4 photo cells + circle (border-radius:50%) + L text overlay
    parts = []
    parts.append("/* ---------------- DESKTOP ---------------- */")
    parts.append("#deck-desktop { display: block; }")
    parts.append(
        f"#deck-desktop .canvas {{\n"
        f"  position: relative;\n"
        f"  width: {canvas_w_css};\n"
        f"  height: {canvas_h_css};\n"
        f"  margin: 0 auto;\n"
        f"  container-type: inline-size;\n"
        f"  overflow: hidden;\n"
        f"  background: var(--bg);\n"
        f"}}"
    )
    parts.append("#deck-desktop .slide { position: absolute; inset: 0; }")
    parts.append(
        "#deck-desktop .photo { position: absolute; overflow: hidden; }\n"
        "#deck-desktop .photo img { width: 100%; height: 100%; object-fit: cover; }"
    )
    for i, p in enumerate(photos):
        parts.append(
            f"#deck-desktop .photo-{i} {{ "
            f"left: {pt_to_pct_x(p.x_pt, slide_w_pt)}; "
            f"top: {pt_to_pct_y(p.y_pt, slide_h_pt)}; "
            f"width: {pt_to_pct_x(p.w_pt, slide_w_pt)}; "
            f"height: {pt_to_pct_y(p.h_pt, slide_h_pt)}; "
            f"}}"
        )
    if logo:
        logo_filter = "  filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            f"#deck-desktop .gif-logo {{\n"
            f"  {positioned_style(logo.x_pt, logo.y_pt, logo.w_pt, logo.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"{logo_filter}"
            f"}}\n"
            f"#deck-desktop .gif-logo img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )
    if circle:
        parts.append(
            f"#deck-desktop .circle {{\n"
            f"  {positioned_style(circle.x_pt, circle.y_pt, circle.w_pt, circle.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  background: #{circle_color};\n"
            f"  border-radius: 50%;\n"
            f"}}"
        )
    if text_sp and (main_run or sub_run):
        parts.append(
            f"#deck-desktop .L {{\n"
            f"  position: absolute;\n"
            f"  left: {pt_to_pct_x(text_sp.x_pt, slide_w_pt)};\n"
            f"  top: {pt_to_pct_y(text_sp.y_pt, slide_h_pt)};\n"
            f"  width: {pt_to_pct_x(text_sp.w_pt, slide_w_pt)};\n"
            f"  height: {pt_to_pct_y(text_sp.h_pt, slide_h_pt)};\n"
            f"  display: flex; flex-direction: column;\n"
            f"  align-items: center; justify-content: center;\n"
            f"  text-align: center;\n"
            f"}}"
        )
    if main_run:
        parts.append(
            f"#deck-desktop .L .t-wnba {{\n"
            f"  font-size: {pt_to_cqw(main_run['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {main_run['weight']};\n"
            f"  line-height: 1.0;\n"
            f"  color: {main_run['color']};\n"
            f"  letter-spacing: 0.01em;\n"
            f"}}"
        )
    if sub_run:
        parts.append(
            f"#deck-desktop .L .t-sub {{\n"
            f"  font-size: {pt_to_cqw(sub_run['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {sub_run['weight']};\n"
            f"  line-height: 1.1;\n"
            f"  color: {sub_run['color']};\n"
            f"  margin-top: 0.2cqw;\n"
            f"}}"
        )

    # Mobile: top-bar with logo, circle-row with badge, photo-grid stacked
    parts.append("\n/* ---------------- MOBILE ---------------- */")
    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    parts.append("  #deck-desktop { display: none; }")
    parts.append("  #deck-mobile { display: block; background: #fff; }")
    parts.append(
        "  #deck-mobile .panel {\n"
        "    min-height: 100vh;\n"
        "    background: #fff;\n"
        "    color: #111;\n"
        "    padding: 0;\n"
        "  }"
    )
    logo_filter_m = "    filter: brightness(0) invert(1);\n" if logo_invert else ""
    parts.append(
        "  #deck-mobile .top-bar {\n"
        "    display: flex; justify-content: flex-end;\n"
        "    padding: 5vw;\n"
        "    background: var(--bg-cyan);\n"
        "  }\n"
        "  #deck-mobile .gif-logo-mobile {\n"
        "    width: 22vw;\n"
        f"{logo_filter_m}"
        "  }\n"
        "  #deck-mobile .gif-logo-mobile img { width: 100%; }"
    )
    parts.append(
        "  #deck-mobile .circle-row {\n"
        "    display: flex; justify-content: center;\n"
        "    padding: 6vw 0 4vw;\n"
        "    background: #fff;\n"
        "  }\n"
        f"  #deck-mobile .circle-mobile {{\n"
        f"    background: #{circle_color};\n"
        f"    border-radius: 50%;\n"
        f"    width: 32vw; height: 32vw;\n"
        f"    display: flex; flex-direction: column;\n"
        f"    align-items: center; justify-content: center;\n"
        f"    color: var(--headline);\n"
        f"  }}\n"
        "  #deck-mobile .circle-mobile .t-wnba { font-size: 8vw; font-weight: 700; line-height: 1.0; }\n"
        "  #deck-mobile .circle-mobile .t-sub { font-size: 4vw; font-weight: 400; margin-top: 1vw; }"
    )
    parts.append(
        "  #deck-mobile .photo-grid-mobile { display: flex; flex-direction: column; }\n"
        "  #deck-mobile .photo-grid-mobile .photo-m { width: 100%; }\n"
        "  #deck-mobile .photo-grid-mobile .photo-m img { width: 100%; height: auto; display: block; }"
    )
    parts.append("}")
    slide_css = "\n".join(parts)

    # Body
    desktop_lines = []
    for i, url in enumerate(photo_urls):
        if url:
            desktop_lines.append(f'      <div class="photo photo-{i}"><img src="{url}" alt=""></div>')
    if logo_url:
        desktop_lines.append(f'      <div class="gif-logo"><img src="{logo_url}" alt=""></div>')
    if circle:
        desktop_lines.append('      <div class="circle"></div>')
    if text_sp and (main_run or sub_run):
        desktop_lines.append('      <div class="L">')
        if main_run:
            desktop_lines.append(f'        <div class="t t-wnba">{_escape(main_run["text"])}</div>')
        if sub_run:
            desktop_lines.append(f'        <div class="t t-sub">{_escape(sub_run["text"])}</div>')
        desktop_lines.append('      </div>')

    mobile_lines = []
    mobile_lines.append('    <div class="top-bar">')
    if logo_url:
        mobile_lines.append(f'      <div class="gif-logo-mobile"><img src="{logo_url}" alt=""></div>')
    mobile_lines.append('    </div>')
    if main_run or sub_run:
        mobile_lines.append('    <div class="circle-row">')
        mobile_lines.append('      <div class="circle-mobile L L-mobile">')
        if main_run:
            mobile_lines.append(f'        <div class="t t-wnba">{_escape(main_run["text"])}</div>')
        if sub_run:
            mobile_lines.append(f'        <div class="t t-sub">{_escape(sub_run["text"])}</div>')
        mobile_lines.append('      </div>')
        mobile_lines.append('    </div>')
    if photo_urls:
        mobile_lines.append('    <div class="photo-grid-mobile">')
        for url in photo_urls:
            if url:
                mobile_lines.append(f'      <div class="photo-m"><img src="{url}" alt=""></div>')
        mobile_lines.append('    </div>')

    body_html = (
        "<!-- DESKTOP -->\n"
        "<section id=\"deck-desktop\">\n"
        "  <div class=\"canvas\">\n"
        "    <div class=\"slide\">\n"
        f"{chr(10).join(desktop_lines)}\n"
        "    </div>\n"
        "  </div>\n"
        "</section>\n\n"
        "<!-- MOBILE -->\n"
        "<section id=\"deck-mobile\">\n"
        "  <div class=\"panel\">\n"
        f"{chr(10).join(mobile_lines)}\n"
        "  </div>\n"
        "</section>"
    )

    return render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg":         canvas_bg,
            "bg-cyan":    deck_brand_color or "#000000",
            "headline":   "#FFFFFF",
            "font-cond":  '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=body_html,
        slide_css=slide_css,
    )


def _classify_badge(shapes, slide_w, slide_h):
    """Return (photos[], logo_pic, circle_sp, text_sp) for the badge variant."""
    photos = []
    logo = None
    circle = None
    text_sp = None
    for s in shapes:
        if s.x_pt is None:
            continue
        if s.kind == "pic":
            if is_logo_pic(s, slide_w, slide_h):
                logo = s
            else:
                photos.append(s)
        elif s.kind == "sp":
            srgb = s.element.find("p:spPr/a:solidFill/a:srgbClr", NS)
            if srgb is not None and circle is None:
                circle = s
            elif _has_text(s.element):
                if text_sp is None:
                    text_sp = s
    # Photos sorted by (y, x) so the grid CSS classes (.photo-0..3) match
    # left-to-right, top-to-bottom reading order.
    photos.sort(key=lambda p: (round(p.y_pt), round(p.x_pt)))
    return photos, logo, circle, text_sp


def _solidfill_hex(sp_elem) -> str:
    srgb = sp_elem.find("p:spPr/a:solidFill/a:srgbClr", NS)
    return srgb.get("val", "000000").upper() if srgb is not None else "000000"


def _extract_badge_runs(text_sp: FlatShape, theme: Theme, headline_class: str = "") -> list:
    """Two runs expected: main (e.g. 'WNBA' bold) + sub (e.g. 'Theme design' regular).

    `headline_class` is the manifest hint. Slide 5's `t-wnba` is a deck-author
    override: the bundle hand-picked 22pt for the inherited-size bold run on
    that one slide. The OOXML inheritance chain converges on 18pt for every
    inherited-size run in the deck (master <p:txStyles>/<p:otherStyle>/<a:lvl1pPr>
    and presentation <p:defaultTextStyle> both declare sz=1800), so this
    can't be derived structurally. See NOTES.md entry on inherited-size
    convergence.
    """
    tf = parse_text_frame(text_sp.element, theme)
    if tf is None:
        return []
    runs = []
    for p in tf.paragraphs:
        for r in p.runs:
            if not r.text.strip():
                continue
            weight = TYPEFACE_WEIGHT.get((r.style.typeface or "").lower(), 500)
            if r.style.bold:
                weight = 700
            cal = calibrate_size_pt(r.style.typeface, r.style.size_pt, r.style.bold)
            if (
                headline_class == "t-wnba"
                and r.style.size_pt is None
                and r.style.bold
            ):
                cal = 22.0
            color = color_with_alpha(r.style.color_hex or "FFFFFF", r.style.color_alpha)
            runs.append({
                "text": r.text.strip(),
                "size_pt": cal,
                "weight": weight,
                "color": color,
            })
    return runs
