"""Media showcase template — photo/video grids with optional badge overlay.

Phase 1B scope: the 4-photo 2×2 grid + cyan circle + small text label
sub-variant (slide 6). Slide-7 video and slides 11/12/14 small-corner-label
sub-variants come in a follow-up; for now only the badge-overlay path is
wired up.

Structurally near-identical to section_divider's BADGE variant — same
shape archetype (4 photos, logo top-right, solid-fill circle, txBox label).
The differences are stylistic and per-slide, so the inherited-bold size
and the bold-run CSS class are driven by the manifest hint
`headline_class` (e.g. slide 6 → "t-bold" → 24pt). See NOTES.md entry on
inherited-size convergence: every inherited-size run resolves to 18pt
through the OOXML chain, so any deviation from 18pt is a deck-author
override that lives in the manifest, not in PPTX walk logic.
"""
from __future__ import annotations

import base64
from pathlib import Path
from typing import Optional

from pptx.slide import Slide

from ..css import pt_to_pct_x, pt_to_pct_y, pt_to_cqw, color_with_alpha
from ..desktop import canvas_aspect_css, positioned_style
from ..html import render_page
from ...layout.detect import SlideClass
from ...parse.font_calibration import calibrate_size_pt
from ...parse.images import extract_image_ref
from ._shared import (
    extract_headline_runs,
    inline_data_url,
    inline_optimized_data_url,
    is_logo_pic,
    parse_gradient,
)
from ...parse.media import is_video_pic, extract_video
from ...parse.shapes import FlatShape, flatten_slide
from ...parse.slide import NS
from ...parse.text import parse_text_frame
from ...parse.theme import Theme

TYPEFACE_WEIGHT = {
    "univers condensed":        500,
    "univers condensed light":  300,
    "univers condensed bold":   700,
    "barlow condensed":         500,
    "barlow condensed light":   300,
    "barlow condensed bold":    700,
}

# Manifest hint → inherited-bold pt for media_showcase overlay labels.
# Keys are the `headline_class` string the deck author put in the manifest.
# Each entry is a deck-author override that the OOXML chain can't produce
# (every inherited-size run in this deck resolves to 18pt at the master).
INHERITED_BOLD_PT_BY_HINT = {
    "t-bold": 24.0,   # slide 6 (Wicked 2 / Promo), slide 9 (Concept / Adapts)
}


def render_media_showcase(
    slide: Slide,
    theme: Theme,
    slide_index: int,
    slide_class: SlideClass,
    deck_name: str,
    slide_w_pt: float,
    slide_h_pt: float,
    deck_brand_color: str = "#000000",
    media_dir: Optional[Path] = None,
) -> tuple:
    """Render a media_showcase slide.

    Returns (html, aux_files) where aux_files is a list of
    (filename_relative_to_html, blob) tuples the caller writes alongside
    the HTML. Empty list when no auxiliary files are needed (photo-grid
    slides). Video slides emit one entry: the .mp4 binary, since inlining
    multi-megabyte video as a data URL would balloon the HTML.

    Phase 1B variants:
      * photo-grid badge-overlay (slide 6) — 4 pics in 2×2
      * video full-bleed (slide 7)         — single <p:pic>+<p14:media>
    """
    flat = list(flatten_slide(slide))
    video_pic = next(
        (s for s in flat if s.kind == "pic" and is_video_pic(s.element)),
        None,
    )
    if video_pic is not None:
        return _render_video_variant(
            slide, theme, flat, video_pic,
            slide_index, slide_class, deck_name,
            slide_w_pt, slide_h_pt, deck_brand_color,
        )
    # Large-headline overlay sub-variant: a sp with <a:gradFill> is the
    # backdrop tinting layer for an 88pt stacked headline (slides 15, 17).
    # No other media_showcase variant uses gradFill — badge-overlay uses
    # solidFill+ellipse, video uses <p:pic>+<p14:media>. Detection by
    # gradient presence is unambiguous within media_showcase.
    has_gradient = any(
        s.kind == "sp" and s.element.find("p:spPr/a:gradFill", NS) is not None
        for s in flat
    )
    if has_gradient:
        return _render_large_headline_variant(
            slide, theme, flat,
            slide_index, slide_class, deck_name,
            slide_w_pt, slide_h_pt, deck_brand_color, media_dir,
        )
    # Sub-variant hint dispatch — for layouts the OOXML can't disambiguate
    # from the badge-overlay default. Currently: "titled-grid" (slides 19,
    # 20 — title text above photo grid, no circle, no logo). Default
    # fallback (no hint) preserves photo-grid badge behavior for slides
    # 6, 9, 11, 12.
    if slide_class.hints.get("media_showcase_layout") == "titled-grid":
        return _render_titled_grid_variant(
            slide, theme, flat,
            slide_index, slide_class, deck_name,
            slide_w_pt, slide_h_pt, deck_brand_color, media_dir,
        )
    return _render_photo_grid_variant(
        slide, theme, flat,
        slide_index, slide_class, deck_name,
        slide_w_pt, slide_h_pt, deck_brand_color, media_dir,
    )


def _render_photo_grid_variant(
    slide, theme, flat, slide_index, slide_class,
    deck_name, slide_w_pt, slide_h_pt, deck_brand_color, media_dir,
) -> tuple:
    photos, logo, circle, text_sp = _classify(flat, slide_w_pt, slide_h_pt)

    photo_urls = [inline_optimized_data_url(p, slide_class, slide, media_dir) for p in photos]
    logo_url = inline_data_url(logo, slide) if logo else None
    circle_color = _solidfill_hex(circle.element) if circle else deck_brand_color
    logo_invert = bool(slide_class.hints.get("logo_invert", False))

    headline_class = slide_class.hints.get("headline_class", "t-bold")
    inherited_bold_pt = INHERITED_BOLD_PT_BY_HINT.get(headline_class, 24.0)
    text_runs = _extract_overlay_runs(text_sp, theme, inherited_bold_pt) if text_sp else []

    # Paragraph count drives the mobile circle/label sizing branch (compact
    # vs default). Derived from text_sp's <a:p> children, ignoring blank
    # paragraphs. 2 paras → default 38vw circle (slide 6); ≥3 paras → compact
    # 44vw circle with smaller text + padding (slides 11, 12 — more lines
    # need more room).
    n_label_paragraphs = 0
    if text_sp is not None:
        tf = parse_text_frame(text_sp.element, theme)
        if tf is not None:
            n_label_paragraphs = sum(
                1 for p in tf.paragraphs
                if any(r.text.strip() for r in p.runs)
            )

    # Mobile photo layout — hint > count-based default. Default catches the
    # common case (5+ photos arranged 2-col); hint overrides when the
    # default is wrong for the slide's photo aspects (e.g. slide 11 has 8
    # photos including 3 portraits, cramped in grid-2col → "stack" hint).
    mobile_photo_layout = slide_class.hints.get("mobile_photo_layout") or (
        "grid-2col" if len(photos) >= 5 else "stack"
    )

    canvas_bg = slide_class.hints.get("canvas_bg") or deck_brand_color or "#000000"

    slide_css = _build_css(
        slide_w_pt, slide_h_pt,
        photos, logo, circle, text_sp,
        text_runs, circle_color, headline_class, logo_invert,
        n_label_paragraphs,
        mobile_photo_layout,
    )
    body_html = _build_body(photo_urls, logo_url, circle, text_runs, headline_class)

    html = render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg":        canvas_bg,
            "bg-cyan":   deck_brand_color or "#000000",
            "headline":  "#FFFFFF",
            "font-cond": '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=body_html,
        slide_css=slide_css,
    )
    return (html, [])


# ────────────────────────────────────────────────────────────────────────────
# Role detection

def _classify(shapes, slide_w, slide_h):
    """Return (photos[], logo_pic, circle_sp, text_sp).

    Same heuristics as section_divider's _classify_badge — small top-right pic
    is the brand logo; remaining pics are content photos sorted into reading
    order; sp with srgbClr solidFill is the circle; sp with text is the
    overlay label.

    Logo is FIRST-MATCH-ONLY: in document order, the first pic that
    satisfies `is_logo_pic` becomes the logo; any subsequent pic that
    happens to fall in the top-right + small-width zone (e.g. slide 11's
    pic[4] at 86.4%×13.8%, a small product photo) is treated as a photo.
    Without this, two-logo-match overwrites the real GIF logo with the
    false-positive photo and the real logo gets lost entirely. Tightening
    `is_logo_pic` itself would risk regressing locked slides 6/7 whose
    real logos sit at the geometric edge of the threshold.
    """
    photos = []
    logo = None
    circle = None
    text_sp = None
    for s in shapes:
        if s.x_pt is None:
            continue
        if s.kind == "pic":
            if logo is None and is_logo_pic(s, slide_w, slide_h):
                logo = s
            else:
                photos.append(s)
        elif s.kind == "sp":
            srgb = s.element.find("p:spPr/a:solidFill/a:srgbClr", NS)
            if srgb is not None and circle is None:
                circle = s
            elif _has_text(s.element):
                if text_sp is None:
                    text_sp = s
    photos.sort(key=lambda p: (round(p.y_pt), round(p.x_pt)))
    return photos, logo, circle, text_sp


def _has_text(sp_elem) -> bool:
    for t in sp_elem.findall(".//a:t", NS):
        if t.text and t.text.strip():
            return True
    return False


# ────────────────────────────────────────────────────────────────────────────
# Run extraction

def _extract_overlay_runs(text_sp: FlatShape, theme: Theme, inherited_bold_pt: float) -> list:
    """One dict per run, in document order. Class is bold→'t-bold', else 't-sub'."""
    tf = parse_text_frame(text_sp.element, theme)
    if tf is None:
        return []
    runs = []
    for p in tf.paragraphs:
        for r in p.runs:
            if not r.text.strip():
                continue
            weight = TYPEFACE_WEIGHT.get((r.style.typeface or "").lower(), 500)
            if r.style.bold:
                weight = 700
            cal = calibrate_size_pt(r.style.typeface, r.style.size_pt, r.style.bold)
            if r.style.size_pt is None and r.style.bold:
                cal = inherited_bold_pt
            color = color_with_alpha(r.style.color_hex or "FFFFFF", r.style.color_alpha)
            runs.append({
                "text": r.text.strip(),
                "size_pt": cal,
                "weight": weight,
                "color": color,
                "klass": "t-bold" if r.style.bold else "t-sub",
            })
    return runs


# ────────────────────────────────────────────────────────────────────────────
# CSS + body assembly

def _build_css(slide_w_pt, slide_h_pt, photos, logo, circle, text_sp,
               text_runs, circle_color, headline_class,
               logo_invert: bool = False,
               n_label_paragraphs: int = 0,
               mobile_photo_layout: str = "stack") -> str:
    canvas_w_css, canvas_h_css = canvas_aspect_css(slide_w_pt, slide_h_pt)
    parts = []

    parts.append("/* ---------------- DESKTOP ---------------- */")
    parts.append("#deck-desktop { display: block; }")
    parts.append(
        f"#deck-desktop .canvas {{\n"
        f"  position: relative;\n"
        f"  width: {canvas_w_css};\n"
        f"  height: {canvas_h_css};\n"
        f"  margin: 0 auto;\n"
        f"  container-type: inline-size;\n"
        f"  overflow: hidden;\n"
        f"  background: var(--bg);\n"
        f"}}"
    )
    parts.append("#deck-desktop .slide { position: absolute; inset: 0; }")
    parts.append(
        "#deck-desktop .photo { position: absolute; overflow: hidden; }\n"
        "#deck-desktop .photo img { width: 100%; height: 100%; object-fit: cover; }"
    )
    for i, p in enumerate(photos):
        parts.append(
            f"#deck-desktop .photo-{i} {{ "
            f"left: {pt_to_pct_x(p.x_pt, slide_w_pt)}; "
            f"top: {pt_to_pct_y(p.y_pt, slide_h_pt)}; "
            f"width: {pt_to_pct_x(p.w_pt, slide_w_pt)}; "
            f"height: {pt_to_pct_y(p.h_pt, slide_h_pt)}; "
            f"}}"
        )
    if logo:
        logo_filter = "  filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            f"#deck-desktop .gif-logo {{\n"
            f"  {positioned_style(logo.x_pt, logo.y_pt, logo.w_pt, logo.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"{logo_filter}"
            f"}}\n"
            f"#deck-desktop .gif-logo img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )
    if circle:
        parts.append(
            f"#deck-desktop .circle {{\n"
            f"  {positioned_style(circle.x_pt, circle.y_pt, circle.w_pt, circle.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  background: #{circle_color};\n"
            f"  border-radius: 50%;\n"
            f"}}"
        )
        # The .L overlay sits ON the circle (same rect) and centers its text via flex.
        parts.append(
            f"#deck-desktop .L {{\n"
            f"  {positioned_style(circle.x_pt, circle.y_pt, circle.w_pt, circle.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  display: flex; flex-direction: column;\n"
            f"  align-items: center; justify-content: center;\n"
            f"  text-align: center;\n"
            f"  padding: 1.5cqw;\n"
            f"}}"
        )

    bold_run = next((r for r in text_runs if r["klass"] == "t-bold"), None)
    sub_run = next((r for r in text_runs if r["klass"] == "t-sub"), None)
    if bold_run:
        parts.append(
            f"#deck-desktop .L .t-bold {{\n"
            f"  font-size: {pt_to_cqw(bold_run['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {bold_run['weight']};\n"
            f"  line-height: 1.0;\n"
            f"  color: {bold_run['color']};\n"
            f"  letter-spacing: 0.005em;\n"
            f"}}"
        )
    if sub_run:
        parts.append(
            f"#deck-desktop .L .t-sub {{\n"
            f"  font-size: {pt_to_cqw(sub_run['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {sub_run['weight']};\n"
            f"  line-height: 1.1;\n"
            f"  color: {sub_run['color']};\n"
            f"  margin-top: 0.4cqw;\n"
            f"}}"
        )

    # Mobile: top-bar with logo, circle-row badge, photo-grid stacked
    parts.append("\n/* ---------------- MOBILE ---------------- */")
    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    parts.append("  #deck-desktop { display: none; }")
    parts.append("  #deck-mobile { display: block; background: #fff; }")
    parts.append(
        "  #deck-mobile .panel {\n"
        "    min-height: 100vh;\n"
        "    background: #fff;\n"
        "    color: #111;\n"
        "    padding: 0;\n"
        "  }"
    )
    # Mobile top-bar background is always var(--bg-cyan), so a present
    # logo is always rendered against cyan. The brightness(0)+invert(1)
    # filter (logo → solid white) is therefore unconditional on mobile
    # whenever a logo exists. The `logo_invert` hint controls DESKTOP
    # only — desktop canvas color varies per slide. See NOTES.md
    # "logo_invert manifest hint" entry for the bundle-author error
    # context that motivates the desktop hint.
    logo_filter_m = "    filter: brightness(0) invert(1);\n" if logo is not None else ""
    # When no logo, the top-bar would collapse to padding-only height (~10vw)
    # and visually weak. Bundle slide 12 sets explicit height: 14vw to give
    # it presence; we replicate.
    top_bar_height = "    height: 14vw;\n" if logo is None else ""
    parts.append(
        "  #deck-mobile .top-bar {\n"
        "    display: flex; justify-content: flex-end;\n"
        "    padding: 5vw;\n"
        "    background: var(--bg-cyan);\n"
        f"{top_bar_height}"
        "  }\n"
        "  #deck-mobile .gif-logo-mobile {\n"
        "    width: 22vw;\n"
        f"{logo_filter_m}"
        "  }\n"
        "  #deck-mobile .gif-logo-mobile img { width: 100%; }"
    )

    # Circle/label mobile sizing — paragraph-count driven (see NOTES.md
    # auto-detect rule): ≥3 paragraphs → compact (44vw circle, smaller text,
    # padding, larger line-heights); else default (38vw, slide-6 baseline).
    compact = n_label_paragraphs >= 3
    if compact:
        circle_block = (
            "  #deck-mobile .circle-row {\n"
            "    display: flex; justify-content: center;\n"
            "    padding: 6vw 0 4vw;\n"
            "    background: #fff;\n"
            "  }\n"
            f"  #deck-mobile .circle-mobile {{\n"
            f"    background: #{circle_color};\n"
            f"    border-radius: 50%;\n"
            f"    width: 44vw; height: 44vw;\n"
            f"    display: flex; flex-direction: column;\n"
            f"    align-items: center; justify-content: center;\n"
            f"    color: var(--headline);\n"
            f"    text-align: center;\n"
            f"    padding: 4vw;\n"
            f"  }}\n"
            "  #deck-mobile .circle-mobile .t-bold { font-size: 4.4vw; font-weight: 700; line-height: 1.1; }\n"
            "  #deck-mobile .circle-mobile .t-sub { font-size: 3.6vw; font-weight: 400; margin-top: 1vw; line-height: 1.2; }"
        )
    else:
        circle_block = (
            "  #deck-mobile .circle-row {\n"
            "    display: flex; justify-content: center;\n"
            "    padding: 6vw 0 4vw;\n"
            "    background: #fff;\n"
            "  }\n"
            f"  #deck-mobile .circle-mobile {{\n"
            f"    background: #{circle_color};\n"
            f"    border-radius: 50%;\n"
            f"    width: 38vw; height: 38vw;\n"
            f"    display: flex; flex-direction: column;\n"
            f"    align-items: center; justify-content: center;\n"
            f"    color: var(--headline);\n"
            f"  }}\n"
            "  #deck-mobile .circle-mobile .t-bold { font-size: 6vw; font-weight: 700; line-height: 1.0; }\n"
            "  #deck-mobile .circle-mobile .t-sub { font-size: 3.4vw; font-weight: 400; margin-top: 1vw; }"
        )
    parts.append(circle_block)

    # Mobile photo arrangement — manifest hint authoritative; caller
    # resolves count-based default ("grid-2col" if ≥5 photos else "stack")
    # before passing in. Hint overrides for portrait-photo slides (e.g.
    # slide 11) where the count-based default picks grid-2col but
    # photos are cramped under aspect-ratio:1/1 + object-fit:contain.
    if mobile_photo_layout == "grid-2col":
        parts.append(
            "  #deck-mobile .photo-grid-mobile {\n"
            "    display: grid;\n"
            "    grid-template-columns: 1fr 1fr;\n"
            "    gap: 1vw; padding: 1vw; background: #fff;\n"
            "  }\n"
            "  #deck-mobile .photo-grid-mobile .photo-m {\n"
            "    width: 100%; background: #fff;\n"
            "    aspect-ratio: 1/1; overflow: hidden;\n"
            "  }\n"
            "  #deck-mobile .photo-grid-mobile .photo-m img {\n"
            "    width: 100%; height: 100%;\n"
            "    object-fit: contain; display: block;\n"
            "  }"
        )
    else:
        parts.append(
            "  #deck-mobile .photo-grid-mobile { display: flex; flex-direction: column; }\n"
            "  #deck-mobile .photo-grid-mobile .photo-m { width: 100%; }\n"
            "  #deck-mobile .photo-grid-mobile .photo-m img { width: 100%; height: auto; display: block; }"
        )
    parts.append("}")
    return "\n".join(parts)


def _build_body(photo_urls, logo_url, circle, text_runs, headline_class) -> str:
    desktop_lines = []
    for i, url in enumerate(photo_urls):
        if url:
            desktop_lines.append(f'      <div class="photo photo-{i}"><img src="{url}" alt=""></div>')
    if logo_url:
        desktop_lines.append(f'      <div class="gif-logo"><img src="{logo_url}" alt=""></div>')
    if circle:
        desktop_lines.append('      <div class="circle"></div>')
    if text_runs:
        desktop_lines.append('      <div class="L">')
        for r in text_runs:
            desktop_lines.append(f'        <div class="t {r["klass"]}">{_escape(r["text"])}</div>')
        desktop_lines.append('      </div>')

    mobile_lines = []
    mobile_lines.append('    <div class="top-bar">')
    if logo_url:
        mobile_lines.append(f'      <div class="gif-logo-mobile"><img src="{logo_url}" alt=""></div>')
    mobile_lines.append('    </div>')
    if text_runs:
        mobile_lines.append('    <div class="circle-row">')
        mobile_lines.append('      <div class="circle-mobile L L-mobile">')
        for r in text_runs:
            mobile_lines.append(f'        <div class="t {r["klass"]}">{_escape(r["text"])}</div>')
        mobile_lines.append('      </div>')
        mobile_lines.append('    </div>')
    if photo_urls:
        mobile_lines.append('    <div class="photo-grid-mobile">')
        for url in photo_urls:
            if url:
                mobile_lines.append(f'      <div class="photo-m"><img src="{url}" alt=""></div>')
        mobile_lines.append('    </div>')

    return (
        "<!-- DESKTOP -->\n"
        "<section id=\"deck-desktop\">\n"
        "  <div class=\"canvas\">\n"
        "    <div class=\"slide\">\n"
        f"{chr(10).join(desktop_lines)}\n"
        "    </div>\n"
        "  </div>\n"
        "</section>\n\n"
        "<!-- MOBILE -->\n"
        "<section id=\"deck-mobile\">\n"
        "  <div class=\"panel\">\n"
        f"{chr(10).join(mobile_lines)}\n"
        "  </div>\n"
        "</section>"
    )


# ────────────────────────────────────────────────────────────────────────────
# Titled-grid sub-variant (slides 19, 20)
#
# Photos arranged in a grid (4×2 for slide 19, 1×5 for slide 20) with a
# title text above the grid — no circle badge, no logo (defensively
# supported, but slides 19/20 have none). The title is a top-level
# <p:sp> TextBox positioned by its own OOXML coords (NOT a circle-overlay
# label like the badge variant). Mobile: title appears above the
# photo-grid-mobile div, photos stack or grid-2col per the shared
# mobile_photo_layout hint.
#
# Dispatched explicitly via the manifest hint
# `media_showcase_layout: "titled-grid"` (see render_media_showcase) —
# the OOXML topology (text shape with no circle) overlaps the no-text
# photo-grid case enough that explicit hint dispatch is safer than
# auto-detection.

def _classify_titled_grid(shapes, slide_w, slide_h):
    """Return (photos[], logo_pic, title_sp).

    Mirror of _classify but skips circle classification — titled-grid
    has no circle. Photos sorted into reading order (y, then x).
    First top-right small <p:pic> matching is_logo_pic is the logo;
    remaining pics are content. Slides 19/20 have no logo, so the
    classifier returns (photos, None, title_sp), but a future
    titled-grid slide could include one.
    """
    photos = []
    logo = None
    title_sp = None
    for s in shapes:
        if s.x_pt is None:
            continue
        if s.kind == "pic":
            if logo is None and is_logo_pic(s, slide_w, slide_h):
                logo = s
            else:
                photos.append(s)
        elif s.kind == "sp" and _has_text(s.element):
            if title_sp is None:
                title_sp = s
    photos.sort(key=lambda p: (round(p.y_pt), round(p.x_pt)))
    return photos, logo, title_sp


def _render_titled_grid_variant(
    slide, theme, flat, slide_index, slide_class,
    deck_name, slide_w_pt, slide_h_pt, deck_brand_color, media_dir,
) -> tuple:
    """Render the titled-grid sub-variant.

    Returns (html, []). Photos inline as optimized WebP via
    inline_optimized_data_url (transport rule for media_showcase
    photo_grid showcase — same as photo-grid + large-headline variants).
    Logo (if present) inlines as raw PPTX blob via inline_data_url.
    Title positioned by its own OOXML x/y/w/h via the .title class.
    """
    photos, logo, title_sp = _classify_titled_grid(flat, slide_w_pt, slide_h_pt)

    photo_urls = [inline_optimized_data_url(p, slide_class, slide, media_dir) for p in photos]
    logo_url = inline_data_url(logo, slide) if logo else None
    logo_invert = bool(slide_class.hints.get("logo_invert", False))

    headline_class = slide_class.hints.get("headline_class", "t-bold")
    inherited_bold_pt = INHERITED_BOLD_PT_BY_HINT.get(headline_class, 24.0)
    title_runs = _extract_overlay_runs(title_sp, theme, inherited_bold_pt) if title_sp else []

    mobile_photo_layout = slide_class.hints.get("mobile_photo_layout") or (
        "grid-2col" if len(photos) >= 5 else "stack"
    )

    canvas_bg = slide_class.hints.get("canvas_bg") or deck_brand_color or "#000000"

    slide_css = _build_css_titled(
        slide_w_pt, slide_h_pt,
        photos, logo, title_sp,
        title_runs, logo_invert,
        mobile_photo_layout,
    )
    body_html = _build_body_titled(photo_urls, logo_url, title_sp, title_runs)

    html = render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg":        canvas_bg,
            "bg-cyan":   deck_brand_color or "#000000",
            "headline":  "#FFFFFF",
            "font-cond": '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=body_html,
        slide_css=slide_css,
    )
    return (html, [])


def _build_css_titled(slide_w_pt, slide_h_pt, photos, logo, title_sp,
                      title_runs,
                      logo_invert: bool = False,
                      mobile_photo_layout: str = "stack") -> str:
    canvas_w_css, canvas_h_css = canvas_aspect_css(slide_w_pt, slide_h_pt)
    parts = []

    parts.append("/* ---------------- DESKTOP ---------------- */")
    parts.append("#deck-desktop { display: block; }")
    parts.append(
        f"#deck-desktop .canvas {{\n"
        f"  position: relative;\n"
        f"  width: {canvas_w_css};\n"
        f"  height: {canvas_h_css};\n"
        f"  margin: 0 auto;\n"
        f"  container-type: inline-size;\n"
        f"  overflow: hidden;\n"
        f"  background: var(--bg);\n"
        f"}}"
    )
    parts.append("#deck-desktop .slide { position: absolute; inset: 0; }")
    parts.append(
        "#deck-desktop .photo { position: absolute; overflow: hidden; }\n"
        "#deck-desktop .photo img { width: 100%; height: 100%; object-fit: cover; }"
    )
    for i, p in enumerate(photos):
        parts.append(
            f"#deck-desktop .photo-{i} {{ "
            f"left: {pt_to_pct_x(p.x_pt, slide_w_pt)}; "
            f"top: {pt_to_pct_y(p.y_pt, slide_h_pt)}; "
            f"width: {pt_to_pct_x(p.w_pt, slide_w_pt)}; "
            f"height: {pt_to_pct_y(p.h_pt, slide_h_pt)}; "
            f"}}"
        )
    if logo:
        logo_filter = "  filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            f"#deck-desktop .gif-logo {{\n"
            f"  {positioned_style(logo.x_pt, logo.y_pt, logo.w_pt, logo.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"{logo_filter}"
            f"}}\n"
            f"#deck-desktop .gif-logo img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )
    if title_sp:
        parts.append(
            f"#deck-desktop .title {{\n"
            f"  {positioned_style(title_sp.x_pt, title_sp.y_pt, title_sp.w_pt, title_sp.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  display: flex; flex-direction: column;\n"
            f"  align-items: center; justify-content: center;\n"
            f"  text-align: center;\n"
            f"}}"
        )

    bold_run = next((r for r in title_runs if r["klass"] == "t-bold"), None)
    sub_run = next((r for r in title_runs if r["klass"] == "t-sub"), None)
    if bold_run:
        parts.append(
            f"#deck-desktop .title .t-bold {{\n"
            f"  font-size: {pt_to_cqw(bold_run['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {bold_run['weight']};\n"
            f"  line-height: 1.0;\n"
            f"  color: {bold_run['color']};\n"
            f"  letter-spacing: 0.005em;\n"
            f"}}"
        )
    if sub_run:
        parts.append(
            f"#deck-desktop .title .t-sub {{\n"
            f"  font-size: {pt_to_cqw(sub_run['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {sub_run['weight']};\n"
            f"  line-height: 1.1;\n"
            f"  color: {sub_run['color']};\n"
            f"  margin-top: 0.4cqw;\n"
            f"}}"
        )

    # Mobile: title above photo grid, no top-bar/circle-row unless logo is present
    parts.append("\n/* ---------------- MOBILE ---------------- */")
    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    parts.append("  #deck-desktop { display: none; }")
    parts.append("  #deck-mobile { display: block; background: #fff; }")
    parts.append(
        "  #deck-mobile .panel {\n"
        "    min-height: 100vh;\n"
        "    background: #fff;\n"
        "    color: #111;\n"
        "    padding: 0;\n"
        "  }"
    )
    if logo:
        # Logo, if present, gets a mobile top-bar identical to the photo-grid
        # variant — cyan background with brightness(0)+invert(1) filter so
        # the logo reads white on cyan. See NOTES.md "logo_invert manifest
        # hint" entry for context.
        parts.append(
            "  #deck-mobile .top-bar {\n"
            "    display: flex; justify-content: flex-end;\n"
            "    padding: 5vw;\n"
            "    background: var(--bg-cyan);\n"
            "  }\n"
            "  #deck-mobile .gif-logo-mobile {\n"
            "    width: 22vw;\n"
            "    filter: brightness(0) invert(1);\n"
            "  }\n"
            "  #deck-mobile .gif-logo-mobile img { width: 100%; }"
        )
    if title_sp and title_runs:
        parts.append(
            "  #deck-mobile .title-mobile {\n"
            "    padding: 6vw 4vw 4vw;\n"
            "    text-align: center;\n"
            "    background: #fff;\n"
            "  }"
        )
        if bold_run:
            parts.append(
                f"  #deck-mobile .title-mobile .t-bold {{\n"
                f"    font-size: 4.4vw;\n"
                f"    font-weight: {bold_run['weight']};\n"
                f"    line-height: 1.2;\n"
                f"    color: {bold_run['color']};\n"
                f"  }}"
            )
        if sub_run:
            parts.append(
                f"  #deck-mobile .title-mobile .t-sub {{\n"
                f"    font-size: 3.6vw;\n"
                f"    font-weight: {sub_run['weight']};\n"
                f"    line-height: 1.3;\n"
                f"    color: {sub_run['color']};\n"
                f"    margin-top: 1.5vw;\n"
                f"  }}"
            )
    if mobile_photo_layout == "grid-2col":
        parts.append(
            "  #deck-mobile .photo-grid-mobile {\n"
            "    display: grid;\n"
            "    grid-template-columns: 1fr 1fr;\n"
            "    gap: 1vw; padding: 1vw; background: #fff;\n"
            "  }\n"
            "  #deck-mobile .photo-grid-mobile .photo-m {\n"
            "    width: 100%; background: #fff;\n"
            "    aspect-ratio: 1/1; overflow: hidden;\n"
            "  }\n"
            "  #deck-mobile .photo-grid-mobile .photo-m img {\n"
            "    width: 100%; height: 100%;\n"
            "    object-fit: contain; display: block;\n"
            "  }"
        )
    else:
        parts.append(
            "  #deck-mobile .photo-grid-mobile { display: flex; flex-direction: column; }\n"
            "  #deck-mobile .photo-grid-mobile .photo-m { width: 100%; }\n"
            "  #deck-mobile .photo-grid-mobile .photo-m img { width: 100%; height: auto; display: block; }"
        )
    parts.append("}")
    return "\n".join(parts)


def _build_body_titled(photo_urls, logo_url, title_sp, title_runs) -> str:
    desktop_lines = []
    for i, url in enumerate(photo_urls):
        if url:
            desktop_lines.append(f'      <div class="photo photo-{i}"><img src="{url}" alt=""></div>')
    if logo_url:
        desktop_lines.append(f'      <div class="gif-logo"><img src="{logo_url}" alt=""></div>')
    if title_sp and title_runs:
        desktop_lines.append('      <div class="title">')
        for r in title_runs:
            desktop_lines.append(f'        <div class="t {r["klass"]}">{_escape(r["text"])}</div>')
        desktop_lines.append('      </div>')

    mobile_lines = []
    if logo_url:
        mobile_lines.append('    <div class="top-bar">')
        mobile_lines.append(f'      <div class="gif-logo-mobile"><img src="{logo_url}" alt=""></div>')
        mobile_lines.append('    </div>')
    if title_sp and title_runs:
        mobile_lines.append('    <div class="title-mobile">')
        for r in title_runs:
            mobile_lines.append(f'      <div class="t {r["klass"]}">{_escape(r["text"])}</div>')
        mobile_lines.append('    </div>')
    if photo_urls:
        mobile_lines.append('    <div class="photo-grid-mobile">')
        for url in photo_urls:
            if url:
                mobile_lines.append(f'      <div class="photo-m"><img src="{url}" alt=""></div>')
        mobile_lines.append('    </div>')

    return (
        "<!-- DESKTOP -->\n"
        "<section id=\"deck-desktop\">\n"
        "  <div class=\"canvas\">\n"
        "    <div class=\"slide\">\n"
        f"{chr(10).join(desktop_lines)}\n"
        "    </div>\n"
        "  </div>\n"
        "</section>\n\n"
        "<!-- MOBILE -->\n"
        "<section id=\"deck-mobile\">\n"
        "  <div class=\"panel\">\n"
        f"{chr(10).join(mobile_lines)}\n"
        "  </div>\n"
        "</section>"
    )


# ────────────────────────────────────────────────────────────────────────────
# Helpers

def _data_url(shape: FlatShape, slide: Slide) -> Optional[str]:
    img_ref = extract_image_ref(shape.element, slide)
    if img_ref is None:
        return None
    b64 = base64.b64encode(img_ref.blob).decode("ascii")
    return f"data:{img_ref.content_type};base64,{b64}"


def _solidfill_hex(sp_elem) -> str:
    srgb = sp_elem.find("p:spPr/a:solidFill/a:srgbClr", NS)
    return srgb.get("val", "000000").upper() if srgb is not None else "000000"


def _escape(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


# ────────────────────────────────────────────────────────────────────────────
# Video variant (slide 7)
#
# Mobile rule: keep the bundle's container aspect (4:5) but use
# object-fit: contain instead of cover. The bundle's 4:5 cover crops ~27%
# off each side of a 16:9 video, killing edge-content like the Wicked 2
# wordmarks. Switching to contain alone (with the bundle's container)
# letterboxes the 16:9 video inside the 4:5 box — full-width, dominant,
# no edge cropping, ~14% black bars top + bottom. Initially we tried
# changing the container aspect to 16:9 too, but that shrank the box to
# ~45% the bundle's height; visually the video looked tiny.
#
# Centered-subject videos (e.g. slide 16) may want to opt back into cover;
# add `mobile_video_aspect` / `mobile_video_fit` hints to the manifest at
# that time. See NOTES.md entry on video edge-content.

def _render_video_variant(
    slide, theme, flat, video_pic, slide_index, slide_class,
    deck_name, slide_w_pt, slide_h_pt, deck_brand_color,
) -> tuple:
    """Full-bleed video + circle/text badge overlay + small logo top-right.

    Returns (html, []). The .mp4 itself is owned by the transform stage
    (see ondeck/transform/video.py) — render references it by the
    filename the transform wrote into slide_class.media.video.filename.
    The poster image stays inlined as a data URL on the <video> element.
    """
    photos, logo, circle, text_sp = _classify(flat, slide_w_pt, slide_h_pt)
    photos = [p for p in photos if p is not video_pic]  # video isn't a photo

    vref = extract_video(video_pic.element, slide)
    poster_data_url = None
    if vref and vref.poster_blob and vref.poster_content_type:
        b64 = base64.b64encode(vref.poster_blob).decode("ascii")
        poster_data_url = f"data:{vref.poster_content_type};base64,{b64}"

    video_filename = slide_class.media.get("video", {}).get("filename")
    if not video_filename:
        raise RuntimeError(
            f"slide {slide_index}: manifest is missing media.video.filename. "
            f"Run `python -m ondeck.transform.video <pptx> <manifest> <out_dir>` "
            f"before rendering — the transform stage writes the .mp4 and "
            f"records its filename in the manifest."
        )
    aux_files: list = []  # transform owns the .mp4; render only references it

    logo_url = inline_data_url(logo, slide) if logo else None
    circle_color = _solidfill_hex(circle.element) if circle else deck_brand_color
    logo_invert = bool(slide_class.hints.get("logo_invert", False))

    headline_class = slide_class.hints.get("headline_class", "t-bold")
    inherited_bold_pt = INHERITED_BOLD_PT_BY_HINT.get(headline_class, 24.0)
    text_runs = _extract_overlay_runs(text_sp, theme, inherited_bold_pt) if text_sp else []

    # Bare video: no chrome (no logo, no badge/text). Slide 16 is the
    # canonical case — bundle renders it as a flex-centered black panel
    # with the video alone, no top-bar or circle-row. Slide 7 has logo +
    # circle + text → not bare → existing locked path. See NOTES.md
    # entry on slide-16 audit (and Phase-2 backlog for slide 7).
    bare = (logo is None and not text_runs)

    canvas_bg = slide_class.hints.get("canvas_bg") or deck_brand_color or "#000000"

    slide_css = _build_video_css(
        slide_w_pt, slide_h_pt,
        video_pic, logo, circle,
        text_runs, circle_color, logo_invert,
        bare,
    )
    body_html = _build_video_body(
        video_filename, poster_data_url, logo_url, circle, text_runs,
        bare,
    )

    html = render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg":        canvas_bg,
            "bg-cyan":   deck_brand_color or "#000000",
            "headline":  "#FFFFFF",
            "font-cond": '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=body_html,
        slide_css=slide_css,
    )
    return (html, aux_files)


def _build_video_css(slide_w_pt, slide_h_pt, video_pic, logo, circle, text_runs, circle_color,
                     logo_invert: bool = False, bare: bool = False) -> str:
    canvas_w_css, canvas_h_css = canvas_aspect_css(slide_w_pt, slide_h_pt)
    parts = []

    parts.append("/* ---------------- DESKTOP ---------------- */")
    parts.append("#deck-desktop { display: block; }")
    parts.append(
        f"#deck-desktop .canvas {{\n"
        f"  position: relative;\n"
        f"  width: {canvas_w_css};\n"
        f"  height: {canvas_h_css};\n"
        f"  margin: 0 auto;\n"
        f"  container-type: inline-size;\n"
        f"  overflow: hidden;\n"
        f"  background: var(--bg);\n"
        f"}}"
    )
    parts.append("#deck-desktop .slide { position: absolute; inset: 0; }")
    parts.append(
        f"#deck-desktop .video-wrap {{ {positioned_style(video_pic.x_pt, video_pic.y_pt, video_pic.w_pt, video_pic.h_pt, slide_w_pt, slide_h_pt)} }}\n"
        f"#deck-desktop .video-wrap video {{ width: 100%; height: 100%; object-fit: contain; background: #fff; }}"
    )
    if logo:
        logo_filter = "  filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            f"#deck-desktop .gif-logo {{\n"
            f"  {positioned_style(logo.x_pt, logo.y_pt, logo.w_pt, logo.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  pointer-events: none;\n"
            f"{logo_filter}"
            f"}}\n"
            f"#deck-desktop .gif-logo img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )
    if circle:
        parts.append(
            f"#deck-desktop .circle {{\n"
            f"  {positioned_style(circle.x_pt, circle.y_pt, circle.w_pt, circle.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  background: #{circle_color};\n"
            f"  border-radius: 50%;\n"
            f"  pointer-events: none;\n"
            f"}}"
        )
        parts.append(
            f"#deck-desktop .L {{\n"
            f"  {positioned_style(circle.x_pt, circle.y_pt, circle.w_pt, circle.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  display: flex; flex-direction: column;\n"
            f"  align-items: center; justify-content: center;\n"
            f"  text-align: center;\n"
            f"  padding: 1.5cqw;\n"
            f"  pointer-events: none;\n"
            f"}}"
        )

    bold_run = next((r for r in text_runs if r["klass"] == "t-bold"), None)
    sub_run = next((r for r in text_runs if r["klass"] == "t-sub"), None)
    if bold_run:
        parts.append(
            f"#deck-desktop .L .t-bold {{\n"
            f"  font-size: {pt_to_cqw(bold_run['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {bold_run['weight']};\n"
            f"  line-height: 1.0;\n"
            f"  color: {bold_run['color']};\n"
            f"  letter-spacing: 0.005em;\n"
            f"}}"
        )
    if sub_run:
        parts.append(
            f"#deck-desktop .L .t-sub {{\n"
            f"  font-size: {pt_to_cqw(sub_run['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {sub_run['weight']};\n"
            f"  line-height: 1.1;\n"
            f"  color: {sub_run['color']};\n"
            f"  margin-top: 0.4cqw;\n"
            f"}}"
        )

    # Mobile — two branches:
    #   bare=True  → flex-centered black panel with just the video (slide 16).
    #               No top-bar, no circle-row. Bundle pattern: short and clean.
    #   bare=False → 4-row stacked panel: top-bar (cyan + logo) + circle-row
    #               (cyan circle with text) + video-mobile (16:9 cover). Slide 7
    #               locked path; preserve byte output exactly.
    parts.append("\n/* ---------------- MOBILE ---------------- */")
    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    if bare:
        # Replicate section_divider's locked mobile centering recipe verbatim:
        # full ancestor-chain heights, panel with position:relative + overflow:
        # hidden, content centered via absolute+translateY (NOT flex). The flex
        # variant of this layout was broken on iPhone Safari — flex-formatting
        # context interacts poorly with 100dvh resolution. Absolute centering
        # is the proven pattern across 5 section_divider slides and 5 rounds
        # of iPhone iteration. See NOTES.md "iOS Safari `100dvh` requires
        # explicit height on every ancestor" + the section_divider mobile
        # architecture entry.
        parts.append("  html, body { height: 100%; }")
        parts.append("  #deck-desktop { display: none; }")
        parts.append("  #deck-mobile { display: block; height: 100%; }")
        parts.append(
            "  #deck-mobile .panel {\n"
            "    height: 100svh;\n"
            "    height: 100dvh;\n"
            "    background: #000;\n"
            "    position: relative;\n"
            "    overflow: hidden;\n"
            "  }\n"
            "  #deck-mobile .video-mobile {\n"
            "    position: absolute;\n"
            "    top: 50%;\n"
            "    left: 0; right: 0;\n"
            "    transform: translateY(-50%);\n"
            "    width: 100%;\n"
            "    background: #000;\n"
            "  }\n"
            "  #deck-mobile .video-mobile video {\n"
            "    width: 100%;\n"
            "    aspect-ratio: 16/9;\n"
            "    background: #000;\n"
            "    display: block;\n"
            "    object-fit: contain;\n"
            "  }"
        )
    else:
        parts.append("  html, body { height: 100%; }")
        parts.append("  #deck-desktop { display: none; }")
        parts.append("  #deck-mobile { display: block; height: 100%; background: #fff; }")
        parts.append(
            "  #deck-mobile .panel {\n"
            "    height: 100svh;\n"
            "    height: 100dvh;\n"
            "    background: #fff;\n"
            "    color: #111;\n"
            "    padding: 0;\n"
            "    position: relative;\n"
            "    overflow: hidden;\n"
            "  }"
        )
        # Mobile top-bar is always cyan → mobile logo always inverted when
        # present (see _build_css comment for full context).
        logo_filter_m = "    filter: brightness(0) invert(1);\n" if logo is not None else ""
        parts.append(
            "  #deck-mobile .top-bar {\n"
            "    display: flex; justify-content: flex-end;\n"
            "    padding: 5vw;\n"
            "    background: var(--bg-cyan);\n"
            "  }\n"
            "  #deck-mobile .gif-logo-mobile {\n"
            "    width: 22vw;\n"
            f"{logo_filter_m}"
            "  }\n"
            "  #deck-mobile .gif-logo-mobile img { width: 100%; }"
        )
        parts.append(
            "  #deck-mobile .circle-row {\n"
            "    display: flex; justify-content: center;\n"
            "    padding: 6vw 0 4vw;\n"
            "    background: #fff;\n"
            "  }\n"
            f"  #deck-mobile .circle-mobile {{\n"
            f"    background: #{circle_color};\n"
            f"    border-radius: 50%;\n"
            f"    width: 38vw; height: 38vw;\n"
            f"    display: flex; flex-direction: column;\n"
            f"    align-items: center; justify-content: center;\n"
            f"    color: var(--headline);\n"
            f"  }}\n"
            "  #deck-mobile .circle-mobile .t-bold { font-size: 6vw; font-weight: 700; line-height: 1.0; }\n"
            "  #deck-mobile .circle-mobile .t-sub { font-size: 3.4vw; font-weight: 400; margin-top: 1vw; }"
        )
        parts.append(
            "  #deck-mobile .video-mobile {\n"
            "    padding: 0;\n"
            "    background: #000;\n"
            "  }\n"
            "  #deck-mobile .video-mobile video {\n"
            "    width: 100%;\n"
            "    aspect-ratio: 16/9;\n"
            "    background: #000;\n"
            "    display: block;\n"
            "    object-fit: cover;\n"
            "  }"
        )
    parts.append("}")
    return "\n".join(parts)


def _build_video_body(video_filename, poster_data_url, logo_url, circle, text_runs,
                      bare: bool = False) -> str:
    poster_attr = f' poster="{poster_data_url}"' if poster_data_url else ""
    video_el = (
        f'<video autoplay loop muted playsinline controls preload="auto"{poster_attr}>\n'
        f'          <source src="{video_filename}" type="video/mp4">\n'
        f'          Your browser does not support video playback.\n'
        f'        </video>'
    )

    desktop_lines = []
    desktop_lines.append(f'      <div class="video-wrap">\n        {video_el}\n      </div>')
    if circle:
        desktop_lines.append('      <div class="circle"></div>')
    if text_runs:
        desktop_lines.append('      <div class="L">')
        for r in text_runs:
            desktop_lines.append(f'        <div class="t {r["klass"]}">{_escape(r["text"])}</div>')
        desktop_lines.append('      </div>')
    if logo_url:
        desktop_lines.append(f'      <div class="gif-logo"><img src="{logo_url}" alt=""></div>')

    mobile_video_el = (
        f'<video autoplay loop muted playsinline controls preload="auto"{poster_attr}>\n'
        f'        <source src="{video_filename}" type="video/mp4">\n'
        f'      </video>'
    )

    mobile_lines = []
    if bare:
        # Slide 16 pattern: just the video centered in a black panel, no chrome.
        mobile_lines.append(f'    <div class="video-mobile">\n      {mobile_video_el}\n    </div>')
    else:
        # Slide 7 pattern: top-bar (logo) + optional circle-row + video-mobile.
        mobile_lines.append('    <div class="top-bar">')
        if logo_url:
            mobile_lines.append(f'      <div class="gif-logo-mobile"><img src="{logo_url}" alt=""></div>')
        mobile_lines.append('    </div>')
        if text_runs:
            mobile_lines.append('    <div class="circle-row">')
            mobile_lines.append('      <div class="circle-mobile L L-mobile">')
            for r in text_runs:
                mobile_lines.append(f'        <div class="t {r["klass"]}">{_escape(r["text"])}</div>')
            mobile_lines.append('      </div>')
            mobile_lines.append('    </div>')
        mobile_lines.append(f'    <div class="video-mobile">\n      {mobile_video_el}\n    </div>')

    return (
        "<!-- DESKTOP -->\n"
        "<section id=\"deck-desktop\">\n"
        "  <div class=\"canvas\">\n"
        "    <div class=\"slide\">\n"
        f"{chr(10).join(desktop_lines)}\n"
        "    </div>\n"
        "  </div>\n"
        "</section>\n\n"
        "<!-- MOBILE -->\n"
        "<section id=\"deck-mobile\">\n"
        "  <div class=\"panel\">\n"
        f"{chr(10).join(mobile_lines)}\n"
        "  </div>\n"
        "</section>"
    )


# ────────────────────────────────────────────────────────────────────────────
# Large-headline overlay variant (slides 15, 17)
#
# Structurally distinct from photo-grid badge-overlay:
#   - N photos at OOXML positions (no logo, no circle/badge)
#   - 1 gradient sp covering left portion of canvas (decorative tint)
#   - 1 headline sp at bottom-left, 88pt stacked ALL CAPS, 3-stop drop shadow
#
# Mobile pivots to a full-bleed CSS Grid of all photos with a bottom-fade
# cyan overlay and a bottom-anchored 88pt headline. Different mobile
# composition from badge-overlay (top-bar + circle-row + photo-grid).
#
# Detection: presence of any sp with <a:gradFill> on the slide. Dispatched
# from render_media_showcase before falling through to photo-grid variant.

def _classify_large_headline(shapes, slide_w, slide_h):
    """Return (photos[], gradient_sp, headline_sp, logo_pic).

    Same logo classifier as the badge-overlay variant (first-match-only).
    Photos sorted by (y, x) — left-to-right, top-to-bottom reading order
    drives the desktop .photo-N indexing. Gradient sp is the first sp
    with <a:gradFill>; headline sp is the first sp with text.
    """
    photos = []
    logo = None
    gradient_sp = None
    headline_sp = None
    for s in shapes:
        if s.x_pt is None:
            continue
        if s.kind == "pic":
            if logo is None and is_logo_pic(s, slide_w, slide_h):
                logo = s
            else:
                photos.append(s)
        elif s.kind == "sp":
            if gradient_sp is None and s.element.find("p:spPr/a:gradFill", NS) is not None:
                gradient_sp = s
            elif headline_sp is None and _has_text(s.element):
                headline_sp = s
    photos.sort(key=lambda p: (round(p.y_pt), round(p.x_pt)))
    return photos, gradient_sp, headline_sp, logo


def _render_large_headline_variant(
    slide, theme, flat, slide_index, slide_class,
    deck_name, slide_w_pt, slide_h_pt, deck_brand_color, media_dir,
) -> tuple:
    """Render the large-headline overlay sub-variant.

    Returns (html, []). Photos inline as optimized WebP via
    inline_optimized_data_url (transport rule for media_showcase
    photo_grid showcase — see NOTES.md Operating principles entry
    on the photo_grid carve-out). Logo inlines as raw PPTX blob.
    Gradient color resolves via parse_gradient (handles lumMod/lumOff
    on theme accent1).
    """
    photos, gradient_sp, headline_sp, logo = _classify_large_headline(
        flat, slide_w_pt, slide_h_pt,
    )

    photo_urls = [inline_optimized_data_url(p, slide_class, slide, media_dir) for p in photos]
    logo_url = inline_data_url(logo, slide) if logo else None
    logo_invert = bool(slide_class.hints.get("logo_invert", False))

    overlay_info = parse_gradient(gradient_sp, theme) if gradient_sp else None
    headline_runs = extract_headline_runs(headline_sp, theme, slide_w_pt) if headline_sp else []

    canvas_bg = slide_class.hints.get("canvas_bg") or deck_brand_color or "#000000"

    slide_css = _build_large_headline_css(
        slide_w_pt, slide_h_pt,
        photos, gradient_sp, headline_sp, logo,
        overlay_info, headline_runs, logo_invert,
    )
    body_html = _build_large_headline_body(
        photo_urls, logo_url, gradient_sp, headline_runs,
    )

    overlay_color_css = overlay_info["color_css"] if overlay_info else "rgba(0,0,0,0)"
    overlay_zero_css = overlay_info["color_zero_css"] if overlay_info else "rgba(0,0,0,0)"

    html = render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg":         canvas_bg,
            "bg-cyan":    deck_brand_color or "#000000",
            "headline":   "#FFFFFF",
            "overlay":    overlay_color_css,
            "overlay-0":  overlay_zero_css,
            "font-cond":  '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=body_html,
        slide_css=slide_css,
    )
    return (html, [])


def _build_large_headline_css(slide_w_pt, slide_h_pt,
                              photos, gradient_sp, headline_sp, logo,
                              overlay_info, headline_runs,
                              logo_invert: bool = False) -> str:
    canvas_w_css, canvas_h_css = canvas_aspect_css(slide_w_pt, slide_h_pt)
    parts = []

    parts.append("/* ---------------- DESKTOP ---------------- */")
    parts.append("#deck-desktop { display: block; }")
    parts.append(
        f"#deck-desktop .canvas {{\n"
        f"  position: relative;\n"
        f"  width: {canvas_w_css};\n"
        f"  height: {canvas_h_css};\n"
        f"  margin: 0 auto;\n"
        f"  container-type: inline-size;\n"
        f"  overflow: hidden;\n"
        f"  background: var(--bg);\n"
        f"}}"
    )
    parts.append("#deck-desktop .slide { position: absolute; inset: 0; }")
    parts.append(
        "#deck-desktop .photo { position: absolute; overflow: hidden; }\n"
        "#deck-desktop .photo img { width: 100%; height: 100%; object-fit: cover; display: block; }"
    )
    for i, p in enumerate(photos):
        parts.append(
            f"#deck-desktop .photo-{i} {{ "
            f"left: {pt_to_pct_x(p.x_pt, slide_w_pt)}; "
            f"top: {pt_to_pct_y(p.y_pt, slide_h_pt)}; "
            f"width: {pt_to_pct_x(p.w_pt, slide_w_pt)}; "
            f"height: {pt_to_pct_y(p.h_pt, slide_h_pt)}; "
            f"}}"
        )

    if gradient_sp and overlay_info:
        gradient_css = (
            f"linear-gradient({overlay_info['css_angle_deg']:.0f}deg, "
            + ", ".join(overlay_info["stop_css"])
            + ")"
        )
        parts.append(
            f"#deck-desktop .overlay {{\n"
            f"  {positioned_style(gradient_sp.x_pt, gradient_sp.y_pt, gradient_sp.w_pt, gradient_sp.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"  background: {gradient_css};\n"
            f"}}"
        )

    if logo:
        logo_filter = "  filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            f"#deck-desktop .gif-logo {{\n"
            f"  {positioned_style(logo.x_pt, logo.y_pt, logo.w_pt, logo.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"{logo_filter}"
            f"}}\n"
            f"#deck-desktop .gif-logo img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )

    if headline_sp and headline_runs:
        first = headline_runs[0]
        parts.append(
            f"#deck-desktop .L {{\n"
            f"  position: absolute;\n"
            f"  left: {pt_to_pct_x(headline_sp.x_pt, slide_w_pt)};\n"
            f"  top: {pt_to_pct_y(headline_sp.y_pt, slide_h_pt)};\n"
            f"  width: {pt_to_pct_x(headline_sp.w_pt, slide_w_pt)};\n"
            f"  height: auto;\n"
            f"  display: flex;\n"
            f"  flex-direction: column;\n"
            f"}}"
        )
        shadow_decl = f"\n  text-shadow: {first['shadow_cqw']};" if first.get("shadow_cqw") else ""
        parts.append(
            f"#deck-desktop .L .t {{\n"
            f"  font-size: {pt_to_cqw(first['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {first['weight']};\n"
            f"  letter-spacing: 0.005em;\n"
            f"  line-height: 1.0;\n"
            f"  color: {first['color']};\n"
            f"  text-transform: uppercase;{shadow_decl}\n"
            f"}}"
        )

    # Mobile — full-bleed CSS Grid of photos + cyan bottom-fade overlay +
    # bottom-anchored 88pt headline. Mirrors section_divider mobile pattern
    # but with multiple photos in a grid instead of single full-bleed.
    n = len(photos)
    parts.append("\n/* ---------------- MOBILE ---------------- */")
    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    parts.append("  html, body { height: 100%; }")
    parts.append("  #deck-desktop { display: none; }")
    parts.append("  #deck-mobile { display: block; height: 100%; }")
    parts.append(
        "  #deck-mobile .panel {\n"
        "    height: 100svh;\n"
        "    height: 100dvh;\n"
        "    background: #000;\n"
        "    color: var(--headline);\n"
        "    position: relative;\n"
        "    overflow: hidden;\n"
        "  }"
    )
    # Grid layout — hardcoded by photo count for B2 (slides 15, 17). Future
    # sub-patterns (e.g. 6-up, 9-up) opt in via a manifest hint when they
    # surface; don't add the hint speculatively.
    if n == 4:
        # 2×2 — slide 15 (2 portraits left + 2 landscapes stacked right).
        parts.append(
            "  #deck-mobile .photo-grid {\n"
            "    position: absolute; inset: 0;\n"
            "    display: grid;\n"
            "    grid-template-columns: 1fr 1fr;\n"
            "    grid-template-rows: 1fr 1fr;\n"
            "    gap: 0;\n"
            "    z-index: 0;\n"
            "  }"
        )
    elif n == 5:
        # tall-left + 2×2 right — slide 17 (1 tall portrait spanning column 1,
        # 4 landscape cells in column 2).
        parts.append(
            "  #deck-mobile .photo-grid {\n"
            "    position: absolute; inset: 0;\n"
            "    display: grid;\n"
            "    grid-template-columns: 1fr 1fr;\n"
            "    grid-template-rows: 1fr 1fr 1fr 1fr;\n"
            "    gap: 0;\n"
            "    z-index: 0;\n"
            "  }\n"
            "  #deck-mobile .photo-grid .photo-m:first-child { grid-row: 1 / 5; grid-column: 1; }"
        )
    else:
        # Fallback — single column. No deck currently hits this; if a
        # future deck does and the result is wrong, add a manifest hint.
        parts.append(
            "  #deck-mobile .photo-grid {\n"
            "    position: absolute; inset: 0;\n"
            "    display: flex; flex-direction: column;\n"
            "    z-index: 0;\n"
            "  }"
        )
    parts.append(
        "  #deck-mobile .photo-grid .photo-m { overflow: hidden; }\n"
        "  #deck-mobile .photo-grid .photo-m img {\n"
        "    width: 100%; height: 100%;\n"
        "    object-fit: cover; display: block;\n"
        "  }"
    )
    parts.append(
        "  #deck-mobile .overlay-mobile {\n"
        "    position: absolute; inset: 0;\n"
        "    background: linear-gradient(to bottom,\n"
        f"      var(--overlay-0) 30%,\n"
        f"      {(_alpha_at(overlay_info, 0.85) if overlay_info else 'rgba(0,0,0,0)')} 70%,\n"
        f"      var(--bg-cyan) 100%);\n"
        "    z-index: 1;\n"
        "    pointer-events: none;\n"
        "  }"
    )
    if headline_runs:
        first = headline_runs[0]
        shadow_m = ""
        if first.get("shadow_cqw"):
            shadow_m = "\n    text-shadow: 0.6vw 0.6vw 1.5vw rgba(0,0,0,0.4);"
        parts.append(
            "  #deck-mobile .headline-mobile {\n"
            "    position: absolute;\n"
            "    left: 0; right: 0; bottom: 16dvh;\n"
            "    padding: 0 6vw calc(env(safe-area-inset-bottom));\n"
            "    z-index: 2;\n"
            "  }\n"
            "  #deck-mobile .headline-mobile .t {\n"
            "    font-size: clamp(3rem, 13vw, 4rem);\n"
            f"    font-weight: {first['weight']};\n"
            "    line-height: 0.95;\n"
            "    text-transform: uppercase;\n"
            "    letter-spacing: 0.005em;"
            f"{shadow_m}\n"
            "  }"
        )
    parts.append("}")
    return "\n".join(parts)


def _alpha_at(overlay_info: dict, alpha: float) -> str:
    """Build an rgba(...) string from an overlay_info dict at the given alpha.

    Used by the large-headline mobile bottom-fade overlay's middle stop —
    we need the gradient base color at 85% opacity, but parse_gradient's
    stop_css only carries the alphas declared in OOXML. Reconstruct from
    the resolved hex.
    """
    hex_color = overlay_info["color_css"].lstrip("#")
    r, g, b = (int(hex_color[i : i + 2], 16) for i in (0, 2, 4))
    return f"rgba({r},{g},{b},{alpha})"


def _build_large_headline_body(photo_urls, logo_url, gradient_sp, headline_runs) -> str:
    desktop_lines = []
    for i, url in enumerate(photo_urls):
        if url:
            desktop_lines.append(
                f'      <div class="photo photo-{i}"><img src="{url}" alt=""></div>'
            )
    if gradient_sp:
        desktop_lines.append('      <div class="overlay"></div>')
    if logo_url:
        desktop_lines.append(f'      <div class="gif-logo"><img src="{logo_url}" alt=""></div>')
    if headline_runs:
        desktop_lines.append('      <div class="L">')
        for r in headline_runs:
            desktop_lines.append(f'        <div class="t">{_escape(r["text"])}</div>')
        desktop_lines.append('      </div>')

    mobile_lines = []
    if photo_urls:
        mobile_lines.append('    <div class="photo-grid">')
        for url in photo_urls:
            if url:
                mobile_lines.append(f'      <div class="photo-m"><img src="{url}" alt=""></div>')
        mobile_lines.append('    </div>')
    if gradient_sp:
        mobile_lines.append('    <div class="overlay-mobile"></div>')
    if headline_runs:
        mobile_lines.append('    <div class="L L-mobile headline-mobile">')
        for r in headline_runs:
            mobile_lines.append(f'      <div class="t">{_escape(r["text"])}</div>')
        mobile_lines.append('    </div>')

    return (
        "<!-- DESKTOP -->\n"
        "<section id=\"deck-desktop\">\n"
        "  <div class=\"canvas\">\n"
        "    <div class=\"slide\">\n"
        f"{chr(10).join(desktop_lines)}\n"
        "    </div>\n"
        "  </div>\n"
        "</section>\n\n"
        "<!-- MOBILE -->\n"
        "<section id=\"deck-mobile\">\n"
        "  <div class=\"panel\">\n"
        f"{chr(10).join(mobile_lines)}\n"
        "  </div>\n"
        "</section>"
    )
