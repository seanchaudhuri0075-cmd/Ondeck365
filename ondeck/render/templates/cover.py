"""Cover template renderer — full HTML for a 'cover' slide.

Slide 1 of P&G is the only cover in this deck. Layout pattern:
  - full-bleed background fill (sp at z=0, inline solidFill)
  - hero image left-center (pic, often SVG)
  - 2-line text block right-center (sp headline + sp sub)
  - small logo top-right (pic, often PNG with transparency)

Role detection is geometry-based — bg = full-bleed sp; logo =
top-right small pic; hero = the remaining largest pic; text = the
other sp shapes. The two text shapes are merged into one `.L` block
positioned at the union of their bounds, matching the bundle's
collapsed structure.
"""
from __future__ import annotations

import base64
from typing import Optional

from pptx.slide import Slide

from ..css import pt_to_pct_x, pt_to_pct_y, pt_to_cqw, color_with_alpha
from ..desktop import canvas_aspect_css, positioned_style
from ..html import render_page
from ...layout.detect import SlideClass
from ...parse.font_calibration import calibrate_size_pt
from ...parse.images import extract_image_ref
from ...parse.shapes import FlatShape, flatten_slide
from ...parse.slide import NS
from ...parse.svg import extract_svg_ref
from ...parse.text import parse_text_frame
from ...parse.theme import Theme
from ._shared import image_src, inline_data_url, is_logo_pic


# Empirical font-weight per typeface, derived from the P&G bundle CSS.
# Univers Condensed → 500 (Medium), Light → 300, Bold → 700.
TYPEFACE_WEIGHT = {
    "univers condensed":        500,
    "univers condensed light":  300,
    "univers condensed bold":   700,
    "barlow condensed":         500,
    "barlow condensed light":   300,
    "barlow condensed bold":    700,
}


def render_cover(
    slide: Slide,
    theme: Theme,
    slide_index: int,
    slide_class: SlideClass,
    deck_name: str,
    slide_w_pt: float,
    slide_h_pt: float,
) -> str:
    """Render a cover slide to a complete HTML page."""
    flat = list(flatten_slide(slide))
    bg, hero, logo, text_shapes = _classify(flat, slide_w_pt, slide_h_pt)

    bg_color = _solidfill_hex(bg.element, theme) if bg else "000000"
    hero_url = image_src(hero, slide_class, slide) if hero else None
    logo_url = inline_data_url(logo, slide) if logo else None
    logo_invert = bool(slide_class.hints.get("logo_invert", False))

    # Extract text runs in document order. First run = headline, second = sub.
    runs = []
    for ts in text_shapes:
        tf = parse_text_frame(ts.element, theme)
        if tf is None:
            continue
        for p in tf.paragraphs:
            for r in p.runs:
                if not r.text.strip():
                    continue
                weight = TYPEFACE_WEIGHT.get((r.style.typeface or "").lower(), 500)
                if r.style.bold:
                    weight = 700
                runs.append({
                    "text": r.text,
                    "size_pt": calibrate_size_pt(
                        r.style.typeface, r.style.size_pt, r.style.bold
                    ),
                    "weight": weight,
                    "color": color_with_alpha(
                        r.style.color_hex or "FFFFFF", r.style.color_alpha
                    ),
                })

    headline = runs[0] if runs else None
    sub = runs[1] if len(runs) > 1 else None

    # Combined .L block position = union bounds of the text shapes.
    if text_shapes:
        L_x = min(s.x_pt for s in text_shapes)
        L_y = min(s.y_pt for s in text_shapes)
        L_w = max(s.x_pt + s.w_pt for s in text_shapes) - L_x
    else:
        L_x = L_y = L_w = 0

    slide_css = _build_css(
        slide_w_pt, slide_h_pt,
        hero, logo, text_shapes,
        L_x, L_y, L_w,
        headline, sub, logo_invert,
    )

    body_html = _build_body(hero_url, logo_url, headline, sub, text_shapes)

    return render_page(
        title=_escape(f"{deck_name} — Slide {slide_index}"),
        root_vars={
            "bg":         f"#{bg_color}",
            "headline":   "#FFFFFF",
            "sub":        "rgba(255,255,255,0.38)",
            "font-cond":  '"Barlow Condensed", "Univers Condensed", "Arial Narrow", sans-serif',
        },
        body_html=body_html,
        slide_css=slide_css,
    )


# ────────────────────────────────────────────────────────────────────────────
# CSS assembly

def _build_css(slide_w_pt, slide_h_pt, hero, logo, text_shapes,
               L_x, L_y, L_w, headline, sub, logo_invert: bool = False) -> str:
    canvas_w_css, canvas_h_css = canvas_aspect_css(slide_w_pt, slide_h_pt)
    parts = []

    parts.append("/* ---------------- DESKTOP ---------------- */")
    parts.append("#deck-desktop { display: block; }")
    parts.append(
        f"#deck-desktop .canvas {{\n"
        f"  position: relative;\n"
        f"  width: {canvas_w_css};\n"
        f"  height: {canvas_h_css};\n"
        f"  margin: 0 auto;\n"
        f"  container-type: inline-size;\n"
        f"  overflow: hidden;\n"
        f"  background: var(--bg);\n"
        f"}}"
    )
    parts.append("#deck-desktop .slide { position: absolute; inset: 0; }")

    if hero:
        parts.append(
            f"#deck-desktop .hero {{ {positioned_style(hero.x_pt, hero.y_pt, hero.w_pt, hero.h_pt, slide_w_pt, slide_h_pt)} }}\n"
            f"#deck-desktop .hero img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )

    if logo:
        logo_filter = "  filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            f"#deck-desktop .gif-logo {{\n"
            f"  {positioned_style(logo.x_pt, logo.y_pt, logo.w_pt, logo.h_pt, slide_w_pt, slide_h_pt)}\n"
            f"{logo_filter}"
            f"}}\n"
            f"#deck-desktop .gif-logo img {{ width: 100%; height: 100%; object-fit: contain; }}"
        )

    if text_shapes:
        parts.append(
            f"#deck-desktop .L {{\n"
            f"  position: absolute;\n"
            f"  left: {pt_to_pct_x(L_x, slide_w_pt)};\n"
            f"  top: {pt_to_pct_y(L_y, slide_h_pt)};\n"
            f"  width: {pt_to_pct_x(L_w, slide_w_pt)};\n"
            f"  height: auto;\n"
            f"  display: flex;\n"
            f"  flex-direction: column;\n"
            f"}}"
        )

    if headline:
        parts.append(
            f"#deck-desktop .L .t-headline {{\n"
            f"  font-size: {pt_to_cqw(headline['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {headline['weight']};\n"
            f"  letter-spacing: 0.01em;\n"
            f"  line-height: 1.1;\n"
            f"  color: {headline['color']};\n"
            f"  text-transform: uppercase;\n"
            f"}}"
        )

    if sub:
        parts.append(
            f"#deck-desktop .L .t-sub {{\n"
            f"  font-size: {pt_to_cqw(sub['size_pt'], slide_w_pt)};\n"
            f"  font-weight: {sub['weight']};\n"
            f"  letter-spacing: 0.01em;\n"
            f"  line-height: 1.05;\n"
            f"  color: {sub['color']};\n"
            f"  margin-top: 0.4cqw;\n"
            f"  max-width: 18cqw;\n"
            f"}}"
        )

    parts.append("\n/* ---------------- MOBILE ---------------- */")
    parts.append("#deck-mobile { display: none; }")
    parts.append("@media (max-width: 768px) {")
    parts.append("  #deck-desktop { display: none; }")
    parts.append("  #deck-mobile { display: block; }")
    parts.append(
        "  #deck-mobile .panel {\n"
        "    min-height: 100vh;\n"
        "    background: var(--bg);\n"
        "    padding: 6vw 6vw 8vw;\n"
        "    display: flex;\n"
        "    flex-direction: column;\n"
        "  }"
    )
    if logo:
        logo_filter_m = "    filter: brightness(0) invert(1);\n" if logo_invert else ""
        parts.append(
            "  #deck-mobile .gif-logo-mobile {\n"
            "    align-self: flex-end; width: 22vw; margin-bottom: 8vw;\n"
            f"{logo_filter_m}"
            "  }\n"
            "  #deck-mobile .gif-logo-mobile img { width: 100%; }"
        )
    if hero:
        parts.append(
            "  #deck-mobile .hero-mobile { width: 70vw; align-self: center; margin: 4vw 0 8vw; }\n"
            "  #deck-mobile .hero-mobile img { width: 100%; }"
        )
    if headline:
        parts.append(
            f"  #deck-mobile .L-mobile .t-headline {{\n"
            f"    font-size: 11vw;\n"
            f"    font-weight: {headline['weight']};\n"
            f"    color: {headline['color']};\n"
            f"    text-transform: uppercase;\n"
            f"    line-height: 1.1;\n"
            f"    letter-spacing: 0.01em;\n"
            f"  }}"
        )
    if sub:
        parts.append(
            f"  #deck-mobile .L-mobile .t-sub {{\n"
            f"    font-size: 16vw;\n"
            f"    font-weight: {sub['weight']};\n"
            f"    color: {sub['color']};\n"
            f"    line-height: 1.05;\n"
            f"    letter-spacing: 0.01em;\n"
            f"    margin-top: 1.5vw;\n"
            f"  }}"
        )
    parts.append("}")
    return "\n".join(parts)


# ────────────────────────────────────────────────────────────────────────────
# Body assembly

def _build_body(hero_url, logo_url, headline, sub, text_shapes) -> str:
    desktop_lines = []
    if hero_url:
        desktop_lines.append(f'      <div class="hero"><img src="{hero_url}" alt=""></div>')
    if logo_url:
        desktop_lines.append(f'      <div class="gif-logo"><img src="{logo_url}" alt=""></div>')
    if text_shapes and (headline or sub):
        desktop_lines.append('      <div class="L">')
        if headline:
            desktop_lines.append(f'        <div class="t t-headline">{_escape(headline["text"])}</div>')
        if sub:
            desktop_lines.append(f'        <div class="t t-sub">{_escape(sub["text"])}</div>')
        desktop_lines.append('      </div>')

    mobile_lines = []
    if logo_url:
        mobile_lines.append(f'    <div class="gif-logo-mobile"><img src="{logo_url}" alt=""></div>')
    if hero_url:
        mobile_lines.append(f'    <div class="hero-mobile"><img src="{hero_url}" alt=""></div>')
    if text_shapes and (headline or sub):
        mobile_lines.append('    <div class="L L-mobile">')
        if headline:
            mobile_lines.append(f'      <div class="t t-headline">{_escape(headline["text"])}</div>')
        if sub:
            mobile_lines.append(f'      <div class="t t-sub">{_escape(sub["text"])}</div>')
        mobile_lines.append('    </div>')

    desktop_html = "\n".join(desktop_lines)
    mobile_html = "\n".join(mobile_lines)

    return (
        "<!-- DESKTOP -->\n"
        "<section id=\"deck-desktop\">\n"
        "  <div class=\"canvas\">\n"
        "    <div class=\"slide\">\n"
        f"{desktop_html}\n"
        "    </div>\n"
        "  </div>\n"
        "</section>\n\n"
        "<!-- MOBILE -->\n"
        "<section id=\"deck-mobile\">\n"
        "  <div class=\"panel\">\n"
        f"{mobile_html}\n"
        "  </div>\n"
        "</section>"
    )


# ────────────────────────────────────────────────────────────────────────────
# Role detection + asset extraction

def _classify(shapes, slide_w, slide_h):
    """Return (bg_sp, hero_pic, logo_pic, text_shapes) by geometry."""
    bg = None
    hero = None
    logo = None
    text_shapes = []
    for s in shapes:
        if s.x_pt is None:
            continue
        if s.kind == "sp":
            if (
                s.x_pt <= 1 and s.y_pt <= 1
                and s.w_pt >= slide_w * 0.95
                and s.h_pt >= slide_h * 0.95
            ):
                bg = s
            else:
                text_shapes.append(s)
        elif s.kind == "pic":
            if is_logo_pic(s, slide_w, slide_h):
                logo = s
            elif hero is None or s.w_pt > (hero.w_pt or 0):
                hero = s
    return bg, hero, logo, text_shapes


def _solidfill_hex(sp_elem, theme: Theme) -> str:
    """Pull the inline solidFill hex from a <p:sp>'s spPr. Falls back to '000000'."""
    srgb = sp_elem.find("p:spPr/a:solidFill/a:srgbClr", NS)
    if srgb is not None:
        return srgb.get("val", "000000").upper()
    scheme = sp_elem.find("p:spPr/a:solidFill/a:schemeClr", NS)
    if scheme is not None:
        try:
            return theme.resolve(scheme.get("val"))
        except AttributeError:
            pass
    return "000000"


def _hero_data_url(hero: FlatShape, slide: Slide) -> Optional[str]:
    """SVG when present (preferred for the hero), else raster."""
    svg_ref = extract_svg_ref(hero.element, slide)
    if svg_ref is not None:
        b64 = base64.b64encode(svg_ref.blob).decode("ascii")
        return f"data:{svg_ref.content_type};base64,{b64}"
    img_ref = extract_image_ref(hero.element, slide)
    if img_ref is not None:
        b64 = base64.b64encode(img_ref.blob).decode("ascii")
        return f"data:{img_ref.content_type};base64,{b64}"
    return None


def _logo_data_url(logo: FlatShape, slide: Slide) -> Optional[str]:
    img_ref = extract_image_ref(logo.element, slide)
    if img_ref is not None:
        b64 = base64.b64encode(img_ref.blob).decode("ascii")
        return f"data:{img_ref.content_type};base64,{b64}"
    return None


def _escape(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
