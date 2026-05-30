"""Entry-point reader for a .pptx file.

Owns: opening the file, exposing parts (slides, theme, media), reading
slide dimensions. Does NOT own rendering, layout detection, or shape
walking — those modules consume the parts this module exposes.

Modules that need raw XML access (SVG-only blips, group transforms)
reach into `slide.element` themselves; this module deliberately stays
thin so it doesn't lock downstream code into python-pptx idioms.
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterator

import pptx as _pptx
from pptx.presentation import Presentation
from pptx.slide import Slide

EMU_PER_PT = 9525  # 1 point = 9525 English Metric Units


class Pptx:
    def __init__(self, path: str | Path):
        self.path = Path(path)
        if not self.path.exists():
            raise FileNotFoundError(self.path)
        self._prs: Presentation = _pptx.Presentation(str(self.path))

    @property
    def slide_size_pt(self) -> tuple[float, float]:
        """Slide width/height in points, read from <p:sldSz>.

        Per the spec: never hardcode dimensions. P&G is 1280x720pt;
        FrameTag and others differ.
        """
        return (
            self._prs.slide_width / EMU_PER_PT,
            self._prs.slide_height / EMU_PER_PT,
        )

    @property
    def slide_count(self) -> int:
        return len(self._prs.slides)

    def slides(self) -> Iterator[Slide]:
        return iter(self._prs.slides)

    @property
    def raw(self) -> Presentation:
        """Escape hatch: the underlying python-pptx Presentation."""
        return self._prs
